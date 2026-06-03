from __future__ import annotations

from io import BytesIO

import pandas as pd
from openpyxl import load_workbook
from pptx import Presentation

from utils.cdsfn_live import (
    CDSFN_BLOCKS,
    build_balanco_dre_dmpl_individual_pptx,
    build_credit_package_excel_cdsfn,
    build_display_table_cdsfn,
    build_excel_display_export_cdsfn,
    build_hierarchy_frame_cdsfn,
    build_excel_export_cdsfn,
    combine_normalized_blocks_cdsfn,
    combine_reference_periods_cdsfn,
    extract_metadata_cdsfn,
    fetch_documento_cdsfn,
    fetch_instituicoes_cdsfn,
    list_blocks_cdsfn,
    list_reference_periods_cdsfn,
    load_instituicoes_csv_cdsfn,
    normalize_long_cdsfn,
    pivot_wide_cdsfn,
    validate_json_cdsfn,
)


def _sample_payload() -> dict:
    return {
        "@cnpj": "05503849",
        "@codigoDocumento": "9011",
        "@tipoRemessa": "I",
        "@unidadeMedida": "1000",
        "@dataBase": "122025",
        "datasBaseReferencia": [
            {"@id": "dt1", "@data": "A122025"},
            {"@id": "dt2", "@data": "S122025"},
        ],
        "BalancoPatrimonial": {
            "contas": [
                {
                    "@id": "conta1",
                    "@nivel": "1",
                    "@descricao": "Ativo",
                    "@contaPai": "",
                    "valoresIndividualizados": [{"@dtBase": "dt1", "@valor": 656719.0}],
                },
                {
                    "@id": "conta2",
                    "@nivel": "1.1",
                    "@descricao": "Disponibilidades",
                    "@contaPai": "1",
                    "valoresIndividualizados": [{"@dtBase": "dt1", "@valor": 100.0}],
                },
                {
                    "@id": "conta3",
                    "@nivel": "2",
                    "@descricao": "Passivo",
                    "@contaPai": "",
                    "valoresIndividualizados": [{"@dtBase": "dt1", "@valor": 1000.0}],
                },
                {
                    "@id": "conta4",
                    "@nivel": "2.1",
                    "@descricao": "Patrimônio líquido",
                    "@contaPai": "2",
                    "valoresIndividualizados": [{"@dtBase": "dt1", "@valor": 500.0}],
                },
            ]
        },
        "DemonstracaoDoResultado": {
            "contas": [
                {
                    "@id": "conta-header",
                    "@nivel": "1",
                    "@descricao": "Outros resultados abrangentes",
                    "@contaPai": "",
                },
                {
                    "@id": "conta3",
                    "@nivel": "1.1.1",
                    "@descricao": "Receita com operações de crédito",
                    "@contaPai": "1.1",
                    "valoresIndividualizados": [
                        {"@dtBase": "dt2", "@valor": 211791.0},
                        {"@dtBase": "dt1", "@valor": 217098.0},
                    ],
                }
            ]
        },
    }


def _sample_payload_second_document() -> dict:
    payload = _sample_payload()
    payload["@dataBase"] = "122026"
    payload["datasBaseReferencia"] = [
        {"@id": "dt1", "@data": "A122026"},
        {"@id": "dt2", "@data": "S122026"},
        {"@id": "dt3", "@data": "A122025"},
    ]
    payload["BalancoPatrimonial"]["contas"][0]["valoresIndividualizados"] = [{"@dtBase": "dt1", "@valor": 700000.0}]
    payload["BalancoPatrimonial"]["contas"][1]["valoresIndividualizados"] = [{"@dtBase": "dt1", "@valor": 120.0}]
    payload["BalancoPatrimonial"]["contas"][2]["valoresIndividualizados"] = [{"@dtBase": "dt1", "@valor": 1100.0}]
    payload["BalancoPatrimonial"]["contas"][3]["valoresIndividualizados"] = [{"@dtBase": "dt1", "@valor": 550.0}]
    payload["DemonstracaoDoResultado"]["contas"][1]["valoresIndividualizados"] = [
        {"@dtBase": "dt2", "@valor": 220000.0},
        {"@dtBase": "dt3", "@valor": 217098.0},
    ]
    return payload


def _sample_payload_credit_package() -> dict:
    payload = _sample_payload()
    payload["DemonstracaoDoResultadoAbrangente"] = {
        "contas": [
            {
                "@id": "dra1",
                "@nivel": "1",
                "@descricao": "Resultado abrangente do período",
                "@contaPai": "",
                "valoresIndividualizados": [{"@dtBase": "dt1", "@valor": 215000.0}],
            }
        ]
    }
    payload["DemonstracaoDosFluxosDeCaixa"] = {
        "contas": [
            {
                "@id": "dfc1",
                "@nivel": "1",
                "@descricao": "Caixa líquido das atividades operacionais",
                "@contaPai": "",
                "valoresIndividualizados": [{"@dtBase": "dt1", "@valor": 35000.0}],
            }
        ]
    }
    payload["DemonstracaoDasMutacoesDoPatrimonioLiquido"] = {
        "contas": [
            {
                "@id": "dmpl1",
                "@nivel": "1",
                "@descricao": "Saldo no início do período",
                "@contaPai": "",
                "valoresIndividualizados": [{"@dtBase": "dt1", "@valor": 480.0}],
            }
        ]
    }
    return payload


def _sample_demonstrativos(payloads: list[dict], periodos_ref: list[str]) -> dict[str, pd.DataFrame]:
    demonstrativos: dict[str, pd.DataFrame] = {}
    for block_key, meta in CDSFN_BLOCKS.items():
        if not any(block_key in payload for payload in payloads):
            continue
        df_long = combine_normalized_blocks_cdsfn(payloads, block_key)
        df_view, _ = build_hierarchy_frame_cdsfn(df_long, block_key)
        for periodo_ref in periodos_ref:
            if periodo_ref not in df_view.columns:
                df_view[periodo_ref] = pd.NA
        demonstrativos[meta["sigla"]] = df_view
    return demonstrativos


def _pptx_text(prs: Presentation) -> str:
    chunks: list[str] = []
    for slide in prs.slides:
        chunks.append(_slide_text(slide))
    return "\n".join(chunks)


def _slide_text(slide) -> str:
    chunks: list[str] = []
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            chunks.append(shape.text)
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                for cell in row.cells:
                    chunks.append(cell.text)
    return "\n".join(chunks)


def _count_slides_containing(prs: Presentation, text: str) -> int:
    return sum(1 for slide in prs.slides if text in _slide_text(slide))


def _table_cells(prs: Presentation):
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    for cell in row.cells:
                        yield cell


class _DummyResponse:
    def __init__(self, body, status_code: int = 200, url: str = ""):
        self._body = body
        self.status_code = status_code
        self.url = url
        self.text = body if isinstance(body, str) else ""
        self.content = body.encode("utf-8") if isinstance(body, str) else b"{}"

    def json(self):
        return self._body

    def raise_for_status(self):
        if self.status_code >= 400:
            raise RuntimeError(f"http {self.status_code}")


class _DummyClient:
    def __init__(self):
        self.posts = []
        self.gets = []

    def post(self, url, json=None, headers=None, timeout=None):
        self.posts.append({"url": url, "json": json, "headers": headers, "timeout": timeout})
        pagina = json["numeroPagina"]
        if pagina == 0:
            return _DummyResponse(
                {
                    "content": [
                        {"nome": "Banco A", "cnpj": "12345678", "idBacen": "1"},
                        {"nome": "Banco B", "cnpj": "87654321", "idBacen": "2"},
                    ],
                    "number": 0,
                    "totalPages": 2,
                    "last": False,
                }
            )
        return _DummyResponse(
            {
                "content": [{"nome": "Banco C", "cnpj": "11112222", "idBacen": "3"}],
                "number": 1,
                "totalPages": 2,
                "last": True,
            }
        )

    def get(self, url, params=None, headers=None, timeout=None):
        self.gets.append({"url": url, "params": params, "headers": headers, "timeout": timeout})
        return _DummyResponse(_sample_payload(), url="https://example.com/final")


def test_validate_and_extract_metadata_cdsfn():
    payload = _sample_payload()
    validate_json_cdsfn(payload)
    metadata = extract_metadata_cdsfn(payload)
    assert metadata["cnpj"] == "05503849"
    assert metadata["codigo_documento"] == "9011"
    assert metadata["datas_base_map"]["dt1"] == "A122025"


def test_list_blocks_and_normalize_long_cdsfn():
    payload = _sample_payload()
    blocks = list_blocks_cdsfn(payload)
    assert [item["sigla"] for item in blocks] == ["BP", "DRE"]

    df_long = normalize_long_cdsfn(payload, "DemonstracaoDoResultado")
    assert list(df_long.columns) == [
        "cnpj",
        "codigo_documento",
        "demonstracao",
        "bloco_origem",
        "conta_id",
        "nivel",
        "descricao",
        "conta_pai",
        "ordem_conta",
        "dt_base",
        "dt_base_referencia",
        "valor",
        "unidade_medida",
    ]
    assert len(df_long) == 3
    assert set(df_long["dt_base_referencia"].dropna()) == {"A122025", "S122025"}
    assert df_long[df_long["conta_id"] == "conta-header"]["valor"].isna().all()


def test_pivot_wide_cdsfn_preserves_dates_as_columns():
    df_long = normalize_long_cdsfn(_sample_payload(), "DemonstracaoDoResultado")
    df_wide = pivot_wide_cdsfn(df_long)
    assert "A122025" in df_wide.columns
    assert "S122025" in df_wide.columns
    assert df_wide.iloc[0]["nivel"] == "1.1.1"


def test_fetch_instituicoes_cdsfn_iterates_pages_until_last():
    client = _DummyClient()
    df = fetch_instituicoes_cdsfn(page_size=500, session=client)
    assert len(client.posts) == 2
    assert list(df["nome"]) == ["Banco A", "Banco B", "Banco C"]
    assert client.posts[0]["json"]["incluirAgencias"] is False
    assert client.posts[0]["json"]["tamanhoPagina"] == 500


def test_fetch_documento_cdsfn_builds_expected_url_and_params():
    client = _DummyClient()
    payload, final_url = fetch_documento_cdsfn("05503849", "202512", session=client)
    assert payload["@codigoDocumento"] == "9011"
    assert final_url == "https://example.com/final"
    assert client.gets[0]["url"].endswith("/202512-9011-05503849.json")
    assert client.gets[0]["params"] == {"cnpj": "05503849", "anoMes": "202512"}


def test_fetch_documento_cdsfn_reports_empty_body_as_unpublished():
    class EmptyClient:
        def get(self, url, params=None, headers=None, timeout=None):
            return _DummyResponse("", url="https://example.com/empty")

    try:
        fetch_documento_cdsfn("60701190", "202603", session=EmptyClient())
    except FileNotFoundError as exc:
        assert "sem JSON publicado" in str(exc)
        assert "60701190" in str(exc)
        assert "202603" in str(exc)
    else:
        raise AssertionError("esperava FileNotFoundError para resposta vazia do CDSFN")


def test_build_excel_export_cdsfn_returns_bytes():
    payload = _sample_payload()
    metadata = extract_metadata_cdsfn(payload)
    df_long = normalize_long_cdsfn(payload, "BalancoPatrimonial")
    df_wide = pivot_wide_cdsfn(df_long)
    excel_bytes = build_excel_export_cdsfn(metadata, df_long, df_wide)
    assert isinstance(excel_bytes, bytes)
    assert len(excel_bytes) > 100


def test_load_instituicoes_csv_cdsfn_parses_bootstrap_file(tmp_path):
    csv_path = tmp_path / "instituicoes.csv"
    csv_path.write_text(
        "\n".join(
            [
                "SEGMENTO;-;",
                "NOME DA INSTITUICAO;-;",
                "CNPJ;-;",
                "PAIS;Todos;",
                "UF;-;",
                "MUNICIPIO;-;",
                ";",
                "CNPJ;NOME DA INSTITUICAO;MUNICIPIO;UF;SITUAÇÃO;",
                '5503849;"CLOUDWALK FINANCEIRA S.A."; "SAO PAULO";"SP";"Autorizada em Atividade";',
                '60394079;"BANCO ITAUBANK S.A.";"SAO PAULO";"SP";"Autorizada em Atividade";',
            ]
        ),
        encoding="latin1",
    )
    df = load_instituicoes_csv_cdsfn(csv_path)
    assert list(df["cnpj"]) == ["60394079", "05503849"]
    assert "label" in df.columns


def test_build_hierarchy_frame_cdsfn_splits_balance_sections():
    df_long = normalize_long_cdsfn(_sample_payload(), "BalancoPatrimonial")
    df_view, value_cols = build_hierarchy_frame_cdsfn(df_long, "BalancoPatrimonial")
    assert "A122025" in value_cols
    assert set(df_view["section"]) >= {"Ativo", "Passivo", "Patrimônio líquido"}
    ativo_row = df_view[df_view["descricao"] == "Ativo"].iloc[0]
    assert int(ativo_row["depth"]) == 0


def test_list_reference_periods_cdsfn_returns_sorted_unique_refs():
    refs = list_reference_periods_cdsfn(_sample_payload())
    assert [item["raw"] for item in refs] == ["S122025", "A122025"]
    assert refs[0]["label"] == "S dez/25"


def test_build_display_table_cdsfn_uses_selected_columns_only():
    df_long = normalize_long_cdsfn(_sample_payload(), "DemonstracaoDoResultado")
    df_display = build_display_table_cdsfn(
        df_long,
        "DemonstracaoDoResultado",
        ["S122025"],
        column_labels={"S122025": "S dez/25"},
    )
    assert list(df_display.columns) == ["Conta", "S dez/25"]
    assert "Receita com operações de crédito" in df_display["Conta"].iloc[-1]


def test_build_excel_display_export_cdsfn_returns_bytes():
    payload = _sample_payload()
    df_long = normalize_long_cdsfn(payload, "BalancoPatrimonial")
    df_view, value_cols = build_hierarchy_frame_cdsfn(df_long, "BalancoPatrimonial")
    excel_bytes = build_excel_display_export_cdsfn(
        df_view,
        ["A122025"],
        column_labels={"A122025": "A dez/25"},
        sheet_name="BP",
    )
    assert isinstance(excel_bytes, bytes)
    assert len(excel_bytes) > 100
    wb = load_workbook(BytesIO(excel_bytes))
    ws = wb["BP"]
    assert ws["A1"].value == "Conta"
    assert ws["B1"].value == "A dez/25"
    assert ws["A2"].font.bold is True
    assert ws["A2"].fill.fgColor.rgb in {"FFF6F6F6", "00F6F6F6"}
    assert ws["B2"].value == 656719
    assert ws["B2"].number_format == "#,##0.00"


def test_build_credit_package_excel_cdsfn_returns_multisheet_workbook():
    payload = _sample_payload_credit_package()
    excel_bytes = build_credit_package_excel_cdsfn(
        payloads=[payload],
        institution_label="Banco Exemplo",
        cnpj="05503849",
        periodos_ref_sel=["A122025"],
        column_labels={"A122025": "A dez/25"},
        reference_label="Anual (A)",
        document_periods=["202512"],
    )

    assert isinstance(excel_bytes, bytes)
    assert len(excel_bytes) > 100

    wb = load_workbook(BytesIO(excel_bytes))
    assert wb.sheetnames == ["BP", "DRE", "DRA", "DFC", "DMPL"]

    ws_bp = wb["BP"]
    assert ws_bp["A1"].value == "BP - Balanço Patrimonial"
    assert "Instituição: Banco Exemplo" in ws_bp["A2"].value
    assert "Fonte: Documento 9011 JSON do Banco Central" in ws_bp["A3"].value
    assert "Competências: 202512" in ws_bp["A3"].value
    assert "Unidade monetária: 1000" in ws_bp["A4"].value
    assert "GMT-3" in ws_bp["A4"].value
    assert ws_bp["A5"].value == "Conta"
    assert ws_bp["B5"].value == "A dez/25"
    assert ws_bp["B6"].value == 656719

    ws_dra = wb["DRA"]
    assert ws_dra["A6"].value == "1 Resultado abrangente do período"
    assert ws_dra["B6"].value == 215000


def test_build_balanco_dre_dmpl_individual_pptx_returns_readable_deck():
    payload = _sample_payload_credit_package()
    demonstrativos = _sample_demonstrativos([payload], ["A122025"])

    pptx_bytes = build_balanco_dre_dmpl_individual_pptx(
        institution_name="Banco Exemplo",
        cnpj="05503849",
        periodo_atual="A122025",
        referencia="Anual (A)",
        remessa="I",
        unidade="1000",
        demonstrativos=demonstrativos,
        column_labels={"A122025": "A dez/25"},
        document_periods=["202512"],
    )

    assert isinstance(pptx_bytes, bytes)
    assert len(pptx_bytes) > 1000
    prs = Presentation(BytesIO(pptx_bytes))
    assert len(prs.slides) >= 6
    text = _pptx_text(prs)
    assert "Balanço, DRE e DMPL" in text
    assert "BP — Balanço Patrimonial" in text
    assert "DRE — Demonstração do Resultado" in text
    assert "Unidade: R$ mil" in text
    assert "Fonte: Banco Central do Brasil" in text
    assert _count_slides_containing(prs, "BP — Balanço Patrimonial") == 1
    assert _count_slides_containing(prs, "DRE — Demonstração do Resultado") == 1
    assert all(cell.text_frame.word_wrap is False for cell in _table_cells(prs))


def test_build_balanco_dre_dmpl_individual_pptx_summary_cards_are_fixed_order():
    payload = _sample_payload_credit_package()
    demonstrativos = _sample_demonstrativos([payload], ["A122025"])

    pptx_bytes = build_balanco_dre_dmpl_individual_pptx(
        institution_name="Banco Exemplo",
        cnpj="05503849",
        periodo_atual="A122025",
        referencia="Anual (A)",
        remessa="I",
        unidade="1000",
        demonstrativos=demonstrativos,
        column_labels={"A122025": "A dez/25"},
        document_periods=["202512"],
    )

    prs = Presentation(BytesIO(pptx_bytes))
    summary_text = _slide_text(prs.slides[1])
    expected = [
        "ATIVO TOTAL",
        "OPERAÇÕES DE CRÉDITO",
        "PROVISÃO PARA PERDAS ESPERADAS",
        "PATRIMÔNIO LÍQUIDO",
    ]
    positions = [summary_text.index(item) for item in expected]
    assert positions == sorted(positions)
    assert "RESULTADO ABRANGENTE" not in summary_text
    assert "CAPITAL SOCIAL" not in summary_text
    assert "RESERVAS DE LUCROS" not in summary_text


def test_build_balanco_dre_dmpl_individual_pptx_includes_comparative_columns():
    payloads = [_sample_payload_second_document(), _sample_payload()]
    demonstrativos = _sample_demonstrativos(payloads, ["A122026", "A122025"])

    pptx_bytes = build_balanco_dre_dmpl_individual_pptx(
        institution_name="Banco Exemplo",
        cnpj="05503849",
        periodo_atual="A122026",
        periodo_anterior="A122025",
        referencia="Anual (A)",
        remessa="I",
        unidade="1000",
        demonstrativos=demonstrativos,
        column_labels={"A122026": "A dez/26", "A122025": "A dez/25"},
        document_periods=["202612", "202512"],
    )

    prs = Presentation(BytesIO(pptx_bytes))
    text = _pptx_text(prs)
    assert "Var. Abs." in text
    assert "Var. %" in text
    assert "43.281" in text


def test_build_balanco_dre_dmpl_individual_pptx_marks_missing_statements():
    demonstrativos = _sample_demonstrativos([_sample_payload()], ["A122025"])

    pptx_bytes = build_balanco_dre_dmpl_individual_pptx(
        institution_name="Banco Exemplo",
        cnpj="05503849",
        periodo_atual="A122025",
        referencia="Anual (A)",
        remessa="I",
        unidade="1000",
        demonstrativos=demonstrativos,
        column_labels={"A122025": "A dez/25"},
        document_periods=["202512"],
    )

    prs = Presentation(BytesIO(pptx_bytes))
    text = _pptx_text(prs)
    assert "DMPL — Demonstração das Mutações do Patrimônio Líquido" in text
    assert "Demonstração não disponível para os filtros selecionados." in text


def test_build_balanco_dre_dmpl_individual_pptx_keeps_dfc_on_one_complete_slide():
    payload = _sample_payload_credit_package()
    payload["DemonstracaoDosFluxosDeCaixa"]["contas"] = [
        {
            "@id": f"dfc{i}",
            "@nivel": f"1.{i}",
            "@descricao": f"Fluxo linha {i:02d}",
            "@contaPai": "1",
            "valoresIndividualizados": [{"@dtBase": "dt1", "@valor": float(i * 10)}],
        }
        for i in range(1, 33)
    ]
    demonstrativos = _sample_demonstrativos([payload], ["A122025"])

    pptx_bytes = build_balanco_dre_dmpl_individual_pptx(
        institution_name="Banco Exemplo",
        cnpj="05503849",
        periodo_atual="A122025",
        referencia="Anual (A)",
        remessa="I",
        unidade="1000",
        demonstrativos=demonstrativos,
        column_labels={"A122025": "A dez/25"},
        document_periods=["202512"],
    )

    prs = Presentation(BytesIO(pptx_bytes))
    assert _count_slides_containing(prs, "DFC — Demonstração dos Fluxos de Caixa") == 1
    assert "Fluxo linha 32" in _pptx_text(prs)


def test_combine_reference_periods_cdsfn_dedupes_and_sorts():
    refs = combine_reference_periods_cdsfn([_sample_payload(), _sample_payload_second_document()])
    assert [item["raw"] for item in refs][:3] == ["S122026", "A122026", "S122025"]


def test_combine_normalized_blocks_cdsfn_merges_two_documents():
    df_long = combine_normalized_blocks_cdsfn(
        [_sample_payload(), _sample_payload_second_document()],
        "DemonstracaoDoResultado",
    )
    assert set(df_long["dt_base_referencia"].dropna()) >= {"A122025", "S122025", "S122026"}
    receita = df_long[df_long["descricao"] == "Receita com operações de crédito"]
    assert len(receita) == 3


def test_combine_normalized_blocks_cdsfn_raises_on_conflicting_duplicate_values():
    payload_a = _sample_payload()
    payload_b = _sample_payload()
    payload_b["DemonstracaoDoResultado"]["contas"][1]["valoresIndividualizados"][1]["@valor"] = 999999.0
    try:
        combine_normalized_blocks_cdsfn([payload_a, payload_b], "DemonstracaoDoResultado")
    except ValueError as exc:
        assert "Conflito de valores" in str(exc)
    else:
        raise AssertionError("Era esperado conflito explícito ao combinar valores divergentes.")
