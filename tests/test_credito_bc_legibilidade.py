"""Garantias de leitura da seção "Estatísticas Crédito BC".

Cada teste aqui trava um defeito medido no app e corrigido: rótulo de barra em
nove tamanhos diferentes na mesma barra, rótulos de fim de linha desenhados no
mesmo pixel, empilhado exportado como linha e card de dois eixos achatado num
eixo só.
"""

from __future__ import annotations

import itertools
import math
from io import BytesIO

import pandas as pd
import pytest
from pptx import Presentation
from pptx.oxml.ns import qn

from utils.sgs_credit_analytics import (
    COR_ROTULO_SOBRE,
    ITAU_ORANGE,
    PALETA_LINHA,
    PALETA_PREENCHIMENTO,
    TAMANHO_ROTULO_BARRA_PX,
    TAMANHO_ROTULO_PX,
    _altura_area_plotagem,
    formatar_numero,
    _espalhar_em_pixels,
    bar_line_figure,
    line_figure,
    stacked_figure,
)
from utils.sgs_credit_pptx_export import exportar_figuras_pptx, figura_para_painel


# =============================================================================
# APOIO
# =============================================================================

def _luminancia(cor: str) -> float:
    canais = [int(cor.lstrip("#")[i:i + 2], 16) / 255 for i in (0, 2, 4)]
    linear = [c / 12.92 if c <= 0.04045 else ((c + 0.055) / 1.055) ** 2.4 for c in canais]
    return 0.2126 * linear[0] + 0.7152 * linear[1] + 0.0722 * linear[2]


def _contraste(a: str, b: str) -> float:
    la, lb = _luminancia(a), _luminancia(b)
    return (max(la, lb) + 0.05) / (min(la, lb) + 0.05)


def _lab(cor: str) -> tuple[float, float, float]:
    canais = [int(cor.lstrip("#")[i:i + 2], 16) / 255 for i in (0, 2, 4)]
    r, g, b = [c / 12.92 if c <= 0.04045 else ((c + 0.055) / 1.055) ** 2.4 for c in canais]
    x = (0.4124 * r + 0.3576 * g + 0.1805 * b) / 0.95047
    y = 0.2126 * r + 0.7152 * g + 0.0722 * b
    z = (0.0193 * r + 0.1192 * g + 0.9505 * b) / 1.08883

    def f(t: float) -> float:
        return t ** (1 / 3) if t > 216 / 24389 else (841 / 108) * t + 4 / 29

    fx, fy, fz = f(x), f(y), f(z)
    return (116 * fy - 16, 500 * (fx - fy), 200 * (fy - fz))


def _delta_e(a: str, b: str) -> float:
    return math.sqrt(sum((x - y) ** 2 for x, y in zip(_lab(a), _lab(b))))


def _wide(meses: int = 12, series: int = 9) -> pd.DataFrame:
    index = pd.date_range("2025-08-31", periods=meses, freq="ME")
    return pd.DataFrame(
        {
            f"serie_{i}": [10.0 + i * 0.4 + passo * 0.05 for passo in range(meses)]
            for i in range(series)
        },
        index=index,
    )


# =============================================================================
# PALETA
# =============================================================================

def test_paleta_de_linha_e_legivel_no_branco_e_distinguivel_entre_si():
    """Cinco cores, todas visíveis sobre o papel e nenhuma confundível.

    A paleta anterior tinha dez cores das quais oito eram cinzas com ΔE entre
    6,5 e 10,9 — 20 dos 45 pares eram indistinguíveis.
    """
    for cor in PALETA_LINHA:
        assert _contraste(cor, "#FFFFFF") >= 3.0, f"{cor} some no papel branco"
    for a, b in itertools.combinations(PALETA_LINHA, 2):
        assert _delta_e(a, b) >= 25.0, f"{a} e {b} são confundíveis (ΔE {_delta_e(a, b):.1f})"


def test_todo_preenchimento_declara_cor_de_rotulo_que_passa_no_contraste():
    """Preto ou branco declarado, nunca a escolha automática da biblioteca.

    Sozinha, a biblioteca escolhia #444444 sobre o laranja: 3,19:1, reprovado.
    """
    for preenchimento in PALETA_PREENCHIMENTO:
        rotulo = COR_ROTULO_SOBRE[preenchimento]
        assert rotulo in {"#FFFFFF", "#141414"}
        assert _contraste(preenchimento, rotulo) >= 4.5, (
            f"rótulo {rotulo} sobre {preenchimento} tem contraste insuficiente"
        )


def test_fatias_vizinhas_do_empilhado_nunca_repetem_cor():
    for posicao in range(len(PALETA_PREENCHIMENTO) * 2):
        atual = PALETA_PREENCHIMENTO[posicao % len(PALETA_PREENCHIMENTO)]
        seguinte = PALETA_PREENCHIMENTO[(posicao + 1) % len(PALETA_PREENCHIMENTO)]
        assert _delta_e(atual, seguinte) >= 25.0


# =============================================================================
# RÓTULO DE BARRA
# =============================================================================

def test_rotulo_de_barra_tem_tamanho_unico_e_nunca_deita():
    """Trava de tamanho ligada e ângulo fixo em zero.

    Sem isso a biblioteca desenhava, na mesma barra, rótulos de 0,2 px a 11,4
    px e girava um deles em 90 graus.
    """
    wide = _wide(series=6)
    fig = stacked_figure(
        wide, [f"serie_{i}" for i in range(6)],
        title="Empilhado", y_title="R$ bi",
    )
    assert fig.layout.uniformtext.mode == "hide"
    assert fig.layout.uniformtext.minsize == TAMANHO_ROTULO_PX
    for trace in fig.data:
        assert trace.textangle == 0
        assert trace.insidetextfont.size == TAMANHO_ROTULO_PX
        assert trace.insidetextfont.color == COR_ROTULO_SOBRE[trace.marker.color]


def test_total_do_empilhado_e_anotacao_e_nao_serie_invisivel():
    """O total como série de texto fazia o exportador trocar barra por linha."""
    wide = _wide(series=4)
    total = wide[[f"serie_{i}" for i in range(4)]].sum(axis=1)
    fig = stacked_figure(
        wide, [f"serie_{i}" for i in range(4)],
        title="Empilhado", y_title="R$ bi", total=total,
    )
    assert {trace.type for trace in fig.data} == {"bar"}
    assert len(fig.data) == 4
    assert any(anotacao.text for anotacao in fig.layout.annotations)


# =============================================================================
# RÓTULO DE FIM DE LINHA
# =============================================================================

def test_rotulos_de_fim_de_linha_nunca_se_sobrepoem():
    """Nove séries, nove rótulos, nenhum par colado.

    No app, "8,8%" e "6,3%" eram desenhados no mesmo pixel: a rotina calculava
    o espaçamento em unidades de dado e convertia para pixel usando uma altura
    de área fixada em 285 px, quando a área real era 323 px.
    """
    wide = _wide(series=9)
    fig = line_figure(
        wide, [f"serie_{i}" for i in range(9)],
        title="Nove linhas", y_title="%", suffix="%",
    )
    anotacoes = list(fig.layout.annotations)
    assert len(anotacoes) == 9

    altura_area = _altura_area_plotagem(fig)
    menor, maior = fig.layout.yaxis.range
    pixels_por_unidade = altura_area / (maior - menor)
    posicoes = sorted(
        (maior - anotacao.y) * pixels_por_unidade + anotacao.ay
        for anotacao in anotacoes
    )
    for anterior, seguinte in zip(posicoes, posicoes[1:]):
        assert seguinte - anterior >= TAMANHO_ROTULO_PX, (
            f"rótulos a {seguinte - anterior:.1f} px um do outro"
        )


def test_area_de_plotagem_e_deterministica():
    """A conta dos rótulos só fecha se a geometria for conhecida."""
    fig = line_figure(_wide(series=3), ["serie_0", "serie_1", "serie_2"],
                      title="Três linhas", y_title="%")
    assert fig.layout.margin.autoexpand is False
    esperado = fig.layout.height - fig.layout.margin.t - fig.layout.margin.b
    assert _altura_area_plotagem(fig) == pytest.approx(esperado)


def test_espalhar_preserva_a_ordem_mesmo_quando_nao_cabe():
    alvos = [10.0, 11.0, 12.0, 13.0, 14.0, 15.0]
    posicoes = _espalhar_em_pixels(alvos, 20.0, 60.0)
    assert posicoes == sorted(posicoes)
    assert min(posicoes) >= 0
    assert max(posicoes) <= 60.0


def test_acima_de_cinco_series_a_cor_repete_com_tracejado_e_o_nome_vai_na_ponta():
    wide = _wide(series=7)
    fig = line_figure(wide, [f"serie_{i}" for i in range(7)],
                      title="Sete linhas", y_title="%")
    assert fig.layout.showlegend is False
    assert fig.data[5].line.color == fig.data[0].line.color
    assert fig.data[5].line.dash == "dash"
    assert fig.data[0].line.dash == "solid"
    # nome + valor, e não só o valor
    assert any("serie_0" in anotacao.text for anotacao in fig.layout.annotations)


# =============================================================================
# EXPORTAÇÃO
# =============================================================================

def test_empilhado_sai_empilhado_no_pptx():
    wide = _wide(series=4)
    total = wide[[f"serie_{i}" for i in range(4)]].sum(axis=1)
    fig = stacked_figure(
        wide, [f"serie_{i}" for i in range(4)],
        title="Empilhado", y_title="R$ bi", total=total,
    )
    painel = figura_para_painel(fig)
    assert painel.tipo_grafico == "column_stacked"
    assert painel.ordem_series == [f"serie_{i}" for i in range(4)]

    blob, _ = exportar_figuras_pptx([fig], titulo_deck="Teste")
    deck = Presentation(BytesIO(blob))
    grafico = [s.chart for s in deck.slides[0].shapes if s.has_chart][0]
    plot_area = grafico._chartSpace.chart.plotArea
    assert plot_area.find(qn("c:barChart")) is not None
    assert plot_area.find(qn("c:lineChart")) is None


def test_card_de_volume_e_prazo_mantem_os_dois_eixos_no_pptx():
    """R$ 246-353 bi e 23-34 meses não podem dividir o mesmo eixo."""
    index = pd.date_range("2025-08-31", periods=12, freq="ME")
    wide = pd.DataFrame(
        {
            "concessoes_livre_pj": [250_000.0 + i * 8_000 for i in range(12)],
            "prazo_livre_pj": [24.0 + i * 0.8 for i in range(12)],
        },
        index=index,
    )
    fig = bar_line_figure(
        wide, bar_alias="concessoes_livre_pj", line_alias="prazo_livre_pj",
        title="Concessões PJ",
    )
    painel = figura_para_painel(fig)
    assert painel.tipo_grafico == "column_line"
    assert len(painel.series_secundarias) == 1

    blob, meta = exportar_figuras_pptx([fig], titulo_deck="Teste")
    deck = Presentation(BytesIO(blob))
    grafico = [s.chart for s in deck.slides[0].shapes if s.has_chart][0]
    plot_area = grafico._chartSpace.chart.plotArea
    assert plot_area.find(qn("c:barChart")) is not None
    assert plot_area.find(qn("c:lineChart")) is not None
    assert len(plot_area.findall(qn("c:valAx"))) == 2
    assert meta["detalhe"][0]["eixo_secundario"] is True


# =============================================================================
# VALIDADE DO ARQUIVO POWERPOINT
# =============================================================================
# O PowerPoint recusa o arquivo inteiro ("é preciso reparar") quando a ordem
# dos elementos foge da sequência do schema. Foi o que aconteceu com o deck de
# Concessões: o <c:lineChart> do eixo secundário entrava depois dos eixos, e o
# <c:marker> da série movida entrava depois de <c:val>.

_NS = "{http://schemas.openxmlformats.org/drawingml/2006/chart}"

# Grupos de gráfico devem vir todos antes de qualquer eixo (CT_PlotArea).
_GRUPOS = {
    "areaChart", "area3DChart", "lineChart", "line3DChart", "stockChart",
    "radarChart", "scatterChart", "pieChart", "pie3DChart", "doughnutChart",
    "barChart", "bar3DChart", "ofPieChart", "surfaceChart", "surface3DChart",
    "bubbleChart",
}
_EIXOS = {"valAx", "catAx", "dateAx", "serAx"}

_SEQUENCIAS = {
    "lineChart": [
        "grouping", "varyColors", "ser", "dLbls", "dropLines", "hiLowLines",
        "upDownBars", "marker", "smooth", "axId", "extLst",
    ],
    "barChart": [
        "barDir", "grouping", "varyColors", "ser", "dLbls", "gapWidth",
        "overlap", "serLines", "axId", "extLst",
    ],
    "catAx": [
        "axId", "scaling", "delete", "axPos", "majorGridlines", "minorGridlines",
        "title", "numFmt", "majorTickMark", "minorTickMark", "tickLblPos",
        "spPr", "txPr", "crossAx", "crosses", "crossesAt", "auto", "lblAlgn",
        "lblOffset", "tickLblSkip", "tickMarkSkip", "noMultiLvlLbl", "extLst",
    ],
    "valAx": [
        "axId", "scaling", "delete", "axPos", "majorGridlines", "minorGridlines",
        "title", "numFmt", "majorTickMark", "minorTickMark", "tickLblPos",
        "spPr", "txPr", "crossAx", "crosses", "crossesAt", "crossBetween",
        "majorUnit", "minorUnit", "dispUnits", "extLst",
    ],
}
_SEQUENCIA_SER = [
    "idx", "order", "tx", "spPr", "marker", "invertIfNegative",
    "pictureOptions", "dPt", "dLbls", "trendline", "errBars", "cat", "val",
    "shape", "smooth", "extLst",
]


# ST_DLblPos aceito por tipo de grupo. Coluna não aceita r/t/b/l: escrever uma
# dessas numa série de coluna faz o PowerPoint recusar o arquivo inteiro.
_POSICOES_VALIDAS = {
    "barChart": {"ctr", "inBase", "inEnd", "outEnd"},
    "lineChart": {"ctr", "l", "r", "t", "b"},
    "areaChart": {"ctr"},
    "pieChart": {"bestFit", "ctr", "inEnd", "outEnd"},
    "scatterChart": {"ctr", "l", "r", "t", "b"},
}


def _local(elemento) -> str:
    return elemento.tag.replace(_NS, "")


def _em_ordem(elemento, sequencia: list[str]) -> bool:
    posicoes = [
        sequencia.index(_local(filho))
        for filho in elemento
        if _local(filho) in sequencia
    ]
    return posicoes == sorted(posicoes)


def validar_xml_do_grafico(xml: bytes) -> list[str]:
    """Erros de ordem que fazem o PowerPoint pedir reparo. Vazio = válido."""
    from lxml import etree

    erros: list[str] = []
    raiz = etree.fromstring(xml)
    for plot_area in raiz.iter(f"{_NS}plotArea"):
        filhos = [_local(e) for e in plot_area]
        indices_grupo = [i for i, n in enumerate(filhos) if n in _GRUPOS]
        indices_eixo = [i for i, n in enumerate(filhos) if n in _EIXOS]
        if indices_grupo and indices_eixo and max(indices_grupo) > min(indices_eixo):
            erros.append(f"grupo de gráfico depois de eixo: {filhos}")
        for elemento in plot_area:
            nome = _local(elemento)
            if nome in _SEQUENCIAS and not _em_ordem(elemento, _SEQUENCIAS[nome]):
                erros.append(f"<c:{nome}> fora de ordem: {[_local(f) for f in elemento]}")
            validas = _POSICOES_VALIDAS.get(nome)
            for ser in elemento.findall(f"{_NS}ser"):
                if not _em_ordem(ser, _SEQUENCIA_SER):
                    erros.append(
                        f"<c:ser> de <c:{nome}> fora de ordem: "
                        f"{[_local(f) for f in ser]}"
                    )
                if validas is None:
                    continue
                for posicao in ser.iter(f"{_NS}dLblPos"):
                    if posicao.get("val") not in validas:
                        erros.append(
                            f"dLblPos '{posicao.get('val')}' inválido em <c:{nome}>"
                        )
    return erros


def _graficos_do_deck(blob: bytes) -> list[bytes]:
    import zipfile

    arquivo = zipfile.ZipFile(BytesIO(blob))
    return [
        arquivo.read(nome)
        for nome in arquivo.namelist()
        if nome.startswith("ppt/charts/chart") and nome.endswith(".xml")
    ]


def test_deck_de_concessoes_respeita_a_sequencia_do_schema():
    """O card de volume e prazo é o que quebrava o arquivo inteiro."""
    index = pd.date_range("2025-08-31", periods=12, freq="ME")
    wide = pd.DataFrame(
        {
            "concessoes_livre_pj": [250_000.0 + i * 8_000 for i in range(12)],
            "prazo_livre_pj": [24.0 + i * 0.8 for i in range(12)],
        },
        index=index,
    )
    fig = bar_line_figure(
        wide, bar_alias="concessoes_livre_pj", line_alias="prazo_livre_pj",
        title="Concessões PJ",
    )
    blob, _ = exportar_figuras_pptx([fig], titulo_deck="Concessões")
    for xml in _graficos_do_deck(blob):
        assert validar_xml_do_grafico(xml) == []


def test_deck_de_empilhado_e_de_linha_respeitam_a_sequencia_do_schema():
    wide = _wide(series=4)
    total = wide[[f"serie_{i}" for i in range(4)]].sum(axis=1)
    figuras = [
        stacked_figure(
            wide, [f"serie_{i}" for i in range(4)],
            title="Empilhado", y_title="R$ bi", total=total,
        ),
        line_figure(
            _wide(series=7), [f"serie_{i}" for i in range(7)],
            title="Linhas", y_title="%",
        ),
    ]
    blob, _ = exportar_figuras_pptx(figuras, titulo_deck="Misto")
    for xml in _graficos_do_deck(blob):
        assert validar_xml_do_grafico(xml) == []


def test_validador_pega_a_inversao_que_quebrava_o_arquivo():
    """Prova que o teste acima falharia se a regressão voltasse."""
    xml = f"""<c:chartSpace xmlns:c="{_NS[1:-1]}">
      <c:chart><c:plotArea>
        <c:barChart><c:axId val="1"/></c:barChart>
        <c:valAx><c:axId val="2"/></c:valAx>
        <c:lineChart><c:axId val="3"/></c:lineChart>
      </c:plotArea></c:chart>
    </c:chartSpace>""".encode()
    assert validar_xml_do_grafico(xml)


def test_validador_pega_posicao_de_rotulo_invalida_em_coluna():
    """A segunda causa do deck recusado: dLblPos "r" numa série de coluna."""
    xml = f"""<c:chartSpace xmlns:c="{_NS[1:-1]}">
      <c:chart><c:plotArea>
        <c:barChart>
          <c:ser><c:idx val="0"/><c:dLbls><c:dLbl>
            <c:dLblPos val="r"/>
          </c:dLbl></c:dLbls></c:ser>
        </c:barChart>
        <c:valAx><c:axId val="2"/></c:valAx>
      </c:plotArea></c:chart>
    </c:chartSpace>""".encode()
    erros = validar_xml_do_grafico(xml)
    assert any("dLblPos" in erro for erro in erros)


# =============================================================================
# LEITURA DAS BARRAS
# =============================================================================

def test_primeira_barra_nao_e_cortada_pela_faixa_do_eixo():
    """A faixa precisa de meia categoria de folga antes da primeira barra.

    Começando exatamente no centro dela, o Plotly desenhava só a metade
    direita — e o duplo clique "consertava" porque restaura o autorange.
    """
    index = pd.date_range("2025-08-31", periods=12, freq="ME")
    wide = pd.DataFrame(
        {
            "concessoes_livre_pj": [250_000.0 + i * 8_000 for i in range(12)],
            "prazo_livre_pj": [24.0 + i * 0.8 for i in range(12)],
        },
        index=index,
    )
    fig = bar_line_figure(
        wide, bar_alias="concessoes_livre_pj", line_alias="prazo_livre_pj",
        title="Concessões PJ",
    )
    inicio = pd.Timestamp(fig.layout.xaxis.range[0])
    fim = pd.Timestamp(fig.layout.xaxis.range[1])
    assert inicio < index[0], "a faixa começa em cima da primeira barra"
    assert (index[0] - inicio) >= pd.Timedelta(days=10)
    assert fim > index[-1]


def test_barra_larga_rotula_todos_os_periodos_com_fonte_maior():
    index = pd.date_range("2025-08-31", periods=12, freq="ME")
    wide = pd.DataFrame(
        {
            "concessoes_livre_pj": [250_000.0 + i * 8_000 for i in range(12)],
            "prazo_livre_pj": [24.0 + i * 0.8 for i in range(12)],
        },
        index=index,
    )
    fig = bar_line_figure(
        wide, bar_alias="concessoes_livre_pj", line_alias="prazo_livre_pj",
        title="Concessões PJ",
    )
    barra = fig.data[0]
    assert sum(1 for t in barra.text if t) == 12
    assert barra.textfont.size == TAMANHO_ROTULO_BARRA_PX
    assert fig.layout.uniformtext.minsize == TAMANHO_ROTULO_BARRA_PX
    assert barra.textangle == 0


def test_eixo_secundario_e_a_linha_saem_em_laranja():
    """Cor é o que diz que "meses" se lê à direita, sem precisar de nota."""
    index = pd.date_range("2025-08-31", periods=6, freq="ME")
    wide = pd.DataFrame(
        {
            "concessoes_livre_pj": [250_000.0 + i * 8_000 for i in range(6)],
            "prazo_livre_pj": [24.0 + i * 0.8 for i in range(6)],
        },
        index=index,
    )
    fig = bar_line_figure(
        wide, bar_alias="concessoes_livre_pj", line_alias="prazo_livre_pj",
        title="Concessões PJ",
    )
    assert fig.data[1].line.color == ITAU_ORANGE
    assert fig.layout.yaxis2.tickfont.color == ITAU_ORANGE
    assert ITAU_ORANGE in fig.data[1].name  # item da legenda tingido
    # O nome com marcação HTML não pode vazar para o deck.
    assert figura_para_painel(fig).ordem_series[1] == "Prazo médio"


# =============================================================================
# CARD DE MEIA LARGURA
# =============================================================================

def test_card_compacto_encolhe_a_margem_e_mantem_o_nome_na_ponta_da_linha():
    """Meia largura preserva o que torna oito linhas próximas legíveis."""
    wide = _wide(series=8)
    compacto = line_figure(
        wide, [f"serie_{i}" for i in range(8)],
        title="Oito linhas", y_title="%", suffix="%", compacto=True,
    )
    largo = line_figure(
        wide, [f"serie_{i}" for i in range(8)],
        title="Oito linhas", y_title="%", suffix="%",
    )
    assert compacto.layout.height < largo.layout.height
    assert compacto.layout.margin.r < largo.layout.margin.r
    assert compacto.layout.showlegend is False
    assert len(compacto.layout.annotations) == 8
    assert any("serie_0" in a.text for a in compacto.layout.annotations)

    altura_area = _altura_area_plotagem(compacto)
    menor, maior = compacto.layout.yaxis.range
    escala = altura_area / (maior - menor)
    posicoes = sorted((maior - a.y) * escala + a.ay for a in compacto.layout.annotations)
    for anterior, seguinte in zip(posicoes, posicoes[1:]):
        assert seguinte - anterior >= TAMANHO_ROTULO_PX


def test_eixo_categorico_nao_quebra_o_estilo():
    """A barra por UF passa pela mesma função de estilo dos gráficos de tempo.

    Com siglas no eixo X, a extração de datas estourava e derrubava a aba
    "Brasil e regiões" inteira.
    """
    import plotly.graph_objects as go

    from utils.sgs_credit_analytics import _valid_trace_dates, aplicar_estilo

    fig = go.Figure(go.Bar(x=["RR", "SP", "BA"], y=[1.0, 2.0, 3.0]))
    assert _valid_trace_dates(fig) == []
    aplicar_estilo(fig, title="Por UF", y_title="%", legenda=False)
    assert fig.layout.xaxis.tickmode != "array"


# =============================================================================
# NÚMERO NO PADRÃO BRASILEIRO
# =============================================================================

@pytest.mark.parametrize(
    "valor,casas,esperado",
    [
        (7372.0, 1, "7.372,0"),
        (1234567.89, 2, "1.234.567,89"),
        (729.2, 1, "729,2"),
        (4.88, 2, "4,88"),
        (-6.3, 1, "-6,3"),
        (0.2, 1, "0,2"),
    ],
)
def test_ponto_separa_milhar_e_virgula_separa_decimal(valor, casas, esperado):
    assert formatar_numero(valor, casas) == esperado


def test_rotulos_do_grafico_usam_o_padrao_brasileiro():
    index = pd.date_range("2025-08-31", periods=12, freq="ME")
    wide = pd.DataFrame(
        {f"serie_{i}": [1_500_000.0 + i * 500_000 + p * 1000 for p in range(12)]
         for i in range(3)},
        index=index,
    )
    total = wide.sum(axis=1)
    fig = stacked_figure(
        wide, [f"serie_{i}" for i in range(3)],
        title="Saldos", y_title="R$ bi", scale=0.001, total=total,
    )
    assert any("." in (t or "") and "," in (t or "") for t in fig.data[0].text)
    assert "." in fig.layout.annotations[0].text
    assert fig.layout.yaxis.tickformat == ","
    assert fig.layout.separators == ",."


# =============================================================================
# SENSIBILIDADE DO EIXO SECUNDÁRIO
# =============================================================================

def _combo(prazo: list[float]):
    index = pd.date_range("2025-08-31", periods=len(prazo), freq="ME")
    wide = pd.DataFrame(
        {
            "concessoes_livre_pj": [250_000.0 + i * 5_000 for i in range(len(prazo))],
            "prazo_livre_pj": prazo,
        },
        index=index,
    )
    return bar_line_figure(
        wide, bar_alias="concessoes_livre_pj", line_alias="prazo_livre_pj",
        title="Combo",
    )


def _ocupacao_vertical(fig, valores: list[float]) -> float:
    menor, maior = fig.layout.yaxis2.range
    return (max(valores) - min(valores)) / (maior - menor)


def test_variacao_pequena_de_prazo_aparece_pequena():
    """Veículos varia 1 mês em 47. Com escala automática, isso enchia o card.

    Ancorar o eixo secundário em zero faz a altura ocupada acompanhar a
    variação relativa, e não a amplitude absoluta.
    """
    quase_constante = [46.3 + 0.1 * i for i in range(11)] + [47.3]
    muito_variavel = [23.0 + i for i in range(11)] + [34.1]

    plano = _ocupacao_vertical(_combo(quase_constante), quase_constante)
    movido = _ocupacao_vertical(_combo(muito_variavel), muito_variavel)

    assert plano < 0.05, "variação de 2% não pode ocupar mais que 5% da altura"
    assert movido > 0.20, "variação de 40% precisa ser visível"
    assert movido > plano * 5


def test_eixo_secundario_comeca_em_zero_nos_dois_tipos_de_grafico():
    combo = _combo([24.0 + i * 0.4 for i in range(12)])
    assert combo.layout.yaxis.range[0] == 0
    assert combo.layout.yaxis2.range[0] == 0

    index = pd.date_range("2025-08-31", periods=12, freq="ME")
    wide = pd.DataFrame(
        {
            "comprometimento_juros": [10.0 + i * 0.1 for i in range(12)],
            "endividamento_renda": [49.0 + i * 0.08 for i in range(12)],
        },
        index=index,
    )
    linhas = line_figure(
        wide, ["comprometimento_juros", "endividamento_renda"],
        title="Situação", y_title="% da renda", compacto=True,
        secundarios=["endividamento_renda"],
    )
    assert linhas.layout.yaxis2.range[0] == 0


def test_rotulos_dos_dois_eixos_nao_colidem():
    """Endividamento, no eixo da direita, caía em cima dos rótulos da esquerda.

    Cada eixo era espalhado por si e nenhum enxergava o outro.
    """
    index = pd.date_range("2025-08-31", periods=12, freq="ME")
    wide = pd.DataFrame(
        {
            "comprometimento_amortizacao": [18.0] * 12,
            "comprometimento_juros": [10.9] * 12,
            "comprometimento_total_derivado": [28.8] * 12,
            "comprometimento_servico_ex_habitacional": [26.6] * 12,
            "endividamento_renda": [49.8] * 12,
        },
        index=index,
    )
    fig = line_figure(
        wide, list(wide.columns), title="Situação", y_title="% da renda",
        suffix="%", compacto=True, secundarios=["endividamento_renda"],
    )
    assert len(fig.layout.annotations) == 5

    altura = _altura_area_plotagem(fig)
    posicoes = []
    for anotacao in fig.layout.annotations:
        eixo = fig.layout.yaxis2 if anotacao.yref == "y2" else fig.layout.yaxis
        menor, maior = eixo.range
        posicoes.append((maior - anotacao.y) * altura / (maior - menor) + anotacao.ay)
    posicoes.sort()
    for anterior, seguinte in zip(posicoes, posicoes[1:]):
        assert seguinte - anterior >= TAMANHO_ROTULO_PX


def test_eixo_secundario_ganha_margem_propria():
    """Sem margem, a régua da direita era desenhada por cima dos dados."""
    index = pd.date_range("2025-08-31", periods=12, freq="ME")
    wide = pd.DataFrame(
        {"a": [10.0] * 12, "b": [50.0] * 12}, index=index
    )
    com = line_figure(wide, ["a", "b"], title="X", y_title="%", compacto=True,
                      secundarios=["b"])
    sem = line_figure(wide, ["a", "b"], title="X", y_title="%", compacto=True)
    assert com.layout.margin.r > sem.layout.margin.r


# =============================================================================
# RÓTULO EM TODOS OS PERÍODOS
# =============================================================================

def test_barra_empilhada_rotula_todos_os_periodos_por_padrao():
    wide = _wide(series=3)
    fig = stacked_figure(
        wide, [f"serie_{i}" for i in range(3)], title="X", y_title="R$ bi",
    )
    for trace in fig.data:
        assert sum(1 for texto in trace.text if texto) == len(wide)
    assert fig.layout.uniformtext.mode == "hide"
