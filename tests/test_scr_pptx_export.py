"""Testes do export dos painéis do SCR para PPTX.

O que se garante aqui é o contrato visual pedido: gráfico NATIVO do Office
(editável, não imagem), quatro painéis por slide em quadrantes iguais, rótulo
apenas no último período e percentual com duas casas no eixo e nos rótulos.

Os testes inspecionam o XML do arquivo gerado, não a API do python-pptx — é o
XML que o PowerPoint lê.
"""

from __future__ import annotations

from io import BytesIO

import pandas as pd
import pytest

pptx = pytest.importorskip("pptx")
from pptx import Presentation  # noqa: E402

from tabs import scr_inadimplencia as T  # noqa: E402
from utils import scr_pptx_export as X  # noqa: E402


def _painel(titulo="Inad (> 90 d) - Cheque especial", n_meses=6, series=None):
    series = series or {
        "Até 1 salário mínimo": [0.10, 0.11, 0.12, 0.13, 0.14, 0.15],
        "Acima de 20 salários mínimos": [0.04, 0.041, 0.042, 0.043, 0.044, 0.045],
        T.SERIE_TOTAL: [0.07, 0.072, 0.075, 0.078, 0.08, 0.082],
    }
    meses = [f"2026-{m:02d}" for m in range(1, n_meses + 1)]
    linhas = []
    for nome, valores in series.items():
        for mes, valor in zip(meses, valores[:n_meses]):
            linhas.append({
                "data_base": mes, "serie": nome, "valor": valor,
                "denominador": 1_000_000.0,
            })
    ordem = list(series)
    quebra = T.QUEBRAS_POR_KEY["renda"]
    return T.PainelSpec(
        titulo=titulo,
        subtitulo="Por Faixa Salário Mínimo - % carteira",
        fonte="fonte: Banco Central do Brasil",
        produto="Cheque especial",
        series=pd.DataFrame(linhas),
        ordem_series=ordem,
        cores=T.cores_das_series(ordem, quebra),
        tracejadas=[T.SERIE_TOTAL],
        metrica="inadimplencia",
        carteira_final_rs_mil=1_000_000.0,
    )


def _abrir(blob: bytes) -> Presentation:
    return Presentation(BytesIO(blob))


def _indices_dos_rotulos(chart) -> list:
    """Índice do ponto rotulado em cada série, na ordem das séries."""
    from pptx.oxml.ns import qn

    indices = []
    for serie in chart.plots[0].series:
        rotulos = serie._element.findall(".//" + qn("c:dLbl"))
        indices.extend(int(r.find(qn("c:idx")).get("val")) for r in rotulos)
    return indices


# =============================================================================
# ESTRUTURA
# =============================================================================

def test_export_gera_pptx_valido():
    blob, meta = X.exportar_paineis_pptx([_painel()])
    assert blob[:2] == b"PK"          # zip do Office
    assert meta["slides"] == 1
    assert _abrir(blob).slides


def test_export_rejeita_lista_vazia():
    with pytest.raises(ValueError):
        X.exportar_paineis_pptx([])


def test_quatro_paineis_por_slide():
    paineis = [_painel(titulo=f"Painel {i}") for i in range(9)]
    blob, meta = X.exportar_paineis_pptx(paineis)
    # 9 painéis => 3 slides (4 + 4 + 1)
    assert meta["slides"] == 3
    prs = _abrir(blob)
    graficos_por_slide = [
        sum(1 for shape in slide.shapes if shape.has_chart) for slide in prs.slides
    ]
    assert graficos_por_slide == [4, 4, 1]


def test_quadrantes_tem_tamanho_igual_e_nao_se_sobrepoem():
    blob, _ = X.exportar_paineis_pptx([_painel(titulo=f"P{i}") for i in range(4)])
    slide = _abrir(blob).slides[0]
    caixas = [
        (s.left, s.top, s.width, s.height) for s in slide.shapes if s.has_chart
    ]
    assert len(caixas) == 4
    larguras = {w for _, _, w, _ in caixas}
    alturas = {h for _, _, _, h in caixas}
    assert len(larguras) == 1, "quadrantes com larguras diferentes"
    assert len(alturas) == 1, "quadrantes com alturas diferentes"
    # Dois à esquerda e dois à direita, dois em cima e dois embaixo.
    assert len({l for l, _, _, _ in caixas}) == 2
    assert len({t for _, t, _, _ in caixas}) == 2


def test_slide_e_widescreen():
    blob, _ = X.exportar_paineis_pptx([_painel()])
    prs = _abrir(blob)
    assert round(prs.slide_width / prs.slide_height, 2) == pytest.approx(1.78, abs=0.01)


def test_paineis_cabem_dentro_do_slide():
    blob, _ = X.exportar_paineis_pptx([_painel(titulo=f"P{i}") for i in range(4)])
    prs = _abrir(blob)
    slide = prs.slides[0]
    for shape in slide.shapes:
        assert shape.left >= 0 and shape.top >= 0
        assert shape.left + shape.width <= prs.slide_width
        assert shape.top + shape.height <= prs.slide_height


# =============================================================================
# GRÁFICO NATIVO E EDITÁVEL
# =============================================================================

def test_grafico_e_nativo_do_office_nao_imagem():
    blob, _ = X.exportar_paineis_pptx([_painel()])
    chart = [s for s in _abrir(blob).slides[0].shapes if s.has_chart][0].chart
    xml = chart._chartSpace.xml
    assert "<c:lineChart>" in xml, "o gráfico precisa ser nativo para ser editável"
    # Os valores ficam embutidos, então o Office consegue reabrir a planilha.
    assert list(chart.plots[0].series[0].values)


def test_series_preservam_nome_curto_e_valores():
    blob, _ = X.exportar_paineis_pptx(
        [_painel()], rotulo_serie_fn=lambda n: T.rotulo_serie(n, T.QUEBRAS_POR_KEY["renda"])
    )
    chart = [s for s in _abrir(blob).slides[0].shapes if s.has_chart][0].chart
    nomes = [s.name for s in chart.plots[0].series]
    assert nomes == ["Até 1", "Acima 20", T.SERIE_TOTAL]
    assert list(chart.plots[0].series[0].values)[-1] == pytest.approx(0.15)


def test_categorias_sao_meses_abreviados():
    blob, _ = X.exportar_paineis_pptx([_painel()])
    chart = [s for s in _abrir(blob).slides[0].shapes if s.has_chart][0].chart
    assert list(chart.plots[0].categories) == [
        "jan.26", "fev.26", "mar.26", "abr.26", "mai.26", "jun.26"
    ]


@pytest.mark.parametrize("entrada,esperado", [
    ("2026-06", "jun.26"), ("2024-09", "set.24"), ("2012-07", "jul.12"),
])
def test_rotulo_mes(entrada, esperado):
    assert X.rotulo_mes(entrada) == esperado


# =============================================================================
# RÓTULO SÓ NO ÚLTIMO PERÍODO
# =============================================================================

def test_rotulo_aparece_apenas_no_ultimo_ponto():
    painel = _painel()
    blob, _ = X.exportar_paineis_pptx([painel])
    chart = [s for s in _abrir(blob).slides[0].shapes if s.has_chart][0].chart
    xml = chart._chartSpace.xml
    n_series = len(painel.ordem_series)
    # Um rótulo por série — e não um por ponto.
    assert xml.count("<c:dLbl>") == n_series
    assert xml.count('showVal val="1"') == n_series
    # E todos apontam para o último período (índice 5 de 6 categorias).
    assert _indices_dos_rotulos(chart) == [5] * n_series
    # O plot inteiro não pode ter rótulos ligados.
    assert not chart.plots[0].has_data_labels


def test_rotulo_usa_showval_e_nao_texto_fixo():
    # showVal mantém o rótulo ligado ao dado: editar a série no Office atualiza
    # o número. Texto literal congelaria o valor.
    blob, _ = X.exportar_paineis_pptx([_painel()])
    xml = [s for s in _abrir(blob).slides[0].shapes if s.has_chart][0].chart._chartSpace.xml
    assert 'showVal val="1"' in xml
    assert "<c:tx>" not in xml.split("<c:dLbl>")[1].split("</c:dLbl>")[0]


def test_rotulo_pula_cauda_sem_dados():
    # Série que termina antes das demais recebe o rótulo no último ponto VÁLIDO.
    painel = _painel(series={
        "Até 1 salário mínimo": [0.10, 0.11, 0.12, 0.13, 0.14, 0.15],
        "Acima de 20 salários mínimos": [0.04, 0.041, 0.042, None, None, None],
    })
    blob, _ = X.exportar_paineis_pptx([painel])
    chart = [s for s in _abrir(blob).slides[0].shapes if s.has_chart][0].chart
    assert chart._chartSpace.xml.count("<c:dLbl>") == 2
    # O índice do rótulo tem que ser lido dentro de cada série: o `<c:idx>` que
    # aparece primeiro no XML é o da própria série, não o do rótulo.
    assert _indices_dos_rotulos(chart) == [5, 2], (
        "a série completa rotula o índice 5; a que termina cedo, o último válido (2)"
    )


# =============================================================================
# FORMATO PERCENTUAL
# =============================================================================

def test_eixo_de_valores_em_percentual_com_duas_casas():
    blob, _ = X.exportar_paineis_pptx([_painel()])
    chart = [s for s in _abrir(blob).slides[0].shapes if s.has_chart][0].chart
    assert chart.value_axis.tick_labels.number_format == "0.00%"
    assert chart.value_axis.tick_labels.number_format_is_linked is False


def test_rotulos_em_percentual_com_duas_casas():
    blob, _ = X.exportar_paineis_pptx([_painel()])
    xml = [s for s in _abrir(blob).slides[0].shapes if s.has_chart][0].chart._chartSpace.xml
    # Um numFmt por rótulo mais o do eixo.
    assert xml.count('formatCode="0.00%"') >= 4


def test_valores_gravados_como_fracao():
    # O formato percentual do Office multiplica por 100: 0.15 vira 15,00%.
    # Gravar 15 produziria 1500,00%.
    blob, _ = X.exportar_paineis_pptx([_painel()])
    chart = [s for s in _abrir(blob).slides[0].shapes if s.has_chart][0].chart
    for serie in chart.plots[0].series:
        for valor in serie.values:
            if valor is not None:
                assert 0 <= valor <= 1


# =============================================================================
# PALETA E ESTILO
# =============================================================================

def test_cores_das_series_seguem_a_paleta_do_painel():
    painel = _painel()
    blob, _ = X.exportar_paineis_pptx([painel])
    chart = [s for s in _abrir(blob).slides[0].shapes if s.has_chart][0].chart
    for posicao, nome in enumerate(painel.ordem_series):
        cor = chart.plots[0].series[posicao].format.line.color.rgb
        assert str(cor) == painel.cores[nome].lstrip("#").upper()


def test_linha_todos_sai_tracejada_e_mais_grossa():
    painel = _painel()
    blob, _ = X.exportar_paineis_pptx([painel])
    chart = [s for s in _abrir(blob).slides[0].shapes if s.has_chart][0].chart
    idx_total = list(painel.ordem_series).index(T.SERIE_TOTAL)
    linha_total = chart.plots[0].series[idx_total].format.line
    linha_faixa = chart.plots[0].series[0].format.line
    assert linha_total.width > linha_faixa.width
    assert "dash" in chart._chartSpace.xml.lower()


def test_titulo_subtitulo_e_fonte_vao_para_o_slide():
    painel = _painel()
    blob, _ = X.exportar_paineis_pptx([painel])
    textos = [
        s.text_frame.text for s in _abrir(blob).slides[0].shapes if s.has_text_frame
    ]
    assert painel.titulo in textos
    assert painel.subtitulo in textos
    assert "fonte: Banco Central do Brasil" in textos


def test_titulo_do_deck_e_opcional():
    blob, _ = X.exportar_paineis_pptx([_painel()], titulo_deck="Deck de crédito")
    textos = [
        s.text_frame.text for s in _abrir(blob).slides[0].shapes if s.has_text_frame
    ]
    assert "Deck de crédito" in textos

    blob_sem, _ = X.exportar_paineis_pptx([_painel()])
    textos_sem = [
        s.text_frame.text for s in _abrir(blob_sem).slides[0].shapes if s.has_text_frame
    ]
    assert "Deck de crédito" not in textos_sem


def test_eixo_de_meses_rala_os_rotulos():
    # 36 meses num quadrante pequeno: sem intervalo, os rótulos empilham.
    painel = _painel(n_meses=6, series={
        "Até 1 salário mínimo": [0.1] * 6, T.SERIE_TOTAL: [0.08] * 6,
    })
    longo = painel.series.copy()
    blob, _ = X.exportar_paineis_pptx([painel])
    chart = [s for s in _abrir(blob).slides[0].shapes if s.has_chart][0].chart
    assert "tickLblSkip" in chart._chartSpace.xml
    assert not longo.empty
