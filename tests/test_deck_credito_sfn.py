"""Deck contínuo de Crédito SFN: as cinco abas em um arquivo.

O deck é deliberadamente simples: caixa de texto e gráfico nativo, nada mais.
Sem slide divisor, sem régua, sem moldura em volta do gráfico.
"""

from __future__ import annotations

import sys
from io import BytesIO
from pathlib import Path

import pandas as pd
import pytest
from pptx import Presentation
from pptx.oxml.ns import qn

import tabs.mercado_credito as MC
from utils.comentarios_credito import comentario
from utils.sgs_credit_analytics import derive_credit_totals, to_wide

sys.path.insert(0, str(Path(__file__).resolve().parent))
from test_credito_bc_legibilidade import (  # noqa: E402
    _graficos_do_deck,
    validar_xml_do_grafico,
)


PROJECT_ROOT = Path(__file__).resolve().parent.parent


@pytest.fixture(scope="module")
def janela() -> pd.DataFrame:
    bruto = pd.read_parquet(
        PROJECT_ROOT / "data" / "bundled" / "mercado_credito_sgs" / "dados.parquet"
    )
    wide = derive_credit_totals(to_wide(bruto))
    analiticas = [c for c in wide.columns if c not in {"cdi_aa", "selic_aa"}]
    periodos = pd.DatetimeIndex(
        wide[analiticas].dropna(how="all").index
    ).sort_values().unique()
    fim = pd.Timestamp(periodos[-1])
    recorte = wide.loc[(wide.index >= fim - pd.DateOffset(months=11)) & (wide.index <= fim)].copy()
    recorte.attrs["full_history"] = wide
    return recorte


@pytest.fixture(scope="module")
def deck(janela):
    blob, meta = MC._deck_credito_sfn(janela)
    return blob, meta, Presentation(BytesIO(blob))


def test_modo_silencioso_monta_as_figuras_sem_desenhar(janela):
    """As cinco abas entram no deck sem o usuário abrir uma por uma."""
    for titulo, _ in MC.SECOES_CREDITO_DECK:
        figuras = MC._figuras_da_subsecao(titulo, janela)
        assert figuras, f"{titulo} não produziu gráfico"
        assert all(figura.data for figura in figuras)
    # O modo silencioso não vaza para o render normal.
    assert MC._SILENCIOSO.get() is False


def test_deck_segue_a_ordem_das_abas_da_tela(deck):
    _, _, apresentacao = deck
    titulos = [
        proxima.text_frame.paragraphs[0].runs[0].text
        for slide in apresentacao.slides
        for proxima in slide.shapes
        if proxima.has_text_frame and proxima.text_frame.paragraphs[0].runs
    ]
    esperada = [titulo for titulo, _ in MC.SECOES_CREDITO_DECK]
    posicoes = [titulos.index(nome) for nome in esperada]
    assert posicoes == sorted(posicoes), "as abas saíram fora da ordem da tela"
    assert MC.CREDIT_SUBSECTIONS == tuple(esperada)


def test_cada_aba_tem_um_slide_de_leitura_antes_dos_graficos(deck):
    _, _, apresentacao = deck
    slides = list(apresentacao.slides)
    for titulo, chave in MC.SECOES_CREDITO_DECK:
        leitura = comentario(chave)
        primeiro_paragrafo = leitura.paragrafos[0]
        indice_texto = next(
            i for i, slide in enumerate(slides)
            if any(primeiro_paragrafo in forma.text for forma in slide.shapes if forma.has_text_frame)
        )
        assert not any(forma.has_chart for forma in slides[indice_texto].shapes)
        indice_grafico = next(
            i for i, slide in enumerate(slides)
            if i > indice_texto and any(forma.has_chart for forma in slide.shapes)
        )
        assert indice_grafico == indice_texto + 1


def test_deck_so_tem_caixa_de_texto_e_grafico(deck):
    """Nada decorativo: nem régua, nem moldura, nem shape de enfeite."""
    _, _, apresentacao = deck
    estranhos = [
        forma
        for slide in apresentacao.slides
        for forma in slide.shapes
        if not forma.has_text_frame and not forma.has_chart
    ]
    assert estranhos == []


def test_grafico_nao_tem_contorno(deck):
    """Sem spPr explícito, o Office desenha a moldura do estilo padrão."""
    blob, _, _ = deck
    for xml in _graficos_do_deck(blob):
        from lxml import etree

        raiz = etree.fromstring(xml)
        sp_pr = raiz.find(qn("c:spPr"))
        assert sp_pr is not None, "gráfico sem spPr herda a moldura do tema"
        assert len(sp_pr.findall(f".//{qn('a:noFill')}")) == 2


def test_deck_inteiro_passa_no_schema(deck):
    blob, _, _ = deck
    erros = [erro for xml in _graficos_do_deck(blob) for erro in validar_xml_do_grafico(xml)]
    assert erros == []


def test_dois_graficos_por_slide(deck):
    _, meta, apresentacao = deck
    assert meta["paineis_por_slide"] == 2
    for slide in apresentacao.slides:
        assert sum(1 for forma in slide.shapes if forma.has_chart) <= 2


def test_capa_traz_titulo_competencia_e_fonte(deck):
    _, _, apresentacao = deck
    capa = apresentacao.slides[0]
    textos = [forma.text for forma in capa.shapes if forma.has_text_frame]
    assert any("Crédito SFN" in texto for texto in textos)
    assert any("janela até" in texto for texto in textos)
    assert any("Banco Central" in texto for texto in textos)
    assert not any(forma.has_chart for forma in capa.shapes)


def test_comentario_sai_em_paragrafos_separados(deck):
    """Linha em branco no JSON vira parágrafo no slide, como na tela."""
    _, _, apresentacao = deck
    leitura = comentario("credito_estoque")
    assert len(leitura.paragrafos) == 2
    slide = next(
        slide for slide in apresentacao.slides
        if any(leitura.paragrafos[0] in f.text for f in slide.shapes if f.has_text_frame)
    )
    corpo = next(
        forma for forma in slide.shapes
        if forma.has_text_frame and leitura.paragrafos[0] in forma.text
    )
    assert len(corpo.text_frame.paragraphs) == len(leitura.paragrafos)


def test_todos_os_dezoito_graficos_de_credito_sfn_entram(deck):
    _, meta, _ = deck
    assert meta["secoes"] == len(MC.SECOES_CREDITO_DECK)
    assert meta["paineis"] == 18
