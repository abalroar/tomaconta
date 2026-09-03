from __future__ import annotations

import ast
from datetime import date
from io import BytesIO
from pathlib import Path

import pandas as pd
import plotly.graph_objects as go
import pytest

from utils.ifdata_cache.sgs_credit import SGSCreditCache
from utils.ifdata_cache.manager import CacheManager
from utils.sgs_credit_analytics import (
    bar_line_figure,
    build_ipca_index,
    coverage_ratio,
    derive_credit_totals,
    eixo_datas_semestral,
    line_figure,
    real_yoy,
    shares,
)
from utils.sgs_credit_pptx_export import exportar_figuras_pptx, figura_para_painel
from utils.sgs_credit_providers import BCBSGSProvider
from utils.sgs_credit_registry import SGS_SERIES, SGS_SERIES_BY_CODE, get_series
from tabs.mercado_credito import _filter_period_range, _real_growth_frame


PROJECT_ROOT = Path(__file__).resolve().parents[1]


def _long_frame(values: dict[str, list[float]]) -> pd.DataFrame:
    dates = pd.date_range("2024-01-31", periods=max(map(len, values.values())), freq="ME")
    rows = []
    for alias, series_values in values.items():
        spec = get_series(alias)
        for data, value in zip(dates, series_values):
            rows.append(
                {
                    "data": data,
                    "codigo": spec.code,
                    "serie": alias,
                    "nome_oficial": spec.official_name,
                    "valor": value,
                    "unidade": spec.unit,
                    "frequencia": spec.frequency,
                    "provedor": spec.provider,
                }
            )
    return pd.DataFrame(rows)


def test_registry_keeps_corrected_codes_and_unique_aliases():
    assert SGS_SERIES["saldo_controle_privado_nacional"].code == 12106
    assert SGS_SERIES["saldo_livre_total"].code == 20542
    assert SGS_SERIES["saldo_direcionado_total"].code == 20593
    assert SGS_SERIES["saldo_pj_total"].code == 20540
    assert SGS_SERIES["saldo_pf_total"].code == 20541
    assert SGS_SERIES["inad_total"].code == 21082
    assert len(SGS_SERIES_BY_CODE) == len({spec.code for spec in SGS_SERIES.values()})


def test_ipca_index_and_real_yoy_follow_workbook_formula():
    index = pd.date_range("2023-01-31", periods=13, freq="ME")
    ipca = pd.Series([1.0] * 13, index=index)
    nominal = pd.Series([100.0] * 12 + [112.682503], index=index)

    deflator = build_ipca_index(ipca)
    growth = real_yoy(nominal, deflator)

    assert round(deflator.iloc[-1], 6) == round(100 * 1.01**13, 6)
    assert abs(growth.iloc[-1]) < 0.001


def test_shares_and_derived_credit_total_preserve_missing_values():
    index = pd.date_range("2024-01-31", periods=2, freq="ME")
    wide = pd.DataFrame(
        {
            "saldo_livre_pj": [10.0, 11.0],
            "saldo_livre_pf": [20.0, 21.0],
            "saldo_direcionado_pj": [30.0, pd.NA],
            "saldo_direcionado_pf": [40.0, 42.0],
        },
        index=index,
    )
    derived = derive_credit_totals(wide)
    assert derived["saldo_sfn_total_derivado"].iloc[0] == 100.0
    assert pd.isna(derived["saldo_sfn_total_derivado"].iloc[1])

    mix = shares(derived, ["saldo_livre_pj", "saldo_livre_pf"], derived["saldo_sfn_total_derivado"])
    assert mix.iloc[0].sum() == 30.0
    assert mix.iloc[1].isna().all()


def test_coverage_is_not_capped_at_one_hundred_percent():
    ratio = coverage_ratio(pd.Series([7.5]), pd.Series([3.0]))
    assert ratio.iloc[0] == 250.0


def test_figures_label_only_last_valid_point_and_combo_has_secondary_axis():
    index = pd.date_range("2025-01-31", periods=3, freq="ME")
    wide = pd.DataFrame(
        {
            "taxa_pf_livre": [20.0, 21.0, 22.0],
            "concessoes_livre_pf_veiculos": [1000.0, 1100.0, 1200.0],
            "prazo_livre_pf_veiculos": [40.0, 41.0, 42.0],
        },
        index=index,
    )
    line = line_figure(wide, ["taxa_pf_livre"], title="Taxa", y_title="%")
    combo = bar_line_figure(
        wide,
        bar_alias="concessoes_livre_pf_veiculos",
        line_alias="prazo_livre_pf_veiculos",
        title="Veículos",
    )

    assert [annotation.text for annotation in line.layout.annotations] == ["22,0"]
    assert line.layout.annotations[0].font.color == line.data[0].line.color
    assert len(combo.data) == 2
    assert combo.data[1].yaxis == "y2"
    assert combo.layout.annotations[0].font.color == combo.data[1].line.color
    assert line.layout.meta["chart_title"] == "Taxa"


def test_eixo_escolhe_passo_adaptativo_e_inclui_ultimo_mes_disponivel():
    datas = pd.date_range("2025-01-31", "2026-07-31", freq="ME")
    valores, rotulos = eixo_datas_semestral(datas)

    assert [data.strftime("%Y-%m") for data in valores] == [
        "2025-01", "2025-03", "2025-05", "2025-07", "2025-09",
        "2025-11", "2026-01", "2026-03", "2026-05", "2026-07",
    ]
    assert rotulos == [
        "Jan/25", "Mar/25", "Mai/25", "Jul/25", "Set/25",
        "Nov/25", "Jan/26", "Mar/26", "Mai/26", "Jul/26",
    ]

    fig = line_figure(
        pd.DataFrame({"taxa_pf_livre": range(len(datas))}, index=datas),
        ["taxa_pf_livre"], title="Taxa", y_title="%",
    )
    assert fig.layout.xaxis.tickangle == 0


def test_eixo_nao_repete_mes_quando_series_fecham_em_dias_distintos():
    valores, rotulos = eixo_datas_semestral([
        "2025-12-29", "2025-12-31", "2026-06-29", "2026-06-30", "2026-07-30"
    ])

    assert [data.strftime("%Y-%m-%d") for data in valores] == [
        "2025-12-31", "2026-06-30", "2026-07-30"
    ]
    assert rotulos == ["Dez/25", "Jun/26", "Jul/26"]


def test_export_sgs_gera_grafico_office_nativo_com_titulo_e_ultimo_rotulo():
    from pptx import Presentation

    index = pd.date_range("2026-01-31", periods=6, freq="ME")
    fig = line_figure(
        pd.DataFrame({"taxa_pf_livre": range(20, 26)}, index=index),
        ["taxa_pf_livre"], title="Taxa PF", y_title="% a.a.", suffix="%",
    )
    blob, meta = exportar_figuras_pptx([fig], titulo_deck="Teste")
    deck = Presentation(BytesIO(blob))
    charts = [shape.chart for shape in deck.slides[0].shapes if shape.has_chart]
    textos = [shape.text for shape in deck.slides[0].shapes if shape.has_text_frame]

    assert blob[:2] == b"PK"
    assert meta["paineis"] == 1
    assert len(charts) == 1
    assert "Taxa PF" in textos
    assert "undefined" not in blob.decode("latin-1", errors="ignore").lower()


def test_export_sgs_aceita_categorias_de_uf_e_rotula_todas_as_barras():
    from pptx import Presentation

    fig = go.Figure(go.Bar(x=["SP", "RJ", "BA"], y=[50.0, 30.0, 20.0], name="Participação"))
    fig.update_layout(meta={
        "chart_title": "Participação por UF",
        "value_format": "0.0%",
        "value_scale": 0.01,
        "label_all_points": True,
        "source": "fonte: Banco Central do Brasil · SCR.data",
    })

    painel = figura_para_painel(fig)
    assert painel.fonte.endswith("SCR.data")

    blob, _ = exportar_figuras_pptx([fig], titulo_deck="Teste")
    deck = Presentation(BytesIO(blob))
    chart = [shape.chart for shape in deck.slides[0].shapes if shape.has_chart][0]

    assert list(chart.plots[0].categories) == ["SP", "RJ", "BA"]
    assert all(point.data_label is not None for point in chart.series[0].points)


def test_line_labels_are_staggered_and_keep_series_order():
    index = pd.date_range("2025-01-31", periods=3, freq="ME")
    wide = pd.DataFrame(
        {
            "taxa_pf_livre": [20.0, 20.9, 21.00],
            "taxa_pj_livre": [19.0, 20.8, 20.99],
            "spread_pf_livre": [18.0, 20.7, 20.98],
        },
        index=index,
    )
    fig = line_figure(
        wide,
        ["taxa_pf_livre", "taxa_pj_livre", "spread_pf_livre"],
        title="Linhas próximas",
        y_title="%",
    )

    assert len(fig.layout.annotations) == 3
    offsets = [annotation.ay for annotation in fig.layout.annotations]
    assert len(set(offsets)) == 3
    assert [annotation.font.color for annotation in fig.layout.annotations] == [
        trace.line.color for trace in fig.data
    ]


def test_latest_period_keeps_each_series_own_last_observation():
    index = pd.date_range("2026-05-31", periods=3, freq="ME")
    wide = pd.DataFrame(
        {"serie_julho": [1.0, 2.0, 3.0], "serie_junho": [4.0, 5.0, pd.NA]},
        index=index,
    )

    filtered = _filter_period_range(wide, pd.Timestamp("2026-05-31"), None)

    assert filtered["serie_julho"].last_valid_index() == pd.Timestamp("2026-07-31")
    assert filtered["serie_junho"].last_valid_index() == pd.Timestamp("2026-06-30")
    assert filtered.attrs["full_history"] is wide


def test_selected_start_keeps_prior_history_for_yoy_calculation():
    index = pd.date_range("2025-01-31", periods=13, freq="ME")
    wide = pd.DataFrame(
        {
            "ipca_mensal": [0.0] * 13,
            "credito_ampliado_total": [100.0] * 12 + [110.0],
        },
        index=index,
    )
    selected = _filter_period_range(wide, index[-1], None)

    growth = _real_growth_frame(selected, ["credito_ampliado_total"])

    assert growth.iloc[0, 0] == pytest.approx(10.0)


class _FakeResponse:
    status_code = 200
    text = ""

    def json(self):
        return [{"data": "31/01/2024", "valor": "1,25"}]


class _FakeSession:
    def get(self, url, params, timeout):
        assert "bcdata.sgs.433" in url
        assert params["dataInicial"] == "01/01/2024"
        return _FakeResponse()


def test_bcb_provider_normalizes_decimal_comma_and_date():
    provider = BCBSGSProvider(session=_FakeSession())
    frame = provider.fetch(get_series("ipca_mensal"), date(2024, 1, 1), date(2024, 1, 31))
    assert frame.iloc[0]["valor"] == 1.25
    assert frame.iloc[0]["data"] == pd.Timestamp("2024-01-31")


class _FlakyJsonResponse(_FakeResponse):
    def __init__(self, valid):
        self.valid = valid

    def json(self):
        if not self.valid:
            raise ValueError("empty response")
        return super().json()


class _FlakyJsonSession:
    def __init__(self):
        self.calls = 0

    def get(self, url, params, timeout):
        self.calls += 1
        return _FlakyJsonResponse(valid=self.calls > 1)


def test_bcb_provider_retries_empty_json_response(monkeypatch):
    monkeypatch.setattr("utils.sgs_credit_providers.time.sleep", lambda _seconds: None)
    session = _FlakyJsonSession()
    provider = BCBSGSProvider(session=session, max_attempts=2)
    frame = provider.fetch(get_series("ipca_mensal"), date(2024, 1, 1), date(2024, 1, 31))
    assert session.calls == 2
    assert frame.iloc[0]["valor"] == 1.25


class _FakeProvider:
    def fetch(self, spec, start, end):
        return pd.DataFrame({"data": [pd.Timestamp("2024-01-31")], "valor": [10.0]})


def test_cache_materializes_long_parquet_with_injected_provider(tmp_path):
    cache = SGSCreditCache(tmp_path, providers={"bcb_sgs": _FakeProvider()})
    result = cache.materialize_history(
        start="2024-01-01",
        end="2024-01-31",
        aliases=["saldo_livre_total", "saldo_direcionado_total"],
        overwrite=True,
        max_workers=2,
    )
    assert result.sucesso is True
    assert cache.arquivo_dados.exists()
    stored = pd.read_parquet(cache.arquivo_dados)
    assert set(stored["serie"]) == {"saldo_livre_total", "saldo_direcionado_total"}
    assert set(stored["provedor"]) == {"bcb_sgs"}


def test_remote_cache_download_uses_cache_buster(monkeypatch, tmp_path):
    frame = _long_frame({"ipca_mensal": [0.5]})
    payload = BytesIO()
    frame.to_parquet(payload, index=False)
    calls = []

    class _ReleaseResponse:
        status_code = 200

        def __init__(self, content=b""):
            self.content = content

        def json(self):
            return {"series": 1, "total_periodos": 1}

    def _fake_get(url, timeout, headers):
        calls.append((url, timeout, headers))
        if "_dados.parquet" in url:
            return _ReleaseResponse(payload.getvalue())
        return _ReleaseResponse()

    monkeypatch.setattr("utils.ifdata_cache.sgs_credit.requests.get", _fake_get)
    result = SGSCreditCache(tmp_path).baixar_remoto()
    assert result.sucesso is True
    assert len(result.dados) == 1
    assert len(calls) == 2
    assert all("cache_bust=" in url for url, _, _ in calls)
    assert all(headers["Cache-Control"] == "no-cache" for _, _, headers in calls)


def test_app_registers_single_sgs_menu_and_dispatch_route():
    source = (PROJECT_ROOT / "app1.py").read_text(encoding="utf-8")
    ast.parse(source)
    assert source.count('"Estatísticas Crédito BC"') >= 2
    assert source.count('elif menu == "Estatísticas Crédito BC":') == 1
    assert 'def _get_sgs_credit_cache(' in source
    assert 'render_mercado_credito(_get_sgs_credit_cache(), get_cache_manager=get_cache_manager)' in source
    assert 'manager.registrar(cache)' in source
    assert '"mercado_credito_sgs": "Estatísticas Crédito BC (BCData/SGS)' in source
    assert "cache_sgs_update.materialize_history(" in source


def test_credit_module_has_period_range_info_popovers_and_no_expectations_registry():
    source = (PROJECT_ROOT / "tabs" / "mercado_credito.py").read_text(encoding="utf-8")
    assert '"Período inicial"' in source
    assert '"Período final"' in source
    assert '"Mais recente"' in source
    assert 'st.popover("i"' in source
    assert "Expectativas do Mercado de Crédito" not in source
    assert "Registry de séries SGS" not in source
    assert "DateOffset(months=11)" in source
    assert '"Baixar PPTX desta aba"' in source
    assert "reversed(periods)" in source
    assert "formatar_competencia" in source
    assert "border-radius: 50%" in source
    assert '"SCR.data"' in source
    assert "_render_source_footer" not in source
    assert 'fig.update_layout(title_text=""' in source
    assert "fig.update_layout(title=None" not in source


def test_navigation_has_dedicated_bcb_group_and_no_top_level_scr():
    source = (PROJECT_ROOT / "app1.py").read_text(encoding="utf-8")
    tree = ast.parse(source)
    assignments = {
        target.id: ast.literal_eval(node.value)
        for node in tree.body
        if isinstance(node, ast.Assign)
        for target in node.targets
        if isinstance(target, ast.Name) and target.id in {"MENU_PRINCIPAL", "MENU_BCB"}
    }
    assert assignments["MENU_BCB"] == ["Estatísticas Crédito BC", "Taxas de Juros por Produto"]
    assert "Inadimplência (SCR)" not in assignments["MENU_PRINCIPAL"]
    assert "Estatísticas Crédito BC" not in assignments["MENU_PRINCIPAL"]
    assert "Taxas de Juros por Produto" not in assignments["MENU_PRINCIPAL"]


def test_taxas_por_produto_defaults_to_12_months_and_exports_pptx():
    source = (PROJECT_ROOT / "app1.py").read_text(encoding="utf-8")
    assert 'st.session_state["tj_beta_janela_meses"] = min(12, max_meses_beta)' in source
    assert 'file_name="taxas_juros_por_produto.pptx"' in source
    assert 'key="tj_beta_download_pptx"' in source


def test_cache_manager_registers_sgs_cache(tmp_path):
    manager = CacheManager(base_dir=tmp_path)
    assert isinstance(manager.get_cache("mercado_credito_sgs"), SGSCreditCache)
