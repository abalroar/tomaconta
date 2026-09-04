"""Deck contínuo: todas as abas de Estatísticas Crédito BC em um arquivo.

O deck é deliberadamente simples: caixa de texto e gráfico nativo, nada mais.
Sem slide divisor, sem régua, sem moldura em volta do gráfico.
"""

from __future__ import annotations

import sys
from io import BytesIO
from pathlib import Path

import pandas as pd
import pytest
from pptx import Presentation
from pptx.oxml.ns import qn

import tabs.mercado_credito as MC
from utils.comentarios_credito import comentario
from utils.sgs_credit_analytics import derive_credit_totals, to_wide

sys.path.insert(0, str(Path(__file__).resolve().parent))
from test_credito_bc_legibilidade import (  # noqa: E402
    _NS,
    _graficos_do_deck,
    validar_xml_do_grafico,
)


PROJECT_ROOT = Path(__file__).resolve().parent.parent


@pytest.fixture(scope="module")
def janela() -> pd.DataFrame:
    bruto = pd.read_parquet(
        PROJECT_ROOT / "data" / "bundled" / "mercado_credito_sgs" / "dados.parquet"
    )
    wide = derive_credit_totals(to_wide(bruto))
    analiticas = [c for c in wide.columns if c not in {"cdi_aa", "selic_aa"}]
    periodos = pd.DatetimeIndex(
        wide[analiticas].dropna(how="all").index
    ).sort_values().unique()
    fim = pd.Timestamp(periodos[-1])
    recorte = wide.loc[(wide.index >= fim - pd.DateOffset(months=11)) & (wide.index <= fim)].copy()
    recorte.attrs["full_history"] = wide
    return recorte


@pytest.fixture(scope="module")
def gerenciador():
    from utils.ifdata_cache import CacheManager

    manager = CacheManager()
    return lambda: manager


@pytest.fixture(scope="module")
def deck(janela, gerenciador):
    blob, meta = MC._deck_completo(janela, gerenciador)
    return blob, meta, Presentation(BytesIO(blob))


def test_modo_silencioso_monta_as_figuras_sem_desenhar(janela):
    """Todas as abas entram no deck sem o usuário abrir uma por uma."""
    for titulo, _, render in MC.SECOES_DECK:
        if render is None:
            continue
        figuras = MC._figuras_da_secao(render, janela)
        assert figuras, f"{titulo} não produziu gráfico"
        assert all(figura.data for figura in figuras)
    # O modo silencioso não vaza para o render normal.
    assert MC._SILENCIOSO.get() is False


def test_deck_segue_a_ordem_das_abas_da_tela(deck):
    _, _, apresentacao = deck
    titulos = [
        forma.text_frame.paragraphs[0].runs[0].text
        for slide in apresentacao.slides
        for forma in slide.shapes
        if forma.has_text_frame and forma.text_frame.paragraphs[0].runs
    ]
    esperada = [titulo for titulo, _, _ in MC.SECOES_DECK]
    posicoes = [titulos.index(nome) for nome in esperada]
    assert posicoes == sorted(posicoes), "as abas saíram fora da ordem da tela"
    # A ordem do deck é a ordem do menu da seção.
    assert [t for t, _, _ in MC.SECOES_DECK][0] == MC.MAIN_SECTIONS[0]
    assert MC.CREDIT_SUBSECTIONS == tuple(
        t.split(" · ")[-1] for t, _, _ in MC.SECOES_DECK if t.startswith("Crédito SFN")
    )


def test_leitura_de_cada_aba_fica_acima_dos_graficos_dela(deck):
    """O comentário abre a aba, na mesma lâmina dos primeiros gráficos.

    Comentário longo demais para a faixa ganha lâmina própria; nesse caso os
    gráficos vêm na lâmina seguinte.
    """
    _, _, apresentacao = deck
    slides = list(apresentacao.slides)
    for _, chave, _ in MC.SECOES_DECK:
        leitura = comentario(chave)
        primeiro_paragrafo = leitura.paragrafos[0]
        indice = next(
            i for i, slide in enumerate(slides)
            if any(
                primeiro_paragrafo in forma.text
                for forma in slide.shapes if forma.has_text_frame
            )
        )
        com_grafico = [
            forma for forma in slides[indice].shapes if forma.has_chart
        ]
        if com_grafico:
            topo_comentario = min(
                forma.top for forma in slides[indice].shapes
                if forma.has_text_frame and primeiro_paragrafo in forma.text
            )
            assert topo_comentario < min(forma.top for forma in com_grafico)
        else:
            assert any(forma.has_chart for forma in slides[indice + 1].shapes)


def test_deck_so_tem_caixa_de_texto_e_grafico(deck):
    """Nada decorativo: nem régua, nem moldura, nem shape de enfeite."""
    _, _, apresentacao = deck
    estranhos = [
        forma
        for slide in apresentacao.slides
        for forma in slide.shapes
        if not forma.has_text_frame and not forma.has_chart
    ]
    assert estranhos == []


def test_grafico_nao_tem_contorno(deck):
    """Sem spPr explícito, o Office desenha a moldura do estilo padrão."""
    blob, _, _ = deck
    for xml in _graficos_do_deck(blob):
        from lxml import etree

        raiz = etree.fromstring(xml)
        sp_pr = raiz.find(qn("c:spPr"))
        assert sp_pr is not None, "gráfico sem spPr herda a moldura do tema"
        assert len(sp_pr.findall(f".//{qn('a:noFill')}")) == 2


def test_deck_inteiro_passa_no_schema(deck):
    blob, _, _ = deck
    erros = [erro for xml in _graficos_do_deck(blob) for erro in validar_xml_do_grafico(xml)]
    assert erros == []


def test_grade_de_ate_quatro_graficos_por_slide(deck):
    _, meta, apresentacao = deck
    assert meta["paineis_por_slide"] == 4
    for slide in apresentacao.slides:
        assert sum(1 for forma in slide.shapes if forma.has_chart) <= 4


def test_capa_traz_titulo_competencia_e_fonte(deck):
    _, _, apresentacao = deck
    capa = apresentacao.slides[0]
    textos = [forma.text for forma in capa.shapes if forma.has_text_frame]
    assert any("Estatísticas Crédito BC" in texto for texto in textos)
    assert any("janela até" in texto for texto in textos)
    assert any("Banco Central" in texto for texto in textos)
    assert not any(forma.has_chart for forma in capa.shapes)


def test_comentario_sai_em_paragrafos_separados(deck):
    """Linha em branco no JSON vira parágrafo no slide, como na tela."""
    _, _, apresentacao = deck
    leitura = comentario("credito_estoque")
    assert len(leitura.paragrafos) == 2
    slide = next(
        slide for slide in apresentacao.slides
        if any(leitura.paragrafos[0] in f.text for f in slide.shapes if f.has_text_frame)
    )
    corpo = next(
        forma for forma in slide.shapes
        if forma.has_text_frame and leitura.paragrafos[0] in forma.text
    )
    assert len(corpo.text_frame.paragraphs) == len(leitura.paragrafos)


def test_todas_as_abas_e_submenus_entram(deck):
    _, meta, _ = deck
    assert meta["secoes"] == len(MC.SECOES_DECK)
    # 42 gráficos do SGS mais os da aba de faixa de renda, que traz todas as
    # modalidades PF e a visão regional.
    assert meta["paineis"] >= 48


# =============================================================================
# FORMATAÇÃO PEDIDA
# =============================================================================

def _paragrafos_de_comentario(apresentacao):
    primeiro = comentario("concessoes").paragrafos[0]
    for slide in apresentacao.slides:
        for forma in slide.shapes:
            if forma.has_text_frame and primeiro in forma.text:
                return forma.text_frame.paragraphs
    raise AssertionError("comentário não encontrado no deck")


def test_comentario_sai_com_entrelinha_de_um_e_meio_e_sem_espaco_extra(deck):
    _, _, apresentacao = deck
    paragrafos = _paragrafos_de_comentario(apresentacao)
    assert len(paragrafos) >= 2
    for paragrafo in paragrafos:
        assert paragrafo.line_spacing == 1.5
        assert paragrafo.space_before.pt == 0
        assert paragrafo.space_after.pt == 0


def test_titulo_e_comentario_ficam_proximos(deck):
    """Antes o comentário vivia em outra lâmina, a 0,80 in do título."""
    from utils.scr_pptx_export import ALTURA_TITULO_SECAO, RESPIRO_TITULO

    _, _, apresentacao = deck
    primeiro = comentario("concessoes").paragrafos[0]
    slide = next(
        slide for slide in apresentacao.slides
        if any(primeiro in f.text for f in slide.shapes if f.has_text_frame)
    )
    corpo = next(f for f in slide.shapes if f.has_text_frame and primeiro in f.text)
    titulo = next(
        f for f in slide.shapes
        if f.has_text_frame and f.text.strip() == "Concessões"
    )
    vao = corpo.top - (titulo.top + titulo.height)
    assert 0 <= vao <= int(RESPIRO_TITULO) + 1
    assert titulo.height == ALTURA_TITULO_SECAO


def test_grade_preenche_a_altura_livre_do_slide(deck):
    """A célula cresce para ocupar o que sobra abaixo do comentário.

    Altura fixa deixava metade do slide vazia sempre que a faixa de leitura era
    curta. Dentro de um mesmo slide todos os gráficos continuam iguais.
    """
    _, _, apresentacao = deck
    limite = apresentacao.slide_height - 347472
    for slide in apresentacao.slides:
        graficos = [forma for forma in slide.shapes if forma.has_chart]
        if not graficos:
            continue
        assert len({forma.height for forma in graficos}) == 1
        base = max(forma.top + forma.height for forma in graficos)
        # A última linha encosta na margem inferior: sem sobra desperdiçada.
        assert limite - base <= 91440, "sobrou mais de 0,1 in de espaço vazio"
        assert base <= limite + 1


def test_nenhum_grafico_passa_da_margem_inferior(deck):
    _, _, apresentacao = deck
    limite = apresentacao.slide_height - 347472  # 7,5 in menos a margem de 0,38
    for slide in apresentacao.slides:
        for forma in slide.shapes:
            if forma.has_chart:
                assert forma.top + forma.height <= limite + 1


def test_rotulo_traz_nome_da_serie_junto_do_valor(deck):
    """O rótulo do último ponto identifica a linha sozinho, como na tela."""
    from lxml import etree

    blob, _, _ = deck
    encontrou = False
    for xml in _graficos_do_deck(blob):
        raiz = etree.fromstring(xml)
        # Só os rótulos ponto a ponto; o bloco de padrão do gráfico continua
        # desligado, e é ele que vale para os pontos sem rótulo.
        for grupo in raiz.iter():
            tipo = grupo.tag.replace(_NS, "")
            if tipo not in {"barChart", "lineChart"}:
                continue
            # Na coluna o nome não cabe na largura da barra e fica na legenda;
            # na linha ele vai junto do valor, na faixa da direita.
            esperado = "0" if tipo == "barChart" else "1"
            for rotulo in grupo.iter(qn("c:dLbl")):
                nome = rotulo.find(qn("c:showSerName"))
                valor = rotulo.find(qn("c:showVal"))
                assert valor is not None and valor.get("val") == "1"
                assert nome is not None and nome.get("val") == esperado
                encontrou = True
    assert encontrou, "nenhum rótulo configurado no deck"


def test_faixa_de_renda_traz_todas_as_modalidades_pf(gerenciador):
    """O deck leva as 7 modalidades PF, e não as 4 que a tela abre."""
    from tabs import scr_inadimplencia as scr_spec

    figuras = MC._figuras_faixa_de_renda(gerenciador)
    titulos = [
        (fig.layout.meta or {}).get("chart_title", "") for fig in figuras
    ]
    modalidades_pf = [
        modalidade for modalidade in scr_spec.MODALIDADES_BCB_PF
        if any(modalidade.split(" - ")[-1] in titulo for titulo in titulos)
    ]
    assert len(figuras) > scr_spec.PAINEIS_POR_SLIDE
    assert len(modalidades_pf) >= 5
    # A visão regional entra como barra por UF, já que mapa não é nativo.
    assert any("por UF" in titulo for titulo in titulos)
    assert any("por região" in titulo for titulo in titulos)


def test_mapa_nao_entra_no_deck(gerenciador):
    """Coroplético não existe como gráfico nativo do Office."""
    figuras = MC._figuras_faixa_de_renda(gerenciador)
    for figura in figuras:
        for trace in figura.data:
            assert trace.type not in {"choropleth", "choroplethmapbox", "scattergeo"}


# =============================================================================
# LEGIBILIDADE DOS RÓTULOS NO SLIDE RENDERIZADO
# =============================================================================

def test_rotulo_de_linha_fica_na_faixa_a_direita_da_plotagem(deck):
    """A área de plotagem encolhe e o rótulo do último ponto vai para a sobra.

    Sem isso o rótulo era desenhado sobre as próprias linhas e cortado na
    borda do gráfico.
    """
    from lxml import etree

    from pptx import Presentation
    from utils.scr_pptx_export import (
        GUTTER_ROTULO_LINHA,
        LARGURA_MINIMA_FAIXA_IN,
    )

    blob, _, apresentacao = deck
    conferidos = 0
    for slide in Presentation(BytesIO(blob)).slides:
        for forma in slide.shapes:
            if not forma.has_chart:
                continue
            raiz = forma.chart._chartSpace
            if raiz.find(f".//{qn('c:lineChart')}") is None:
                continue
            layout = raiz.find(
                f".//{qn('c:plotArea')}/{qn('c:layout')}/{qn('c:manualLayout')}"
            )
            assert layout is not None, "gráfico de linhas sem layout manual"
            esquerda = float(layout.find(qn("c:x")).get("val"))
            largura = float(layout.find(qn("c:w")).get("val"))
            faixa = 1.0 - (esquerda + largura)
            # A faixa é medida em polegadas: cabe o maior nome de série, com um
            # piso que não some e um teto que não come o gráfico.
            assert faixa * forma.width / 914400.0 >= LARGURA_MINIMA_FAIXA_IN * 0.95
            assert faixa <= GUTTER_ROTULO_LINHA + 0.001
            conferidos += 1
    assert conferidos, "nenhum gráfico de linhas no deck"


def test_rotulos_finais_nao_se_sobrepoem_na_faixa():
    """O escalonamento respeita a altura de cada nome, inclusive os que quebram."""
    from utils.scr_pptx_export import _escalonar_no_gutter

    finais = [
        ("Desconto de duplicatas/recebíveis", -3.6), ("Capital de giro", -5.4),
        ("Conta garantida", -2.5), ("Aquisição de bens", -9.7), ("ACC", -3.3),
        ("Financiamento à exportação", -10.3), ("Rural PJ", 22.4),
        ("Imobiliário PJ", 8.7), ("BNDES PJ", 5.5),
    ]
    posicoes = _escalonar_no_gutter(
        finais, int(2.5 * 914400), int(6.14 * 914400)
    )
    from utils.scr_pptx_export import FATOR_LINHA_CALIBRI, corpo_do_rotulo

    assert len(posicoes) == len(finais)
    # A folga exigida é a altura do rótulo de cima, não um número fixo: um nome
    # que quebra em duas linhas precisa do dobro do espaço do vizinho de baixo.
    linha = (corpo_do_rotulo(len(finais)) * FATOR_LINHA_CALIBRI / 72) / 2.5
    longos = {"Desconto de duplicatas/recebíveis", "Financiamento à exportação"}
    ordenadas = sorted(posicoes.items(), key=lambda item: item[1])
    for (nome, alto), (_, baixo) in zip(ordenadas, ordenadas[1:]):
        minimo = linha * (2 if nome in longos else 1)
        assert baixo - alto >= minimo * 0.98, f"{nome} invade o rótulo de baixo"
    assert ordenadas[0][1] >= 0.03
    assert ordenadas[-1][1] <= 0.95


def test_barra_empilhada_tem_barra_larga_para_o_rotulo_caber(deck):
    """O PowerPoint recorta o rótulo pela largura da barra."""
    from lxml import etree

    blob, _, _ = deck
    for xml in _graficos_do_deck(blob):
        raiz = etree.fromstring(xml)
        barra = raiz.find(f".//{qn('c:barChart')}")
        if barra is None or raiz.find(f".//{qn('c:lineChart')}") is not None:
            continue
        vao = barra.find(qn("c:gapWidth"))
        assert vao is not None and int(vao.get("val")) <= 50


def test_legenda_so_existe_onde_o_rotulo_nao_traz_o_nome(deck):
    """Linha não leva legenda: o nome já vai no rótulo e sobrava disputa de
    espaço com os meses do eixo."""
    from lxml import etree

    blob, _, _ = deck
    for xml in _graficos_do_deck(blob):
        raiz = etree.fromstring(xml)
        tem_legenda = raiz.find(f".//{qn('c:legend')}") is not None
        so_linhas = (
            raiz.find(f".//{qn('c:lineChart')}") is not None
            and raiz.find(f".//{qn('c:barChart')}") is None
        )
        if so_linhas:
            assert not tem_legenda


def test_eixo_secundario_do_deck_comeca_em_zero(deck):
    """Mesma sensibilidade da tela: prazo que varia pouco aparece plano."""
    from lxml import etree

    blob, _, _ = deck
    achou = False
    for xml in _graficos_do_deck(blob):
        raiz = etree.fromstring(xml)
        eixos = raiz.findall(f".//{qn('c:valAx')}")
        if len(eixos) < 2:
            continue
        minimo = eixos[-1].find(f"{qn('c:scaling')}/{qn('c:min')}")
        assert minimo is not None and float(minimo.get("val")) == 0
        achou = True
    assert achou, "nenhum card de dois eixos no deck"


def test_escala_de_eixo_usa_degrau_redondo_e_ancora_no_zero():
    """A régua do eixo é calculada aqui, e não deixada ao Office.

    O rótulo do último ponto é posicionado por coordenada: se a escala do
    gráfico não for a mesma que a conta do rótulo usa, o nome da série cai
    longe da linha que ele nomeia.
    """
    from utils.scr_pptx_export import escala_de_eixo

    # Série que não fica achatada pelo zero começa no zero.
    assert escala_de_eixo([3.6, 5.8, 6.6]) == (0.0, 7.0)
    # Série estreita no alto da escala mantém o piso, senão vira uma reta.
    piso, teto = escala_de_eixo([430.2, 453.1, 436.2])
    assert piso >= 400 and teto >= 453.1
    # Valor negativo puxa o piso para baixo do zero.
    piso, teto = escala_de_eixo([-10.3, 22.4])
    assert piso <= -10.3 and teto >= 22.4
    # Eixo secundário é sempre ancorado no zero.
    assert escala_de_eixo([0.42, 0.94], base_zero=True)[0] == 0.0


def test_rotulo_fica_na_altura_da_propria_regua():
    """Duas séries em eixos diferentes, cada uma medida na escala do seu eixo.

    Normalizar as duas pelo intervalo dos valores finais mandava o rótulo do
    eixo da direita para o rodapé, ao lado de uma linha que corre no topo.
    """
    from utils.scr_pptx_export import _escalonar_no_gutter

    posicoes = _escalonar_no_gutter(
        [("Capital de giro", 5.8), ("Conta garantida", 4.9),
         ("Desconto de recebíveis", 0.9)],
        int(3.4 * 914400), int(12.0 * 914400),
        secundarias=["Desconto de recebíveis"],
        escalas={False: (0.0, 6.0), True: (0.0, 1.0)},
    )
    # 0,9 de 1,0 é o topo do eixo da direita: o rótulo acompanha a linha.
    assert posicoes["Desconto de recebíveis"] < posicoes["Conta garantida"]
    assert posicoes["Desconto de recebíveis"] < 0.2


def test_eixo_secundario_tem_teto_fixo_e_fio_claro(deck):
    """Teto explícito para casar com o rótulo; fio cinza como o eixo da esquerda."""
    from lxml import etree

    blob, _, _ = deck
    achou = False
    for xml in _graficos_do_deck(blob):
        raiz = etree.fromstring(xml)
        eixos = raiz.findall(f".//{qn('c:valAx')}")
        if len(eixos) < 2:
            continue
        maximo = eixos[-1].find(f"{qn('c:scaling')}/{qn('c:max')}")
        assert maximo is not None and float(maximo.get("val")) > 0
        cor = eixos[-1].find(
            f"{qn('c:spPr')}/{qn('a:ln')}/{qn('a:solidFill')}/{qn('a:srgbClr')}"
        )
        assert cor is not None and cor.get("val") == "E6E6E6"
        achou = True
    assert achou, "nenhum card de dois eixos no deck"


def test_barra_com_muitas_categorias_gira_o_rotulo(deck):
    """27 UFs deixam a barra mais estreita que o número, e o Office o corta."""
    from lxml import etree

    blob, _, _ = deck
    achou = False
    for xml in _graficos_do_deck(blob):
        raiz = etree.fromstring(xml)
        if raiz.find(f".//{qn('c:barChart')}") is None:
            continue
        primeira = raiz.find(f".//{qn('c:ser')}")
        categorias = primeira.findall(f"{qn('c:cat')}//{qn('c:pt')}")
        if len(categorias) <= 14:
            continue
        corpos = raiz.findall(
            f".//{qn('c:dLbl')}/{qn('c:txPr')}/{qn('a:bodyPr')}"
        )
        assert corpos, "barra densa sem rótulo"
        assert all(no.get("rot") == "-5400000" for no in corpos)
        achou = True
    assert achou, "nenhuma barra com muitas categorias no deck"


def test_serie_unica_nao_leva_legenda(deck):
    """Com uma série só o Office lista as categorias na legenda, o que é ruído."""
    from lxml import etree

    blob, _, _ = deck
    for xml in _graficos_do_deck(blob):
        raiz = etree.fromstring(xml)
        series = raiz.findall(f".//{qn('c:ser')}")
        if len(series) == 1:
            assert raiz.find(f".//{qn('c:legend')}") is None
