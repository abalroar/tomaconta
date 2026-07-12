"""
Gera apresentação PowerPoint MVP TomaConta — estilo Itaú BBA.
Cores: Preto #1A1A1A, Laranja #F06611, Branco/Bege #FFFFFF / #FAF6F1
Fonte: Calibri
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

OUTPUT = "TomaConta_MVP_Status.pptx"

# ── Paleta ──────────────────────────────────────────────────────────────────
BLACK   = RGBColor(0x1A, 0x1A, 0x1A)
ORANGE  = RGBColor(0xF0, 0x66, 0x11)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
BEIGE   = RGBColor(0xFA, 0xF6, 0xF1)
DGRAY   = RGBColor(0x40, 0x40, 0x40)
LGRAY   = RGBColor(0xCC, 0xCC, 0xCC)

# Status colors
GREEN_FILL  = RGBColor(0xC6, 0xEF, 0xCE)
GREEN_FONT  = RGBColor(0x37, 0x56, 0x23)
YELLOW_FILL = RGBColor(0xFF, 0xEB, 0x9C)
YELLOW_FONT = RGBColor(0x9C, 0x65, 0x00)
RED_FILL    = RGBColor(0xFF, 0xC7, 0xCE)
RED_FONT    = RGBColor(0x9C, 0x00, 0x06)
GRAY_FILL   = RGBColor(0xED, 0xED, 0xED)
GRAY_FONT   = RGBColor(0x63, 0x63, 0x63)

STATUS_COLORS = {
    "Implementado":              (GREEN_FILL,  GREEN_FONT,  "✅"),
    "Parcialmente Implementado": (YELLOW_FILL, YELLOW_FONT, "🟡"),
    "Pendente":                  (RED_FILL,    RED_FONT,    "🔴"),
    "Não classificável":         (GRAY_FILL,   GRAY_FONT,   "⚠️"),
}

# 16:9
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]  # completely blank


def add_rect(slide, x, y, w, h, fill=None, line=None):
    shape = slide.shapes.add_shape(
        pptx.enum.shapes.MSO_SHAPE_TYPE.AUTO_SHAPE,
        x, y, w, h
    )
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, x, y, w, h, text, size=12, bold=False,
                color=BLACK, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Calibri"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb


def add_shape_text(shape, text, size=11, bold=False, color=BLACK,
                   align=PP_ALIGN.LEFT, wrap=True):
    tf = shape.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Calibri"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — CAPA
# ════════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(blank_layout)

# Full black background
bg = add_rect(s1, 0, 0, W, H, fill=BLACK)

# Orange bottom stripe
add_rect(s1, 0, H - Inches(1.1), W, Inches(1.1), fill=ORANGE)

# Title
add_textbox(s1, Inches(0.7), Inches(1.5), Inches(12), Inches(1.6),
            "TOMA CONTA", size=60, bold=True, color=ORANGE,
            align=PP_ALIGN.LEFT)

# Subtitle
add_textbox(s1, Inches(0.7), Inches(3.2), Inches(10), Inches(0.9),
            "Às quantas anda o MVP?", size=26, bold=False, color=WHITE,
            align=PP_ALIGN.LEFT)

# Date
add_textbox(s1, Inches(0.7), Inches(4.1), Inches(8), Inches(0.6),
            "Abril 2026  ·  Demo com o Diretor: 15/04/2026",
            size=14, color=LGRAY, align=PP_ALIGN.LEFT)

# Bottom stripe text
add_textbox(s1, Inches(0.7), H - Inches(1.0), Inches(10), Inches(0.8),
            "Corporativo | Interno", size=13, bold=False, color=WHITE,
            align=PP_ALIGN.LEFT)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — VISÃO GERAL
# ════════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(blank_layout)
add_rect(s2, 0, 0, W, H, fill=BEIGE)

# Top bar
add_rect(s2, 0, 0, W, Inches(1.0), fill=BLACK)
add_textbox(s2, Inches(0.4), Inches(0.15), Inches(12), Inches(0.7),
            "VISÃO GERAL DO MVP", size=22, bold=True, color=ORANGE,
            align=PP_ALIGN.LEFT)

# Completion big number
add_textbox(s2, Inches(0.4), Inches(1.2), Inches(3), Inches(2),
            "50%", size=80, bold=True, color=ORANGE, align=PP_ALIGN.LEFT)

add_textbox(s2, Inches(0.4), Inches(3.1), Inches(3.5), Inches(0.5),
            "de completion", size=16, bold=False, color=DGRAY,
            align=PP_ALIGN.LEFT)

add_textbox(s2, Inches(0.4), Inches(3.65), Inches(5), Inches(0.45),
            "(15 completos + ½ × 2 parciais) / 32 itens",
            size=11, color=DGRAY, align=PP_ALIGN.LEFT)

# Progress bar
bar_x, bar_y = Inches(0.4), Inches(4.2)
bar_w, bar_h  = Inches(5.0), Inches(0.32)
add_rect(s2, bar_x, bar_y, bar_w, bar_h, fill=LGRAY, line=LGRAY)
add_rect(s2, bar_x, bar_y, int(bar_w * 0.50), bar_h, fill=ORANGE)

# KPI boxes
kpis = [
    ("15", "Implementados", GREEN_FILL,  GREEN_FONT),
    ("2",  "Parcialmente",  YELLOW_FILL, YELLOW_FONT),
    ("15", "Pendentes",     RED_FILL,    RED_FONT),
    ("4",  "N/A ou obs.",   GRAY_FILL,   GRAY_FONT),
]
bx, by, bw, bh = Inches(6.2), Inches(1.3), Inches(1.6), Inches(1.3)
gap = Inches(1.9)
for i, (num, label, fc, fo) in enumerate(kpis):
    box = add_rect(s2, bx + i * gap, by, bw, bh, fill=fc, line=LGRAY)
    add_textbox(s2, bx + i * gap + Inches(0.1), by + Inches(0.05),
                bw - Inches(0.2), Inches(0.75),
                num, size=36, bold=True, color=fo, align=PP_ALIGN.CENTER)
    add_textbox(s2, bx + i * gap + Inches(0.05), by + Inches(0.75),
                bw - Inches(0.1), Inches(0.45),
                label, size=12, bold=False, color=fo, align=PP_ALIGN.CENTER)

# Prioridades demo
add_rect(s2, Inches(6.2), Inches(3.0), Inches(6.7), Inches(3.8),
         fill=WHITE, line=ORANGE)

add_textbox(s2, Inches(6.4), Inches(3.1), Inches(6.3), Inches(0.5),
            "PRIORIDADES PARA A DEMO — 15/04",
            size=12, bold=True, color=ORANGE, align=PP_ALIGN.LEFT)

prios = [
    "🔴  S2  — Corrigir bug 'variação de bastão' no Scatter",
    "🔴  P2/P3/P4  — Métricas crédito e N1 no Peer",
    "🔴  P1 / D1  — Dropdown base de comparação (Peer + DRE)",
    "🔴  R8  — Tratar variação % sobre dado já percentual",
    "🔴  D2/D3  — Reorganizar estrutura da DRE",
    "🔴  F1/F2  — Melhorias de filtros no FGC",
]
for i, p in enumerate(prios):
    add_textbox(s2, Inches(6.4), Inches(3.65) + i * Inches(0.50),
                Inches(6.3), Inches(0.45),
                p, size=11, color=DGRAY, align=PP_ALIGN.LEFT)

# Bottom bar
add_rect(s2, 0, H - Inches(0.4), W, Inches(0.4), fill=ORANGE)
add_textbox(s2, Inches(0.4), H - Inches(0.38), Inches(10), Inches(0.35),
            "Corporativo | Interno", size=10, color=WHITE)


# ════════════════════════════════════════════════════════════════════════════
# Helper: slide com tabela de pedidos por grupo de abas
# ════════════════════════════════════════════════════════════════════════════

ROWS_DATA = [
    ("Rankings",              "R1",    "Implementado",              "Multi-período no gráfico"),
    ("Rankings",              "R2",    "Implementado",              "ROE Trimestral como indicador"),
    ("Rankings",              "R3",    "Implementado",              "Campo 'Ponderar média por' ausente"),
    ("Rankings",              "R4",    "Implementado",              "Linhas pontilhadas média seleção + SFN"),
    ("Rankings",              "R5",    "Implementado",              "Pools Top 5/10/20 pré-selecionados"),
    ("Rankings",              "R6",    "Implementado",              "Campo 'Ordenação' ausente; só direção"),
    ("Rankings",              "R7",    "Implementado",              "Deltas: tri vs tri + acumulado vs acumulado"),
    ("Rankings",              "R8",    "Pendente",                  "Tratar variação % sobre dado em % (ROE/Basileia)"),
    ("Rankings",              "R_exp", "Pendente",                  "Remover export da visão gráfica (Deltas)"),
    ("Rankings",              "R13",   "Implementado",              "Tabela com pools Top N"),
    ("Rankings",              "R14",   "Implementado",              "Encaixe de banco individual no ranking"),
    ("Rankings",              "R15",   "Implementado",              "Export da tabela disponível"),
    ("Rankings",              "R16",   "Pendente",                  "Formatação números redondos na tabela"),
    ("Taxa de Juros",         "T1",    "Implementado",              "Ranking Top 10 aberto com posições"),
    ("FGC",                   "F1",    "Pendente",                  "Top N como clusters no dropdown de instituições"),
    ("FGC",                   "F2",    "Pendente",                  "Múltiplos períodos + toggle tri vs acumulado"),
    ("FGC",                   "F3",    "Pendente",                  "LL Trimestral como variável no FGC"),
    ("Scatter",               "S1",    "Implementado",              "LL Trim disponível para tamanho bolinha"),
    ("Scatter",               "S2",    "Pendente",                  "Bug: gráfico variação de bastão não renderiza"),
    ("Scatter",               "S3",    "Pendente",                  "Variáveis avançadas (backlog v2)"),
    ("DRE",                   "D1",    "Parcialmente Implementado", "Dropdown base lucro existe, falta tri vs acum."),
    ("DRE",                   "D2",    "Pendente",                  "Reorganizar DRE + Res. Intermediação Ajustado"),
    ("DRE",                   "D3",    "Pendente",                  "Desp PDD Outras Operações como subconta"),
    ("DRE",                   "D4",    "Pendente",                  "Análise aderência PDD/NI → levar ao Peer"),
    ("DRE Individual",        "DI1",   "Implementado",              "Exibe nome do banco (não código)"),
    ("Peer",                  "P1",    "Pendente",                  "Dropdown base de comparação"),
    ("Peer",                  "P2",    "Pendente",                  "Est3/Carteira e Est2+3/Carteira"),
    ("Peer",                  "P3",    "Pendente",                  "Perda Esperada / Estágio 2 e 3"),
    ("Peer",                  "P4",    "Pendente",                  "Índice Capital N1 (T1) na tabela"),
    ("Evolução",              "E1",    "Implementado",              "N1 incluso; ordem CET1→N1→BIS"),
    ("Carteira 4966",         "C1",    "Parcialmente Implementado", "Em construção (baixa prioridade)"),
]

# Group into 3 slides (~10-11 rows each)
SLIDE_GROUPS = [
    ("Rankings & Taxa de Juros",        ROWS_DATA[:14]),
    ("FGC, Scatter & DRE",              ROWS_DATA[14:24]),
    ("Peer, Evolução & Carteira 4966",  ROWS_DATA[24:]),
]

COL_WIDTHS = [Inches(1.4), Inches(0.55), Inches(2.2), Inches(8.8)]
COL_LABELS = ["Aba", "#", "Status", "Pedido / Descrição"]

for title, group in SLIDE_GROUPS:
    sl = prs.slides.add_slide(blank_layout)
    add_rect(sl, 0, 0, W, H, fill=BEIGE)

    # Header
    add_rect(sl, 0, 0, W, Inches(0.85), fill=BLACK)
    add_textbox(sl, Inches(0.35), Inches(0.1), Inches(12), Inches(0.65),
                title.upper(), size=20, bold=True, color=ORANGE)

    # Column headers
    row_h = Inches(0.38)
    header_y = Inches(0.9)
    cx = Inches(0.25)
    for cw, cl in zip(COL_WIDTHS, COL_LABELS):
        hdr = add_rect(sl, cx, header_y, cw, row_h, fill=BLACK)
        add_shape_text(hdr, cl, size=10, bold=True, color=WHITE,
                       align=PP_ALIGN.CENTER)
        cx += cw

    # Data rows
    for ri, (aba, num, status, desc) in enumerate(group):
        ry = Inches(0.9) + row_h * (ri + 1) + Inches(0.04) * ri
        fill_c, font_c, icon = STATUS_COLORS.get(status, (BEIGE, DGRAY, ""))

        cx = Inches(0.25)
        vals = [aba, num, f"{icon} {status}", desc]
        for cw, val in zip(COL_WIDTHS, vals):
            is_status = val.startswith(("✅", "🟡", "🔴", "⚠️"))
            bg = fill_c if is_status else WHITE
            fc = font_c if is_status else DGRAY
            cell_shape = add_rect(sl, cx, ry, cw, row_h, fill=bg, line=LGRAY)
            add_shape_text(cell_shape, val, size=9,
                           bold=is_status,
                           color=fc,
                           align=PP_ALIGN.CENTER if is_status else PP_ALIGN.LEFT)
            cx += cw

    # Bottom bar
    add_rect(sl, 0, H - Inches(0.4), W, Inches(0.4), fill=ORANGE)
    add_textbox(sl, Inches(0.35), H - Inches(0.38), Inches(6), Inches(0.35),
                "Corporativo | Interno", size=10, color=WHITE)
    add_textbox(sl, Inches(9), H - Inches(0.38), Inches(4), Inches(0.35),
                "Completion: 50%  |  15 impl · 2 parcial · 15 pendente",
                size=9, color=WHITE, align=PP_ALIGN.RIGHT)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE FINAL — Próximos Passos
# ════════════════════════════════════════════════════════════════════════════
sf = prs.slides.add_slide(blank_layout)
add_rect(sf, 0, 0, W, H, fill=BLACK)
add_rect(sf, 0, 0, W, Inches(1.0), fill=ORANGE)
add_textbox(sf, Inches(0.4), Inches(0.12), Inches(12), Inches(0.75),
            "PRÓXIMOS PASSOS PARA A DEMO — 15/04/2026",
            size=21, bold=True, color=WHITE)

steps = [
    ("ALTA PRIORIDADE",  [
        "S2  — Corrigir bug 'variação de bastão' no Scatter Plot",
        "P2/P3/P4  — Adicionar métricas de crédito e Índice N1 ao Peer",
        "P1 / D1  — Dropdown 'base de comparação' no Peer e na DRE",
        "R8  — Tratar corretamente variação % sobre dado já percentual",
    ], ORANGE, RED_FILL, RED_FONT),
    ("MÉDIA PRIORIDADE", [
        "D2 / D3  — Reorganizar variáveis da DRE e abrir subconta PDD",
        "F1 / F2  — Melhorar filtros da aba Contribuições FGC",
        "R16  — Formatação de números na tabela do Ranking",
    ], LGRAY, YELLOW_FILL, YELLOW_FONT),
    ("BACKLOG (PÓS-DEMO)", [
        "S3  — Variáveis avançadas no Scatter (PDD/Cart, Est3/Cart, …)",
        "D4  — Análise de aderência PDD/NI → Peer (condicional)",
        "R9-R12  — Seção ilegível — realinhar com o chefe",
    ], DGRAY, GRAY_FILL, GRAY_FONT),
]

col_x = [Inches(0.35), Inches(4.7), Inches(9.0)]
for ci, (title_s, items, title_col, bg_col, fg_col) in enumerate(steps):
    tx = col_x[ci]
    add_textbox(sf, tx, Inches(1.15), Inches(4.1), Inches(0.45),
                title_s, size=12, bold=True, color=title_col)
    for ii, item in enumerate(items):
        box = add_rect(sf, tx, Inches(1.65) + ii * Inches(1.22),
                       Inches(4.1), Inches(1.12), fill=bg_col, line=LGRAY)
        add_shape_text(box, item, size=10, color=fg_col, wrap=True)

add_rect(sf, 0, H - Inches(0.4), W, Inches(0.4), fill=ORANGE)
add_textbox(sf, Inches(0.4), H - Inches(0.38), Inches(10), Inches(0.35),
            "Corporativo | Interno", size=10, color=WHITE)

prs.save(OUTPUT)
print(f"✅ Salvo: {OUTPUT}  ({prs.slides.__len__()} slides)")
