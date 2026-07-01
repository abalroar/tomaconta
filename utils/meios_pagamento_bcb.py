"""Dados e visualizações de meios de pagamento do BCB.

Fonte: serviço OData MPV_DadosAbertos do Banco Central do Brasil.
"""

from __future__ import annotations

import io
from datetime import datetime
from typing import Any, Iterable, Mapping, Optional, Sequence

import pandas as pd
import plotly.graph_objects as go
import requests


BCB_MPV_BASE_URL = "https://olinda.bcb.gov.br/olinda/servico/MPV_DadosAbertos/versao/v1/odata"
BCB_MPV_DATASET_URL = "https://dadosabertos.bcb.gov.br/dataset/estatisticas-meios-pagamentos"
BCB_MPV_SITE_URL = "https://www.bcb.gov.br/estatisticas/spbadendos"
BCB_MPV_DOC_URL = "https://olinda.bcb.gov.br/olinda/servico/MPV_DadosAbertos/versao/v1/documentacao"

MONTHLY_RESOURCE = "MeiosdePagamentosMensalDA"
QUARTERLY_RESOURCE = "MeiosdePagamentosTrimestralDA"

DEFAULT_MONTHLY_START = "200501"
DEFAULT_QUARTERLY_START = "20051"

MONTH_NAMES_PT = {
    1: "jan",
    2: "fev",
    3: "mar",
    4: "abr",
    5: "mai",
    6: "jun",
    7: "jul",
    8: "ago",
    9: "set",
    10: "out",
    11: "nov",
    12: "dez",
}

MONTHLY_MODALITIES = ["Pix", "TED", "Boleto", "Cheque", "TEC", "DOC"]
QUARTERLY_MODALITIES = [
    "Pix",
    "TED",
    "Boleto",
    "Cheque",
    "Cartão de Crédito",
    "Cartão de Débito",
    "Cartão Pré-pago",
    "Transferência Intrabancária",
    "Convênios",
    "Débito Direto",
    "Saques",
    "TEC",
    "DOC",
]

MODALITY_FIELD_SUFFIX = {
    "Pix": "Pix",
    "TED": "TED",
    "TEC": "TEC",
    "Cheque": "Cheque",
    "Boleto": "Boleto",
    "DOC": "DOC",
    "Cartão de Crédito": "CartaoCredito",
    "Cartão de Débito": "CartaoDebito",
    "Cartão Pré-pago": "CartaoPrePago",
    "Transferência Intrabancária": "TransIntrabancaria",
    "Convênios": "Convenios",
    "Débito Direto": "DebitoDireto",
    "Saques": "Saques",
}

ITAU_BBA_PALETTE = {
    "Pix": "#EC7000",
    "TED": "#000000",
    "Boleto": "#6B6B6B",
    "Cheque": "#A6A6A6",
    "TEC": "#FDB913",
    "DOC": "#7A4D9E",
    "Cartão de Crédito": "#004B8D",
    "Cartão de Débito": "#2F80B7",
    "Cartão Pré-pago": "#79B8D1",
    "Transferência Intrabancária": "#404040",
    "Convênios": "#B7B7B7",
    "Débito Direto": "#4F2E7F",
    "Saques": "#FFC857",
}


def _normalize_decimal(value: Any) -> float | None:
    if value is None:
        return None
    try:
        return float(value)
    except (TypeError, ValueError):
        text = str(value).strip().replace(".", "").replace(",", ".")
        if not text:
            return None
        try:
            return float(text)
        except ValueError:
            return None


def _bcb_odata_function(resource: str, parameter_name: str, parameter_value: str, *, timeout: int = 45) -> pd.DataFrame:
    endpoint = f"{BCB_MPV_BASE_URL}/{resource}({parameter_name}=@{parameter_name})"
    params = {
        f"@{parameter_name}": f"'{parameter_value}'",
        "$format": "json",
    }
    response = requests.get(endpoint, params=params, timeout=timeout)
    response.raise_for_status()
    payload = response.json()
    values = payload.get("value", [])
    return pd.DataFrame(values)


def fetch_meios_pagamento_mensal(start_yyyymm: str = DEFAULT_MONTHLY_START, *, timeout: int = 45) -> pd.DataFrame:
    """Busca dados mensais a partir de AAAAMM."""
    return _bcb_odata_function(MONTHLY_RESOURCE, "AnoMes", str(start_yyyymm), timeout=timeout)


def fetch_meios_pagamento_trimestral(start_yyyyq: str = DEFAULT_QUARTERLY_START, *, timeout: int = 45) -> pd.DataFrame:
    """Busca dados trimestrais a partir de AAAAT."""
    return _bcb_odata_function(QUARTERLY_RESOURCE, "trimestre", str(start_yyyyq), timeout=timeout)


def parse_monthly_period(value: Any) -> pd.Timestamp:
    text = str(value).strip()
    if len(text) != 6 or not text.isdigit():
        return pd.NaT
    return pd.Timestamp(year=int(text[:4]), month=int(text[4:6]), day=1)


def parse_quarterly_period(value: Any) -> pd.Timestamp:
    text = str(value).strip()
    if not text:
        return pd.NaT
    parsed = pd.to_datetime(text, errors="coerce")
    if pd.notna(parsed):
        return pd.Timestamp(parsed).normalize()
    if len(text) == 5 and text.isdigit():
        year = int(text[:4])
        quarter = int(text[4])
        month = {1: 3, 2: 6, 3: 9, 4: 12}.get(quarter)
        if month:
            return pd.Timestamp(year=year, month=month, day=1) + pd.offsets.MonthEnd(0)
    return pd.NaT


def format_monthly_axis_label(dt: Any) -> str:
    ts = pd.Timestamp(dt)
    if pd.isna(ts):
        return ""
    return f"{ts.month:02d}-{ts.year}"


def format_quarterly_axis_label(dt: Any) -> str:
    ts = pd.Timestamp(dt)
    if pd.isna(ts):
        return ""
    return f"{MONTH_NAMES_PT.get(ts.month, ts.strftime('%b').lower())}-{str(ts.year)[-2:]}"


def normalize_meios_pagamento(
    raw: pd.DataFrame,
    *,
    periodicidade: str,
) -> pd.DataFrame:
    """Converte a base wide do BCB para formato longo analítico."""
    if raw is None or raw.empty:
        return pd.DataFrame(
            columns=[
                "periodicidade",
                "periodo",
                "periodo_codigo",
                "periodo_label",
                "modalidade",
                "metrica",
                "valor",
                "unidade",
            ]
        )

    periodicidade_norm = str(periodicidade).strip().lower()
    if periodicidade_norm.startswith("mens"):
        period_col = "AnoMes"
        parser = parse_monthly_period
        labeler = format_monthly_axis_label
        modalities = MONTHLY_MODALITIES
    elif periodicidade_norm.startswith("trim"):
        period_col = "datatrimestre"
        parser = parse_quarterly_period
        labeler = format_quarterly_axis_label
        modalities = QUARTERLY_MODALITIES
    else:
        raise ValueError("periodicidade deve ser 'mensal' ou 'trimestral'.")

    frame = raw.copy()
    if period_col not in frame.columns:
        return pd.DataFrame()

    frame["periodo"] = frame[period_col].map(parser)
    frame = frame[pd.notna(frame["periodo"])].copy()
    frame["periodo_codigo"] = frame[period_col].astype(str)
    frame["periodo_label"] = frame["periodo"].map(labeler)

    records: list[dict[str, Any]] = []
    for _, row in frame.iterrows():
        for modality in modalities:
            suffix = MODALITY_FIELD_SUFFIX[modality]
            for metric, unit in (("quantidade", "mil transações"), ("valor", "R$ milhão")):
                field = f"{metric}{suffix}"
                if field not in frame.columns:
                    continue
                records.append(
                    {
                        "periodicidade": "Mensal" if periodicidade_norm.startswith("mens") else "Trimestral",
                        "periodo": row["periodo"],
                        "periodo_codigo": row["periodo_codigo"],
                        "periodo_label": row["periodo_label"],
                        "modalidade": modality,
                        "metrica": "Quantidade" if metric == "quantidade" else "Valor",
                        "valor": _normalize_decimal(row.get(field)),
                        "unidade": unit,
                    }
                )

    long_df = pd.DataFrame.from_records(records)
    if long_df.empty:
        return long_df
    long_df["valor"] = pd.to_numeric(long_df["valor"], errors="coerce")
    long_df = long_df.sort_values(["periodo", "modalidade", "metrica"]).reset_index(drop=True)
    return long_df


def load_meios_pagamento_bcb(
    *,
    monthly_start: str = DEFAULT_MONTHLY_START,
    quarterly_start: str = DEFAULT_QUARTERLY_START,
    timeout: int = 45,
) -> tuple[pd.DataFrame, pd.DataFrame]:
    monthly_raw = fetch_meios_pagamento_mensal(monthly_start, timeout=timeout)
    quarterly_raw = fetch_meios_pagamento_trimestral(quarterly_start, timeout=timeout)
    monthly = normalize_meios_pagamento(monthly_raw, periodicidade="mensal")
    quarterly = normalize_meios_pagamento(quarterly_raw, periodicidade="trimestral")
    return monthly, quarterly


def filter_period_range(df: pd.DataFrame, start: Any, end: Any) -> pd.DataFrame:
    if df is None or df.empty:
        return pd.DataFrame(columns=getattr(df, "columns", []))
    start_ts = pd.Timestamp(start)
    end_ts = pd.Timestamp(end)
    return df[(df["periodo"] >= start_ts) & (df["periodo"] <= end_ts)].copy()


def available_periods(df: pd.DataFrame) -> list[pd.Timestamp]:
    if df is None or df.empty or "periodo" not in df.columns:
        return []
    return sorted(pd.to_datetime(df["periodo"].dropna().unique()).tolist())


def compute_share(df: pd.DataFrame, *, metric: str, modalities: Sequence[str]) -> pd.DataFrame:
    if df is None or df.empty:
        return pd.DataFrame(columns=["periodo", "periodo_label", "modalidade", "participacao"])
    subset = df[
        (df["metrica"].eq(metric))
        & (df["modalidade"].isin(list(modalities)))
        & (pd.to_numeric(df["valor"], errors="coerce").fillna(0) >= 0)
    ].copy()
    if subset.empty:
        return pd.DataFrame(columns=["periodo", "periodo_label", "modalidade", "participacao"])
    totals = subset.groupby("periodo")["valor"].transform("sum")
    subset["participacao"] = subset["valor"] / totals.replace(0, pd.NA)
    return subset.dropna(subset=["participacao"]).copy()


def _format_short_number(value: Any, *, percent: bool = False) -> str:
    number = _normalize_decimal(value)
    if number is None or pd.isna(number):
        return "n.d."
    if percent:
        return f"{number * 100:.1f}%".replace(".", ",")
    abs_value = abs(number)
    if abs_value >= 1_000_000:
        return f"{number / 1_000_000:.1f} mi".replace(".", ",")
    if abs_value >= 1_000:
        return f"{number / 1_000:.1f} mil".replace(".", ",")
    return f"{number:,.0f}".replace(",", "X").replace(".", ",").replace("X", ".")


def build_line_figure(
    df: pd.DataFrame,
    *,
    metric: str,
    modalities: Sequence[str],
    title: str,
    unit_label: str,
    periodicidade: str,
    height: int = 470,
) -> go.Figure:
    fig = go.Figure()
    subset = df[(df["metrica"].eq(metric)) & (df["modalidade"].isin(list(modalities)))].copy()
    subset = subset.sort_values(["periodo", "modalidade"])
    for modality in modalities:
        series = subset[subset["modalidade"].eq(modality)].sort_values("periodo")
        if series.empty:
            continue
        values = pd.to_numeric(series["valor"], errors="coerce")
        texts = [""] * len(series)
        if len(texts):
            texts[-1] = _format_short_number(values.iloc[-1])
        color = ITAU_BBA_PALETTE.get(modality, "#666666")
        fig.add_trace(
            go.Scatter(
                x=series["periodo"],
                y=values,
                mode="lines+markers+text",
                name=modality,
                line={"color": color, "width": 2.6},
                marker={"color": color, "size": 5},
                text=texts,
                textposition="middle right",
                textfont={"color": color, "size": 12},
                cliponaxis=False,
                hovertemplate="%{fullData.name}<br>%{x|%d/%m/%Y}<br>%{y:,.2f}<extra></extra>",
            )
        )

    _style_time_series_figure(fig, title=title, yaxis_title=unit_label, periodicidade=periodicidade, height=height)
    return fig


def build_share_figure(
    df: pd.DataFrame,
    *,
    metric: str,
    modalities: Sequence[str],
    title: str,
    periodicidade: str,
    height: int = 430,
) -> go.Figure:
    share = compute_share(df, metric=metric, modalities=modalities)
    fig = go.Figure()
    for modality in modalities:
        series = share[share["modalidade"].eq(modality)].sort_values("periodo")
        if series.empty:
            continue
        values = pd.to_numeric(series["participacao"], errors="coerce")
        texts = [""] * len(series)
        if len(texts):
            texts[-1] = _format_short_number(values.iloc[-1], percent=True)
        color = ITAU_BBA_PALETTE.get(modality, "#666666")
        fig.add_trace(
            go.Scatter(
                x=series["periodo"],
                y=values,
                mode="lines+markers+text",
                name=modality,
                line={"color": color, "width": 2.4},
                marker={"color": color, "size": 5},
                text=texts,
                textposition="middle right",
                textfont={"color": color, "size": 12},
                cliponaxis=False,
                hovertemplate="%{fullData.name}<br>%{x|%d/%m/%Y}<br>%{y:.1%}<extra></extra>",
            )
        )
    _style_time_series_figure(
        fig,
        title=title,
        yaxis_title="Participação no total selecionado",
        periodicidade=periodicidade,
        height=height,
        percent_axis=True,
    )
    return fig


def _style_time_series_figure(
    fig: go.Figure,
    *,
    title: str,
    yaxis_title: str,
    periodicidade: str,
    height: int,
    percent_axis: bool = False,
) -> None:
    fig.update_layout(
        title={"text": title, "x": 0.0, "xanchor": "left", "font": {"size": 18, "color": "#1F1F1F"}},
        template="plotly_white",
        height=height,
        margin={"l": 64, "r": 112, "t": 60, "b": 72},
        legend={"orientation": "h", "yanchor": "bottom", "y": 1.02, "xanchor": "left", "x": 0},
        font={"family": "Arial, sans-serif", "color": "#1F1F1F"},
        hovermode="x unified",
        paper_bgcolor="white",
        plot_bgcolor="white",
    )
    fig.update_yaxes(
        title=yaxis_title,
        tickformat=".0%" if percent_axis else ",.0f",
        gridcolor="#E5E5E5",
        zerolinecolor="#B7B7B7",
    )
    if str(periodicidade).lower().startswith("trim"):
        all_x = []
        for trace in fig.data:
            all_x.extend(list(trace.x))
        ticks = sorted(pd.to_datetime(pd.Series(all_x).dropna().unique()).tolist()) if all_x else []
        fig.update_xaxes(
            tickmode="array",
            tickvals=ticks,
            ticktext=[format_quarterly_axis_label(tick) for tick in ticks],
            tickangle=-35,
            gridcolor="#F2F2F2",
            title=None,
        )
    else:
        fig.update_xaxes(tickformat="%m-%Y", tickangle=-35, gridcolor="#F2F2F2", title=None)


def latest_summary(df: pd.DataFrame, *, metric: str, modalities: Sequence[str]) -> pd.DataFrame:
    if df is None or df.empty:
        return pd.DataFrame(columns=["Modalidade", "Valor", "Participação"])
    subset = df[(df["metrica"].eq(metric)) & (df["modalidade"].isin(list(modalities)))].copy()
    if subset.empty:
        return pd.DataFrame(columns=["Modalidade", "Valor", "Participação"])
    last_period = subset["periodo"].max()
    latest = subset[subset["periodo"].eq(last_period)].copy()
    latest["valor"] = pd.to_numeric(latest["valor"], errors="coerce").fillna(0.0)
    total = latest["valor"].sum()
    latest["Participação"] = latest["valor"] / total if total else pd.NA
    latest["Modalidade"] = latest["modalidade"]
    latest["Valor"] = latest["valor"]
    latest = latest[["Modalidade", "Valor", "Participação"]]
    order = {name: idx for idx, name in enumerate(modalities)}
    latest["_order"] = latest["Modalidade"].map(order).fillna(999)
    return latest.sort_values("_order").drop(columns="_order").reset_index(drop=True)


def build_meios_pagamento_pptx(
    *,
    title: str,
    subtitle: str,
    figures: Mapping[str, bytes | None],
    summaries: Mapping[str, pd.DataFrame],
    source_note: str,
) -> bytes:
    """Gera PPTX simples para download no Streamlit.

    As figuras devem vir como PNG. Quando o ambiente não tiver exportação de
    Plotly instalada, a função mantém slides-resumo tabulares.
    """
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    black = RGBColor(31, 31, 31)
    orange = RGBColor(236, 112, 0)
    gray = RGBColor(107, 107, 107)
    light_gray = RGBColor(229, 229, 229)

    def add_text(slide: Any, text: str, x: float, y: float, w: float, h: float, *, size: int, bold: bool = False, color: RGBColor = black, align: Any = PP_ALIGN.LEFT):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.name = "Arial"
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        return box

    def add_footer(slide: Any, page: int) -> None:
        add_text(slide, "Banco Central do Brasil | Elaboração: Toma Conta", 0.55, 7.1, 6.8, 0.25, size=8, color=gray)
        add_text(slide, str(page), 12.25, 7.1, 0.45, 0.25, size=8, color=gray, align=PP_ALIGN.RIGHT)

    cover = prs.slides.add_slide(prs.slide_layouts[6])
    bar = cover.shapes.add_shape(1, Inches(0.55), Inches(1.55), Inches(0.12), Inches(2.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = orange
    bar.line.color.rgb = orange
    add_text(cover, title, 0.82, 1.55, 9.7, 0.65, size=34, bold=True)
    add_text(cover, subtitle, 0.84, 2.35, 10.4, 0.75, size=15, color=gray)
    add_text(cover, source_note, 0.84, 6.65, 10.8, 0.32, size=9, color=gray)
    add_footer(cover, 1)

    page = 2
    for slide_title, png in figures.items():
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_text(slide, slide_title, 0.55, 0.35, 11.6, 0.42, size=21, bold=True)
        add_text(slide, "Último ponto rotulado na cor da respectiva série; eixo de tempo preserva a periodicidade oficial.", 0.57, 0.83, 11.7, 0.3, size=9, color=gray)
        if png:
            slide.shapes.add_picture(io.BytesIO(png), Inches(0.55), Inches(1.25), width=Inches(12.1), height=Inches(5.55))
        else:
            add_text(slide, "Figura indisponível neste ambiente. Use o dashboard interativo para visualizar o gráfico.", 1.1, 2.8, 10.6, 0.35, size=16, color=gray, align=PP_ALIGN.CENTER)
        add_footer(slide, page)
        page += 1

    if summaries:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_text(slide, "Participação no último período disponível", 0.55, 0.35, 11.6, 0.42, size=21, bold=True)
        add_text(slide, "Leitura executiva para conferir o espaço ocupado por cada modalidade na última observação selecionada.", 0.57, 0.83, 11.7, 0.3, size=9, color=gray)
        y = 1.35
        for label, summary in summaries.items():
            add_text(slide, label, 0.65, y, 5.5, 0.28, size=13, bold=True)
            rows = min(len(summary), 8)
            table = slide.shapes.add_table(rows + 1, 3, Inches(0.65), Inches(y + 0.35), Inches(5.9), Inches(0.35 + 0.28 * (rows + 1))).table
            headers = ["Modalidade", "Valor", "%"]
            for col, header in enumerate(headers):
                cell = table.cell(0, col)
                cell.text = header
                cell.fill.solid()
                cell.fill.fore_color.rgb = black
                cell.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(255, 255, 255)
                cell.text_frame.paragraphs[0].runs[0].font.size = Pt(8)
            for row_idx, (_, row) in enumerate(summary.head(rows).iterrows(), start=1):
                vals = [
                    str(row.get("Modalidade", "")),
                    _format_short_number(row.get("Valor")),
                    _format_short_number(row.get("Participação"), percent=True),
                ]
                for col, val in enumerate(vals):
                    cell = table.cell(row_idx, col)
                    cell.text = val
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(248, 248, 248) if row_idx % 2 else RGBColor(255, 255, 255)
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(8)
                            run.font.name = "Arial"
            y += 2.55
            if y > 5.5:
                break
        divider = slide.shapes.add_shape(1, Inches(6.95), Inches(1.25), Inches(0.02), Inches(5.25))
        divider.fill.solid()
        divider.fill.fore_color.rgb = light_gray
        divider.line.color.rgb = light_gray
        add_text(slide, "Observações de leitura", 7.25, 1.35, 4.8, 0.28, size=13, bold=True)
        add_text(
            slide,
            "1. Pix desloca a comparação por quantidade; por isso o filtro de período é indispensável.\n"
            "2. Valor e quantidade contam histórias diferentes: TED segue relevante em ticket médio, enquanto Pix domina frequência.\n"
            "3. DOC/TEC zerados no histórico recente devem ser mantidos para preservar a série e sinalizar descontinuidade.",
            7.25,
            1.85,
            4.95,
            2.2,
            size=11,
            color=black,
        )
        add_footer(slide, page)

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()
