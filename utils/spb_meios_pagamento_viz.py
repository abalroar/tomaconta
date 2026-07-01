from __future__ import annotations

from io import BytesIO
from typing import Mapping, Sequence

import pandas as pd
import plotly.graph_objects as go


ITAU_BBA_PALETTE = {
    "Pix": "#EC7000",
    "TED": "#111111",
    "Boleto": "#6F6F6F",
    "Cheque": "#B8B8B8",
    "DOC": "#7E3F98",
    "TEC": "#F4C542",
    "Cartão de Crédito": "#003A70",
    "Cartão de Débito": "#1F77B4",
    "Cartão Pré-Pago": "#63A8E2",
    "Transferências Intrabancárias": "#4B4B4B",
    "Convênios": "#F7A35C",
    "Débito Direto": "#9A9A9A",
    "Saques": "#5B2C6F",
    "Cobranças": "#6F6F6F",
    "Cartões": "#003A70",
    "Transferências Interbancárias": "#111111",
}

SPB_INSTRUMENTOS_LABEL = {
    "pix": "Pix",
    "ted": "TED",
    "tec": "TEC",
    "cheque": "Cheque",
    "boleto": "Boleto",
    "doc": "DOC",
    "cartao_credito": "Cartão de Crédito",
    "cartao_debito": "Cartão de Débito",
    "cartao_pre_pago": "Cartão Pré-Pago",
    "trans_intrabancaria": "Transferências Intrabancárias",
    "convenios": "Convênios",
    "debito_direto": "Débito Direto",
    "saques": "Saques",
}

SPB_CONSOLIDADO_MAPA = {
    "Pix": "Pix",
    "Saques": "Saques",
    "Boleto": "Cobranças",
    "Convênios": "Cobranças",
    "Débito Direto": "Cobranças",
    "Cartão de Crédito": "Cartões",
    "Cartão de Débito": "Cartões",
    "Cartão Pré-Pago": "Cartões",
    "TED": "Transferências Interbancárias",
    "DOC": "Transferências Interbancárias",
    "TEC": "Transferências Interbancárias",
    "Cheque": "Cheque",
    "Transferências Intrabancárias": "Transferências Intrabancárias",
}

MONTHLY_DEFAULT_INSTRUMENTS = ["Pix", "TED", "Boleto", "Cheque", "TEC", "DOC"]
QUARTERLY_DEFAULT_INSTRUMENTS = [
    "Pix",
    "TED",
    "Boleto",
    "Cheque",
    "Cartão de Crédito",
    "Cartão de Débito",
    "Cartão Pré-Pago",
]

_MONTHS_ABBR = ["jan", "fev", "mar", "abr", "mai", "jun", "jul", "ago", "set", "out", "nov", "dez"]


def format_ano_mes_label(value) -> str:
    s = str(value)
    if len(s) == 6 and s.isdigit():
        return f"{s[4:6]}-{s[:4]}"
    try:
        dt = pd.Timestamp(value)
        if not pd.isna(dt):
            return f"{dt.month:02d}-{dt.year}"
    except Exception:
        pass
    return s


def format_trimestre_label(value) -> str:
    dt = _period_to_timestamp(value, "trimestral")
    if pd.isna(dt):
        return str(value)
    return f"{_MONTHS_ABBR[dt.month - 1]}-{str(dt.year)[-2:]}"


def _period_to_timestamp(value, periodicidade: str) -> pd.Timestamp:
    s = str(value)
    if periodicidade == "mensal":
        if len(s) == 6 and s.isdigit():
            return pd.Timestamp(year=int(s[:4]), month=int(s[4:6]), day=1)
        return pd.Timestamp(value)

    if len(s) == 5 and s.isdigit():
        quarter = int(s[4])
        return pd.Timestamp(year=int(s[:4]), month=quarter * 3, day=1)
    dt = pd.Timestamp(value)
    quarter_month = ((dt.month - 1) // 3 + 1) * 3
    return pd.Timestamp(year=dt.year, month=quarter_month, day=1)


def melt_nucleo_spb(df: pd.DataFrame, periodo_col: str, periodicidade: str) -> pd.DataFrame:
    columns = ["periodo", "periodo_ordem", "periodo_label", "instrumento", "tipo", "valor"]
    if df.empty or periodo_col not in df.columns:
        return pd.DataFrame(columns=columns)

    rows = []
    for prefix, tipo in (("valor_", "Valor (R$ milhão)"), ("quantidade_", "Quantidade (mil)")):
        value_cols = [col for col in df.columns if col.startswith(prefix)]
        for col in value_cols:
            suffix = col[len(prefix):]
            part = df[[periodo_col, col]].copy()
            part.columns = ["periodo", "valor"]
            part["instrumento"] = SPB_INSTRUMENTOS_LABEL.get(suffix, suffix)
            part["tipo"] = tipo
            rows.append(part)

    if not rows:
        return pd.DataFrame(columns=columns)

    long_df = pd.concat(rows, ignore_index=True)
    long_df["valor"] = pd.to_numeric(long_df["valor"], errors="coerce")
    long_df["periodo_ordem"] = long_df["periodo"].map(lambda value: _period_to_timestamp(value, periodicidade))
    if periodicidade == "mensal":
        long_df["periodo_label"] = long_df["periodo"].map(format_ano_mes_label)
    else:
        long_df["periodo_label"] = long_df["periodo"].map(format_trimestre_label)
    return long_df.dropna(subset=["valor", "periodo_ordem"]).sort_values(["periodo_ordem", "instrumento"])


def period_options(long_df: pd.DataFrame) -> list[pd.Timestamp]:
    if long_df.empty:
        return []
    return (
        long_df[["periodo_ordem"]]
        .drop_duplicates()
        .sort_values("periodo_ordem")["periodo_ordem"]
        .tolist()
    )


def period_label_map(long_df: pd.DataFrame) -> dict[pd.Timestamp, str]:
    if long_df.empty:
        return {}
    labels = (
        long_df[["periodo_ordem", "periodo_label"]]
        .drop_duplicates()
        .sort_values("periodo_ordem")
    )
    return dict(zip(labels["periodo_ordem"], labels["periodo_label"]))


def filter_period_range(long_df: pd.DataFrame, start_period, end_period) -> pd.DataFrame:
    if long_df.empty:
        return long_df.copy()
    return long_df[
        (long_df["periodo_ordem"] >= pd.Timestamp(start_period))
        & (long_df["periodo_ordem"] <= pd.Timestamp(end_period))
    ].copy()


def available_instruments(long_df: pd.DataFrame) -> list[str]:
    if long_df.empty:
        return []
    return sorted(long_df["instrumento"].dropna().unique().tolist(), key=lambda x: _instrument_sort_key(x))


def default_start_index(periods: Sequence[pd.Timestamp], periodicidade: str) -> int:
    if not periods:
        return 0
    cutoff = pd.Timestamp("2020-11-01") if periodicidade == "mensal" else pd.Timestamp("2020-12-01")
    for idx, period in enumerate(periods):
        if pd.Timestamp(period) >= cutoff:
            return idx
    return max(0, len(periods) - 24)


def compute_share(long_df: pd.DataFrame, tipo: str, instruments: Sequence[str]) -> pd.DataFrame:
    df = long_df[(long_df["tipo"] == tipo) & (long_df["instrumento"].isin(instruments))].copy()
    if df.empty:
        return df.assign(participacao=pd.Series(dtype="float64"))
    totals = df.groupby("periodo_ordem", as_index=False)["valor"].sum().rename(columns={"valor": "total_periodo"})
    df = df.merge(totals, on="periodo_ordem", how="left")
    df["participacao"] = df["valor"] / df["total_periodo"].replace(0, pd.NA) * 100
    return df.dropna(subset=["participacao"]).sort_values(["periodo_ordem", "instrumento"])


def latest_summary(long_df: pd.DataFrame, tipo: str, instruments: Sequence[str]) -> pd.DataFrame:
    share_df = compute_share(long_df, tipo, instruments)
    if share_df.empty:
        return pd.DataFrame(columns=["Período", "Instrumento", "Valor", "Participação"])
    latest_period = share_df["periodo_ordem"].max()
    latest = share_df[share_df["periodo_ordem"] == latest_period].copy()
    latest = latest.sort_values("participacao", ascending=False)
    return pd.DataFrame(
        {
            "Período": latest["periodo_label"],
            "Instrumento": latest["instrumento"],
            "Valor": latest["valor"],
            "Participação": latest["participacao"],
        }
    )


def build_line_figure(
    long_df: pd.DataFrame,
    *,
    tipo: str,
    instruments: Sequence[str],
    title: str,
    yaxis_title: str,
    palette: Mapping[str, str] | None = None,
) -> go.Figure:
    palette = palette or ITAU_BBA_PALETTE
    fig = go.Figure()
    df = long_df[(long_df["tipo"] == tipo) & (long_df["instrumento"].isin(instruments))].copy()
    x_order = _x_order(df)

    for instrument in instruments:
        series = df[df["instrumento"] == instrument].sort_values("periodo_ordem")
        if series.empty:
            continue
        text = [""] * len(series)
        text[-1] = _format_value(series.iloc[-1]["valor"], tipo)
        color = palette.get(instrument, "#4B5563")
        fig.add_trace(
            go.Scatter(
                x=series["periodo_label"],
                y=series["valor"],
                mode="lines+markers+text",
                name=instrument,
                text=text,
                textposition="middle right",
                textfont=dict(color=color, size=11),
                line=dict(color=color, width=2.4),
                marker=dict(color=color, size=6),
                cliponaxis=False,
                hovertemplate="<b>%{fullData.name}</b><br>%{x}<br>%{y:,.2f}<extra></extra>",
            )
        )

    _apply_base_layout(fig, title, yaxis_title, x_order)
    return fig


def build_share_figure(
    long_df: pd.DataFrame,
    *,
    tipo: str,
    instruments: Sequence[str],
    title: str,
    palette: Mapping[str, str] | None = None,
) -> go.Figure:
    palette = palette or ITAU_BBA_PALETTE
    fig = go.Figure()
    share_df = compute_share(long_df, tipo, instruments)
    x_order = _x_order(share_df)

    for instrument in instruments:
        series = share_df[share_df["instrumento"] == instrument].sort_values("periodo_ordem")
        if series.empty:
            continue
        text = [""] * len(series)
        text[-1] = f"{series.iloc[-1]['participacao']:.1f}%".replace(".", ",")
        color = palette.get(instrument, "#4B5563")
        fig.add_trace(
            go.Scatter(
                x=series["periodo_label"],
                y=series["participacao"],
                mode="lines+markers+text",
                name=instrument,
                text=text,
                textposition="middle right",
                textfont=dict(color=color, size=11),
                line=dict(color=color, width=2.3),
                marker=dict(color=color, size=6),
                cliponaxis=False,
                hovertemplate="<b>%{fullData.name}</b><br>%{x}<br>%{y:.1f}%<extra></extra>",
            )
        )

    _apply_base_layout(fig, title, "% do total selecionado", x_order)
    fig.update_yaxes(range=[0, 100], ticksuffix="%")
    return fig


def build_spb_pptx(
    *,
    title: str,
    figures: Mapping[str, bytes | None],
    summaries: Mapping[str, pd.DataFrame],
    source_note: str,
) -> bytes:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    def add_title(slide, text: str, top: float = 0.28):
        box = slide.shapes.add_textbox(Inches(0.45), Inches(top), Inches(12.2), Inches(0.38))
        frame = box.text_frame
        frame.clear()
        p = frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RGBColor(17, 17, 17)
        return box

    def add_footer(slide):
        box = slide.shapes.add_textbox(Inches(0.45), Inches(7.05), Inches(12.0), Inches(0.22))
        p = box.text_frame.paragraphs[0]
        p.text = source_note
        p.font.size = Pt(7.5)
        p.font.color.rgb = RGBColor(105, 105, 105)

    blank = prs.slide_layouts[6]
    cover = prs.slides.add_slide(blank)
    cover.background.fill.solid()
    cover.background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    accent = cover.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.12), Inches(7.5))
    accent.fill.solid()
    accent.fill.fore_color.rgb = RGBColor(236, 112, 0)
    accent.line.fill.background()
    title_box = cover.shapes.add_textbox(Inches(0.55), Inches(2.45), Inches(11.8), Inches(0.8))
    p = title_box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = RGBColor(17, 17, 17)
    sub = cover.shapes.add_textbox(Inches(0.58), Inches(3.35), Inches(10.8), Inches(0.55))
    p = sub.text_frame.paragraphs[0]
    p.text = "Banco Central do Brasil | MPV_DadosAbertos"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(85, 85, 85)
    add_footer(cover)

    for chart_title, png in figures.items():
        slide = prs.slides.add_slide(blank)
        add_title(slide, chart_title)
        if png:
            slide.shapes.add_picture(BytesIO(png), Inches(0.45), Inches(0.85), width=Inches(12.25), height=Inches(5.85))
        else:
            box = slide.shapes.add_textbox(Inches(1.0), Inches(2.7), Inches(11.2), Inches(0.6))
            p = box.text_frame.paragraphs[0]
            p.text = "Imagem do gráfico indisponível neste ambiente."
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(105, 105, 105)
        add_footer(slide)

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Última composição disponível")
    left = 0.55
    for section, df in summaries.items():
        box = slide.shapes.add_textbox(Inches(left), Inches(0.95), Inches(5.9), Inches(5.85))
        frame = box.text_frame
        frame.clear()
        header = frame.paragraphs[0]
        header.text = section
        header.font.bold = True
        header.font.size = Pt(13)
        header.font.color.rgb = RGBColor(236, 112, 0)
        for _, row in df.head(12).iterrows():
            p = frame.add_paragraph()
            value_label = _format_value(row["Valor"], "Valor (R$ milhão)" if "Valor" in section else "Quantidade (mil)")
            share_label = f"{row['Participação']:.1f}".replace(".", ",")
            p.text = f"{row['Instrumento']}: {value_label} ({share_label}%)"
            p.font.size = Pt(9.5)
            p.font.color.rgb = RGBColor(35, 35, 35)
        left += 6.25
    add_footer(slide)

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Plano de melhorias aplicado")
    bullets = [
        "Manter mensal e trimestral separados para preservar a periodicidade oficial do BCB.",
        "Permitir seleção de período inicial e final para evitar perda de escala após a entrada do Pix.",
        "Exibir rótulos no último ponto de cada série com a mesma cor da linha.",
        "Usar meses compactos no eixo: mm-aaaa no mensal e mês de fechamento no trimestral.",
        "Mostrar a última participação percentual por instrumento.",
        "Aplicar paleta fixa de alto contraste em estilo Itaú BBA.",
    ]
    box = slide.shapes.add_textbox(Inches(0.75), Inches(1.05), Inches(11.8), Inches(5.7))
    frame = box.text_frame
    frame.clear()
    for idx, bullet in enumerate(bullets):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(15)
        p.font.color.rgb = RGBColor(35, 35, 35)
        p.space_after = Pt(9)
    add_footer(slide)

    output = BytesIO()
    prs.save(output)
    return output.getvalue()


def build_spb_native_pptx(
    *,
    title: str,
    charts: Sequence[Mapping[str, object]],
    summaries: Mapping[str, pd.DataFrame],
    source_note: str,
) -> bytes:
    from pptx import Presentation
    from pptx.chart.data import ChartData
    from pptx.dml.color import RGBColor
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches, Pt

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def rgb(hex_color: str) -> RGBColor:
        value = hex_color.lstrip("#")
        return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))

    def add_title(slide, text: str, top: float = 0.28):
        box = slide.shapes.add_textbox(Inches(0.45), Inches(top), Inches(12.2), Inches(0.4))
        frame = box.text_frame
        frame.clear()
        p = frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = rgb("#111111")

    def add_footer(slide):
        box = slide.shapes.add_textbox(Inches(0.45), Inches(7.06), Inches(12.0), Inches(0.22))
        p = box.text_frame.paragraphs[0]
        p.text = source_note
        p.font.size = Pt(7.5)
        p.font.color.rgb = rgb("#6F6F6F")

    def add_cover():
        slide = prs.slides.add_slide(blank)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = rgb("#FFFFFF")
        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.12), Inches(7.5))
        accent.fill.solid()
        accent.fill.fore_color.rgb = rgb("#EC7000")
        accent.line.fill.background()
        title_box = slide.shapes.add_textbox(Inches(0.55), Inches(2.45), Inches(11.8), Inches(0.85))
        p = title_box.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(30)
        p.font.bold = True
        p.font.color.rgb = rgb("#111111")
        sub = slide.shapes.add_textbox(Inches(0.58), Inches(3.35), Inches(10.8), Inches(0.55))
        p = sub.text_frame.paragraphs[0]
        p.text = "Banco Central do Brasil | MPV_DadosAbertos"
        p.font.size = Pt(14)
        p.font.color.rgb = rgb("#555555")
        add_footer(slide)

    def add_latest_strip(slide, latest_items: list[dict[str, object]], *, share: bool):
        box = slide.shapes.add_textbox(Inches(10.45), Inches(1.05), Inches(2.45), Inches(5.65))
        frame = box.text_frame
        frame.clear()
        header = frame.paragraphs[0]
        header.text = "Último ponto"
        header.font.bold = True
        header.font.size = Pt(11)
        header.font.color.rgb = rgb("#111111")
        for item in latest_items[:11]:
            p = frame.add_paragraph()
            label = f"{item['instrumento']}: {item['label']}"
            p.text = label
            p.font.size = Pt(8.5)
            p.font.color.rgb = rgb(str(item["color"]))
            p.space_after = Pt(2)

    def add_chart_slide(spec: Mapping[str, object]):
        chart_title = str(spec["title"])
        long_df = spec["long_df"]
        tipo = str(spec["tipo"])
        instruments = list(spec["instruments"])
        yaxis_title = str(spec.get("yaxis_title", ""))
        share = bool(spec.get("share", False))
        palette = spec.get("palette") or ITAU_BBA_PALETTE

        slide = prs.slides.add_slide(blank)
        add_title(slide, chart_title)
        categories, series_payload, latest_items = _native_chart_payload(
            long_df,
            tipo=tipo,
            instruments=instruments,
            share=share,
            palette=palette,
        )
        if not categories or not series_payload:
            box = slide.shapes.add_textbox(Inches(0.8), Inches(2.7), Inches(11.6), Inches(0.45))
            p = box.text_frame.paragraphs[0]
            p.text = "Sem dados disponíveis para a janela selecionada."
            p.font.size = Pt(15)
            p.font.color.rgb = rgb("#6F6F6F")
            add_footer(slide)
            return

        chart_data = ChartData()
        chart_data.categories = categories
        for series_name, values, _color in series_payload:
            chart_data.add_series(series_name, values)

        frame = slide.shapes.add_chart(
            XL_CHART_TYPE.LINE_MARKERS,
            Inches(0.45),
            Inches(0.95),
            Inches(9.75),
            Inches(5.55),
            chart_data,
        )
        chart = frame.chart
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.value_axis.has_title = True
        chart.value_axis.axis_title.text_frame.text = yaxis_title
        chart.value_axis.has_major_gridlines = True
        chart.value_axis.tick_labels.number_format = "0%" if share else "#,##0"
        chart.category_axis.tick_labels.font.size = Pt(8)
        chart.value_axis.tick_labels.font.size = Pt(8)

        for idx, series in enumerate(chart.series):
            color = rgb(series_payload[idx][2])
            series.format.line.color.rgb = color
            series.format.line.width = Pt(2)
            try:
                series.marker.format.fill.solid()
                series.marker.format.fill.fore_color.rgb = color
                series.marker.format.line.color.rgb = color
            except Exception:
                pass

        add_latest_strip(slide, latest_items, share=share)
        add_footer(slide)

    def add_summary_slide():
        slide = prs.slides.add_slide(blank)
        add_title(slide, "Última composição disponível")
        left = 0.55
        for section, df in summaries.items():
            box = slide.shapes.add_textbox(Inches(left), Inches(0.95), Inches(5.9), Inches(5.85))
            frame = box.text_frame
            frame.clear()
            header = frame.paragraphs[0]
            header.text = section
            header.font.bold = True
            header.font.size = Pt(13)
            header.font.color.rgb = rgb("#EC7000")
            for _, row in df.head(12).iterrows():
                p = frame.add_paragraph()
                value_label = _format_value(row["Valor"], "Valor (R$ milhão)" if "Valor" in section else "Quantidade (mil)")
                share_label = f"{row['Participação']:.1f}".replace(".", ",")
                p.text = f"{row['Instrumento']}: {value_label} ({share_label}%)"
                p.font.size = Pt(9.5)
                p.font.color.rgb = rgb("#232323")
            left += 6.25
        add_footer(slide)

    def add_plan_slide():
        slide = prs.slides.add_slide(blank)
        add_title(slide, "Plano de melhorias aplicado")
        bullets = [
            "Manter mensal e trimestral separados para preservar a periodicidade oficial do BCB.",
            "Permitir seleção de período inicial e final para evitar perda de escala após a entrada do Pix.",
            "Exibir a última informação das séries no PPT e na interface.",
            "Usar meses compactos no eixo: mm-aaaa no mensal e mês de fechamento no trimestral.",
            "Mostrar a última participação percentual por instrumento.",
            "Aplicar paleta fixa de alto contraste em estilo Itaú BBA.",
        ]
        box = slide.shapes.add_textbox(Inches(0.75), Inches(1.05), Inches(11.8), Inches(5.7))
        frame = box.text_frame
        frame.clear()
        for idx, bullet in enumerate(bullets):
            p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
            p.text = bullet
            p.font.size = Pt(15)
            p.font.color.rgb = rgb("#232323")
            p.space_after = Pt(9)
        add_footer(slide)

    add_cover()
    for chart_spec in charts:
        add_chart_slide(chart_spec)
    add_summary_slide()
    add_plan_slide()

    output = BytesIO()
    prs.save(output)
    return output.getvalue()


def _apply_base_layout(fig: go.Figure, title: str, yaxis_title: str, x_order: Sequence[str]) -> None:
    fig.update_layout(
        title=dict(text=title, font=dict(size=15, color="#111111")),
        height=460,
        plot_bgcolor="white",
        paper_bgcolor="white",
        font=dict(family="IBM Plex Sans, Arial", color="#222222"),
        legend=dict(orientation="h", yanchor="bottom", y=1.02, xanchor="left", x=0),
        margin=dict(l=55, r=115, t=72, b=82),
        xaxis=dict(
            title="",
            type="category",
            categoryorder="array",
            categoryarray=list(x_order),
            tickangle=-45,
            showgrid=False,
            automargin=True,
        ),
        yaxis=dict(title=yaxis_title, gridcolor="#E6E6E6", zerolinecolor="#D8D8D8"),
    )


def _x_order(df: pd.DataFrame) -> list[str]:
    if df.empty:
        return []
    return (
        df[["periodo_ordem", "periodo_label"]]
        .drop_duplicates()
        .sort_values("periodo_ordem")["periodo_label"]
        .tolist()
    )


def _native_chart_payload(
    long_df: pd.DataFrame,
    *,
    tipo: str,
    instruments: Sequence[str],
    share: bool,
    palette: Mapping[str, str],
) -> tuple[list[str], list[tuple[str, tuple[float | None, ...], str]], list[dict[str, object]]]:
    if share:
        df = compute_share(long_df, tipo, instruments)
        value_col = "participacao"
    else:
        df = long_df[(long_df["tipo"] == tipo) & (long_df["instrumento"].isin(instruments))].copy()
        value_col = "valor"
    if df.empty:
        return [], [], []

    order = (
        df[["periodo_ordem", "periodo_label"]]
        .drop_duplicates()
        .sort_values("periodo_ordem")
    )
    categories = order["periodo_label"].tolist()
    payload: list[tuple[str, tuple[float | None, ...], str]] = []
    latest_items: list[dict[str, object]] = []

    for instrument in instruments:
        series = df[df["instrumento"] == instrument].sort_values("periodo_ordem")
        if series.empty:
            continue
        by_label = series.set_index("periodo_label")[value_col]
        values = []
        for label in categories:
            value = by_label.get(label, None)
            if value is None or pd.isna(value):
                values.append(None)
            else:
                number = float(value) / 100 if share else float(value)
                values.append(number)
        color = palette.get(instrument, "#4B5563")
        payload.append((instrument, tuple(values), color))

        latest = series.dropna(subset=[value_col]).iloc[-1]
        latest_value = float(latest[value_col])
        latest_items.append(
            {
                "instrumento": instrument,
                "color": color,
                "label": f"{latest_value:.1f}%".replace(".", ",") if share else _format_value(latest_value, tipo),
            }
        )

    return categories, payload, latest_items


def _instrument_sort_key(instrument: str) -> tuple[int, str]:
    preferred = [
        "Pix",
        "TED",
        "Boleto",
        "Cheque",
        "TEC",
        "DOC",
        "Cartão de Crédito",
        "Cartão de Débito",
        "Cartão Pré-Pago",
        "Saques",
        "Convênios",
        "Débito Direto",
        "Transferências Intrabancárias",
    ]
    try:
        return (preferred.index(instrument), instrument)
    except ValueError:
        return (999, instrument)


def _format_value(value: float, tipo: str) -> str:
    if pd.isna(value):
        return ""
    formatted = f"{float(value):,.0f}".replace(",", ".")
    if tipo == "Valor (R$ milhão)":
        return f"R$ {formatted}"
    return formatted
