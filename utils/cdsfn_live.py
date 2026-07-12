from __future__ import annotations

import json
import re
import unicodedata
from concurrent.futures import ThreadPoolExecutor
from datetime import datetime
from io import BytesIO
from pathlib import Path
from typing import Any
from zoneinfo import ZoneInfo

import pandas as pd
import requests
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt

from utils.formatting import formatar_numero_br


BCB_INFORMES_BASE_URL = "https://www3.bcb.gov.br/informes/rest"
PESSOAS_JURIDICAS_URL = f"{BCB_INFORMES_BASE_URL}/pessoasJuridicas"
BALANCO_DOWNLOAD_URL = f"{BCB_INFORMES_BASE_URL}/balanco//download"

CDSFN_BLOCKS = {
    "BalancoPatrimonial": {"sigla": "BP", "label": "Balanço Patrimonial"},
    "DemonstracaoDoResultado": {"sigla": "DRE", "label": "Demonstração do Resultado"},
    "DemonstracaoDoResultadoAbrangente": {"sigla": "DRA", "label": "Demonstração do Resultado Abrangente"},
    "DemonstracaoDosFluxosDeCaixa": {"sigla": "DFC", "label": "Demonstração dos Fluxos de Caixa"},
    "DemonstracaoDasMutacoesDoPatrimonioLiquido": {"sigla": "DMPL", "label": "Demonstração das Mutações do Patrimônio Líquido"},
}

DEFAULT_HEADERS = {
    "User-Agent": "Mozilla/5.0 (compatible; TomaConta/1.0; +https://github.com/matheusjprates/tomaconta)",
    "Referer": "https://www3.bcb.gov.br/informes/",
}

REQUIRED_TOP_LEVEL_KEYS = ("@cnpj", "@codigoDocumento", "@dataBase")
REQUIRED_ACCOUNT_KEYS = ("@nivel", "@descricao", "@contaPai")
CSV_REQUIRED_COLUMNS = ("CNPJ", "NOME DA INSTITUICAO")


def _safe_response_json(response: Any, contexto: str) -> dict[str, Any]:
    try:
        return response.json()
    except Exception:
        try:
            return json.loads(response.text)
        except Exception as json_exc:
            raise ValueError(f"{contexto}: resposta não é JSON válido.") from json_exc


def _digits_only(value: Any) -> str:
    return re.sub(r"\D", "", str(value or ""))


def _normalize_text(value: Any) -> str:
    texto = unicodedata.normalize("NFKD", str(value or "")).encode("ascii", "ignore").decode("ascii")
    texto = texto.strip().upper()
    texto = re.sub(r"\s+", " ", texto)
    return texto


def _normalize_instituicoes(rows: list[dict[str, Any]]) -> pd.DataFrame:
    normalized: list[dict[str, Any]] = []
    vistos: set[tuple[str, str, str]] = set()

    for item in rows:
        if not isinstance(item, dict):
            continue
        nome = str(item.get("nome") or "").strip()
        cnpj = _digits_only(item.get("cnpj"))
        id_bacen = str(item.get("idBacen") or "").strip()
        if not nome and not cnpj:
            continue
        chave = (nome, cnpj, id_bacen)
        if chave in vistos:
            continue
        vistos.add(chave)
        label = nome or cnpj
        if nome and cnpj:
            label = f"{nome} ({cnpj})"
        normalized.append(
            {
                "nome": nome,
                "cnpj": cnpj,
                "id_bacen": id_bacen,
                "label": label,
                "raw": item,
            }
        )

    if not normalized:
        raise ValueError("API de instituições retornou lista vazia ou sem campos utilizáveis (nome/cnpj).")

    return pd.DataFrame(normalized).sort_values(["nome", "cnpj", "id_bacen"], na_position="last").reset_index(drop=True)


def load_instituicoes_csv_cdsfn(path: str | Path) -> pd.DataFrame:
    caminho = Path(path)
    if not caminho.exists():
        raise FileNotFoundError(f"Arquivo de instituições não encontrado: {caminho}")

    df = pd.read_csv(caminho, sep=";", encoding="latin1", skiprows=7, dtype=str)
    if len(df.columns) and str(df.columns[-1]).startswith("Unnamed"):
        df = df.iloc[:, :-1]

    faltantes = [col for col in CSV_REQUIRED_COLUMNS if col not in df.columns]
    if faltantes:
        raise ValueError(f"CSV de instituições sem colunas obrigatórias: {', '.join(faltantes)}.")

    for col in df.columns:
        df[col] = df[col].fillna("").astype(str).str.strip()

    df["cnpj"] = df["CNPJ"].map(_digits_only).str.zfill(8)
    df["nome"] = df["NOME DA INSTITUICAO"].astype(str).str.strip()
    df["municipio"] = df.get("MUNICIPIO", "").astype(str).str.strip() if "MUNICIPIO" in df.columns else ""
    df["uf"] = df.get("UF", "").astype(str).str.strip() if "UF" in df.columns else ""
    df["situacao"] = df.get("SITUAÇÃO", "").astype(str).str.strip() if "SITUAÇÃO" in df.columns else ""

    df = df[(df["cnpj"].str.len() == 8) & (df["nome"] != "")]
    if df.empty:
        raise ValueError("CSV de instituições não contém linhas válidas com CNPJ raiz e nome.")

    df = df.drop_duplicates(subset=["cnpj", "nome"], keep="first").copy()
    df["label"] = df.apply(
        lambda row: f"{row['nome']} ({row['cnpj']})" if row["cnpj"] else row["nome"],
        axis=1,
    )
    return df.sort_values(["nome", "cnpj"], na_position="last").reset_index(drop=True)


def _normalize_reference_dates(payload: dict[str, Any]) -> tuple[list[dict[str, str]], dict[str, str]]:
    datas_raw = payload.get("datasBaseReferencia") or []
    if datas_raw and not isinstance(datas_raw, list):
        raise ValueError("Campo 'datasBaseReferencia' inválido: esperado lista.")

    datas: list[dict[str, str]] = []
    mapping: dict[str, str] = {}
    for idx, item in enumerate(datas_raw):
        if not isinstance(item, dict):
            raise ValueError(f"datasBaseReferencia[{idx}] inválido: esperado objeto.")
        dt_id = str(item.get("@id") or "").strip()
        dt_data = str(item.get("@data") or "").strip()
        if not dt_id:
            raise ValueError(f"datasBaseReferencia[{idx}] sem '@id'.")
        datas.append({"id": dt_id, "data": dt_data})
        mapping[dt_id] = dt_data or dt_id
    return datas, mapping


def _level_sort_key(nivel: Any) -> tuple[int, ...]:
    texto = str(nivel or "").strip()
    if not texto:
        return (10**9,)
    partes = []
    for parte in texto.split("."):
        try:
            partes.append(int(parte))
        except ValueError:
            partes.append(10**6)
    return tuple(partes)


def _reference_sort_key(ref: str) -> tuple[int, int, str]:
    texto = str(ref or "").strip().upper()
    match = re.search(r"([A-Z])?(\d{2})(\d{4})$", texto)
    if not match:
        return (0, 0, texto)
    prefixo, mes, ano = match.groups()
    return (int(ano), int(mes), prefixo or "")


def _reference_prefix(ref: str) -> str:
    texto = str(ref or "").strip().upper()
    if not texto:
        return ""
    if texto[0].isalpha():
        return texto[0]
    return ""


def _reference_display_label(ref: str) -> str:
    texto = str(ref or "").strip().upper()
    match = re.search(r"([A-Z])?(\d{2})(\d{4})$", texto)
    if not match:
        return texto
    prefixo, mes, ano = match.groups()
    mapa_mes = {
        "01": "jan",
        "02": "fev",
        "03": "mar",
        "04": "abr",
        "05": "mai",
        "06": "jun",
        "07": "jul",
        "08": "ago",
        "09": "set",
        "10": "out",
        "11": "nov",
        "12": "dez",
    }
    mes_txt = mapa_mes.get(mes, mes)
    if prefixo:
        return f"{prefixo} {mes_txt}/{ano[-2:]}"
    return f"{mes_txt}/{ano[-2:]}"


def _format_excel_display_value_cdsfn(valor: Any) -> str:
    valor_num = pd.to_numeric(valor, errors="coerce")
    if pd.isna(valor_num):
        return "—"
    casas = 0 if float(valor_num).is_integer() else 2
    return formatar_numero_br(float(valor_num), casas=casas)


def _cdsfn_export_timestamp_gmt3() -> str:
    dt = datetime.now(ZoneInfo("America/Sao_Paulo"))
    return dt.strftime("%Y-%m-%d %H:%M:%S GMT-3")


def _build_parent_map(df_long: pd.DataFrame) -> dict[str, str]:
    base = (
        df_long[["nivel", "conta_pai"]]
        .drop_duplicates()
        .fillna("")
    )
    return {str(row["nivel"]).strip(): str(row["conta_pai"]).strip() for _, row in base.iterrows()}


def _build_description_map(df_long: pd.DataFrame) -> dict[str, str]:
    base = (
        df_long[["nivel", "descricao"]]
        .drop_duplicates()
        .fillna("")
    )
    return {str(row["nivel"]).strip(): str(row["descricao"]).strip() for _, row in base.iterrows()}


def _ancestor_chain(nivel: str, parent_map: dict[str, str]) -> list[str]:
    chain = [nivel]
    current = parent_map.get(nivel, "")
    seen = {nivel}
    while current and current not in seen:
        chain.append(current)
        seen.add(current)
        current = parent_map.get(current, "")
    return chain


def _compute_depth(nivel: str, parent_map: dict[str, str]) -> int:
    cadeia = _ancestor_chain(nivel, parent_map)
    chain_depth = max(0, len(cadeia) - 1)
    dot_depth = max(0, len(str(nivel or "").split(".")) - 1)
    return max(chain_depth, dot_depth)


def _resolve_balance_section(nivel: str, parent_map: dict[str, str], desc_map: dict[str, str]) -> str:
    texts = [_normalize_text(desc_map.get(item, "")) for item in _ancestor_chain(nivel, parent_map)]
    if any("PATRIMON" in texto for texto in texts):
        return "Patrimônio líquido"
    if any("PASSIVO" in texto for texto in texts):
        return "Passivo"
    if any("ATIVO" in texto for texto in texts):
        return "Ativo"
    return "Outras contas"


def list_reference_periods_cdsfn(payload: dict[str, Any]) -> list[dict[str, Any]]:
    metadata = extract_metadata_cdsfn(payload)
    refs = []
    vistos: set[str] = set()
    for item in metadata.get("datas_base_referencia", []):
        raw = str(item.get("data") or item.get("id") or "").strip()
        if not raw or raw in vistos:
            continue
        vistos.add(raw)
        refs.append(
            {
                "id": str(item.get("id") or "").strip(),
                "raw": raw,
                "prefixo": _reference_prefix(raw),
                "label": _reference_display_label(raw),
                "sort_key": _reference_sort_key(raw),
            }
        )
    refs = sorted(refs, key=lambda item: item["sort_key"], reverse=True)
    return refs


def build_hierarchy_frame_cdsfn(df_long: pd.DataFrame, block_key: str) -> tuple[pd.DataFrame, list[str]]:
    if df_long.empty:
        return pd.DataFrame(), []

    base_cols = ["demonstracao", "nivel", "descricao", "conta_pai", "conta_id", "ordem_conta"]
    df_wide = pivot_wide_cdsfn(df_long)
    value_columns = [col for col in df_wide.columns if col not in base_cols]

    meta = df_long[base_cols].drop_duplicates().copy()
    parent_map = _build_parent_map(df_long)
    desc_map = _build_description_map(df_long)
    child_levels = {parent for parent in parent_map.values() if parent}

    meta["depth"] = meta["nivel"].astype(str).map(lambda nivel: _compute_depth(nivel, parent_map))
    meta["has_children"] = meta["nivel"].astype(str).isin(child_levels)
    meta["section"] = (
        meta["nivel"].astype(str).map(lambda nivel: _resolve_balance_section(nivel, parent_map, desc_map))
        if block_key == "BalancoPatrimonial"
        else "Demonstração"
    )
    meta["_sort_key"] = meta["nivel"].astype(str).map(_level_sort_key)

    merged = meta.merge(df_wide, on=base_cols, how="left")
    merged = merged.sort_values(["ordem_conta", "_sort_key", "descricao"], na_position="last").reset_index(drop=True)
    merged = merged.drop(columns=["_sort_key"])
    return merged, value_columns


def build_display_table_cdsfn(
    df_long: pd.DataFrame,
    block_key: str,
    periodos_referencia: list[str],
    *,
    column_labels: dict[str, str] | None = None,
) -> pd.DataFrame:
    df_view, value_columns = build_hierarchy_frame_cdsfn(df_long, block_key)
    if df_view.empty:
        return pd.DataFrame()

    periodos_validos = [p for p in periodos_referencia]
    for periodo in periodos_validos:
        if periodo not in df_view.columns:
            df_view[periodo] = pd.NA
    colunas = ["descricao", "nivel", "depth", "has_children"] + periodos_validos
    df_export = df_view[colunas].copy()

    def _label_conta(row) -> str:
        depth = int(row.get("depth", 0) or 0)
        descricao = str(row.get("descricao") or "")
        return f"{'   ' * depth}{descricao}"

    df_export.insert(0, "Conta", df_export.apply(_label_conta, axis=1))
    df_export = df_export.drop(columns=["descricao", "nivel", "depth", "has_children"])
    if column_labels:
        df_export = df_export.rename(columns=column_labels)
    return df_export


def combine_reference_periods_cdsfn(payloads: list[dict[str, Any]]) -> list[dict[str, Any]]:
    refs_map: dict[str, dict[str, Any]] = {}
    for payload in payloads:
        for item in list_reference_periods_cdsfn(payload):
            raw = str(item.get("raw") or "").strip()
            if not raw:
                continue
            existente = refs_map.get(raw)
            if existente is None:
                refs_map[raw] = item
                continue
            if (
                existente.get("prefixo") != item.get("prefixo")
                or existente.get("label") != item.get("label")
            ):
                raise ValueError(f"Referência inconsistente encontrada entre documentos para '{raw}'.")
    return sorted(refs_map.values(), key=lambda item: item["sort_key"], reverse=True)


def combine_normalized_blocks_cdsfn(payloads: list[dict[str, Any]], block_key: str) -> pd.DataFrame:
    frames: list[pd.DataFrame] = []
    for payload in payloads:
        if block_key not in payload:
            continue
        frames.append(normalize_long_cdsfn(payload, block_key))

    if not frames:
        raise KeyError(f"Bloco '{block_key}' não existe em nenhum dos documentos carregados.")

    if len(frames) == 1:
        return frames[0].copy()

    df_all = pd.concat(frames, ignore_index=True)
    key_cols = [
        "cnpj",
        "codigo_documento",
        "demonstracao",
        "bloco_origem",
        "conta_id",
        "nivel",
        "descricao",
        "conta_pai",
        "dt_base_referencia",
    ]
    merged_rows: list[dict[str, Any]] = []

    for _, grupo in df_all.groupby(key_cols, dropna=False, sort=False):
        valores_validos = pd.to_numeric(grupo["valor"], errors="coerce").dropna().unique().tolist()
        if len(valores_validos) > 1:
            descricao = str(grupo["descricao"].iloc[0] or "")
            dt_ref = str(grupo["dt_base_referencia"].iloc[0] or "sem_data")
            raise ValueError(
                f"Conflito de valores entre documentos para '{descricao}' em '{dt_ref}'."
            )
        base = grupo.sort_values(["ordem_conta", "dt_base"], na_position="last").iloc[0].to_dict()
        base["ordem_conta"] = int(pd.to_numeric(grupo["ordem_conta"], errors="coerce").min())
        if valores_validos:
            base["valor"] = valores_validos[0]
        else:
            base["valor"] = pd.NA
        dt_bases_validas = [str(item).strip() for item in grupo["dt_base"].dropna().tolist() if str(item).strip()]
        base["dt_base"] = dt_bases_validas[0] if dt_bases_validas else None
        merged_rows.append(base)

    return pd.DataFrame(merged_rows).sort_values(
        by=["ordem_conta", "nivel", "descricao", "dt_base_referencia", "dt_base"],
        key=lambda s: s.map(_level_sort_key) if s.name == "nivel" else s,
        na_position="last",
    ).reset_index(drop=True)


def fetch_instituicoes_cdsfn(
    page_size: int = 500,
    *,
    timeout: int = 60,
    session: Any | None = None,
) -> pd.DataFrame:
    if page_size <= 0:
        raise ValueError("page_size deve ser positivo.")

    client = session or requests
    rows: list[dict[str, Any]] = []

    def _fetch_page(page: int) -> dict[str, Any]:
        payload = {
            "segmento": None,
            "nome": "",
            "cnpj": "",
            "pais": "",
            "estado": None,
            "municipio": None,
            "incluirAgencias": False,
            "incluirInstituicoesLiquidacao": False,
            "numeroPagina": page,
            "tamanhoPagina": page_size,
        }
        response = client.post(
            PESSOAS_JURIDICAS_URL,
            json=payload,
            headers=DEFAULT_HEADERS,
            timeout=timeout,
        )
        response.raise_for_status()
        return _safe_response_json(response, f"consulta de instituições CDSFN página {page}")

    primeira_pagina = _fetch_page(0)
    content = primeira_pagina.get("content")
    if not isinstance(content, list):
        raise ValueError("Resposta da API de instituições sem campo 'content' em formato lista.")
    rows.extend(content)

    total_pages = primeira_pagina.get("totalPages")
    if not isinstance(total_pages, int):
        last = primeira_pagina.get("last")
        if isinstance(last, bool) and last:
            return _normalize_instituicoes(rows)
        raise ValueError("Resposta da API de instituições sem paginação válida ('totalPages' ausente).")

    if total_pages <= 1:
        return _normalize_instituicoes(rows)

    max_workers = min(8, total_pages - 1)
    with ThreadPoolExecutor(max_workers=max_workers) as executor:
        for body in executor.map(_fetch_page, range(1, total_pages)):
            page_content = body.get("content")
            if not isinstance(page_content, list):
                raise ValueError("Resposta de paginação da API sem campo 'content' em formato lista.")
            rows.extend(page_content)

    return _normalize_instituicoes(rows)


def fetch_documento_cdsfn(
    cnpj: str,
    ano_mes: str,
    codigo_documento: str = "9011",
    *,
    timeout: int = 60,
    session: Any | None = None,
) -> tuple[dict[str, Any], str]:
    cnpj_norm = _digits_only(cnpj)
    ano_mes_norm = re.sub(r"\D", "", str(ano_mes or ""))
    codigo_norm = re.sub(r"\D", "", str(codigo_documento or ""))

    if not cnpj_norm:
        raise ValueError("CNPJ inválido para consulta do documento CDSFN.")
    if not re.fullmatch(r"\d{6}", ano_mes_norm):
        raise ValueError("anoMes inválido. Use o formato YYYYMM.")
    if not codigo_norm:
        raise ValueError("codigo_documento inválido.")

    url = f"{BALANCO_DOWNLOAD_URL}/{ano_mes_norm}-{codigo_norm}-{cnpj_norm}.json"
    client = session or requests
    response = client.get(
        url,
        params={"cnpj": cnpj_norm, "anoMes": ano_mes_norm},
        headers=DEFAULT_HEADERS,
        timeout=timeout,
    )
    if response.status_code == 404:
        raise FileNotFoundError(f"Documento {codigo_norm} não encontrado para CNPJ {cnpj_norm} em {ano_mes_norm}.")
    response.raise_for_status()
    response_text = getattr(response, "text", None)
    response_content = getattr(response, "content", None)
    if response_text is not None and not str(response_text).strip() and response_content is not None and not response_content:
        raise FileNotFoundError(
            f"Documento {codigo_norm} sem JSON publicado para CNPJ {cnpj_norm} em {ano_mes_norm}."
        )
    payload = _safe_response_json(response, "download do documento CDSFN")
    return payload, getattr(response, "url", url)


def validate_json_cdsfn(payload: Any) -> None:
    if not isinstance(payload, dict):
        raise ValueError("Documento CDSFN inválido: raiz não é objeto JSON.")

    faltantes = [campo for campo in REQUIRED_TOP_LEVEL_KEYS if campo not in payload]
    if faltantes:
        raise ValueError(f"Documento CDSFN inválido: metadados ausentes: {', '.join(faltantes)}.")

    _normalize_reference_dates(payload)

    blocos_presentes = [chave for chave in CDSFN_BLOCKS if chave in payload]
    if not blocos_presentes:
        raise ValueError("Documento CDSFN inválido: nenhum bloco financeiro reconhecido encontrado.")

    for bloco in blocos_presentes:
        body = payload.get(bloco)
        if not isinstance(body, dict):
            raise ValueError(f"Bloco '{bloco}' inválido: esperado objeto.")
        contas = body.get("contas")
        if not isinstance(contas, list):
            raise ValueError(f"Bloco '{bloco}' inválido: campo 'contas' ausente ou não é lista.")
        for idx, conta in enumerate(contas):
            if not isinstance(conta, dict):
                raise ValueError(f"Bloco '{bloco}' inválido: conta #{idx} não é objeto.")
            faltantes_conta = [campo for campo in REQUIRED_ACCOUNT_KEYS if campo not in conta]
            if faltantes_conta:
                raise ValueError(
                    f"Bloco '{bloco}' inválido: conta #{idx} sem campos obrigatórios: {', '.join(faltantes_conta)}."
                )
            valores = conta.get("valoresIndividualizados")
            if valores is None:
                continue
            if not isinstance(valores, list):
                raise ValueError(f"Bloco '{bloco}' inválido: 'valoresIndividualizados' da conta #{idx} não é lista.")
            for jdx, valor in enumerate(valores):
                if not isinstance(valor, dict):
                    raise ValueError(f"Bloco '{bloco}' inválido: valor #{jdx} da conta #{idx} não é objeto.")
                if "@dtBase" not in valor or "@valor" not in valor:
                    raise ValueError(
                        f"Bloco '{bloco}' inválido: valor #{jdx} da conta #{idx} sem '@dtBase' e/ou '@valor'."
                    )


def extract_metadata_cdsfn(payload: dict[str, Any]) -> dict[str, Any]:
    validate_json_cdsfn(payload)
    datas, mapping = _normalize_reference_dates(payload)
    return {
        "cnpj": str(payload.get("@cnpj") or "").strip(),
        "codigo_documento": str(payload.get("@codigoDocumento") or "").strip(),
        "tipo_remessa": str(payload.get("@tipoRemessa") or "").strip(),
        "unidade_medida": str(payload.get("@unidadeMedida") or "").strip(),
        "data_base": str(payload.get("@dataBase") or "").strip(),
        "datas_base_referencia": datas,
        "datas_base_map": mapping,
    }


def list_blocks_cdsfn(payload: dict[str, Any]) -> list[dict[str, Any]]:
    validate_json_cdsfn(payload)
    blocks: list[dict[str, Any]] = []
    for block_key, meta in CDSFN_BLOCKS.items():
        if block_key not in payload:
            continue
        contas = payload.get(block_key, {}).get("contas") or []
        blocks.append(
            {
                "block_key": block_key,
                "sigla": meta["sigla"],
                "label": meta["label"],
                "contas": len(contas),
            }
        )
    return blocks


def normalize_long_cdsfn(payload: dict[str, Any], block_key: str) -> pd.DataFrame:
    validate_json_cdsfn(payload)
    if block_key not in payload:
        raise KeyError(f"Bloco '{block_key}' não existe no documento consultado.")

    metadata = extract_metadata_cdsfn(payload)
    block_meta = CDSFN_BLOCKS.get(block_key, {"sigla": block_key, "label": block_key})
    contas = payload.get(block_key, {}).get("contas") or []

    rows: list[dict[str, Any]] = []
    for idx, conta in enumerate(contas):
        conta_id = str(conta.get("@id") or "").strip()
        nivel = str(conta.get("@nivel") or "").strip()
        descricao = str(conta.get("@descricao") or "").strip()
        conta_pai = str(conta.get("@contaPai") or "").strip()
        valores = conta.get("valoresIndividualizados")

        if not isinstance(valores, list) or not valores:
            rows.append(
                {
                    "cnpj": metadata["cnpj"],
                    "codigo_documento": metadata["codigo_documento"],
                    "demonstracao": block_meta["sigla"],
                    "bloco_origem": block_key,
                    "conta_id": conta_id,
                    "nivel": nivel,
                    "descricao": descricao,
                    "conta_pai": conta_pai,
                    "ordem_conta": idx,
                    "dt_base": None,
                    "dt_base_referencia": None,
                    "valor": pd.NA,
                    "unidade_medida": metadata["unidade_medida"],
                }
            )
            continue

        for valor_item in valores:
            dt_base = str(valor_item.get("@dtBase") or "").strip() or None
            rows.append(
                {
                    "cnpj": metadata["cnpj"],
                    "codigo_documento": metadata["codigo_documento"],
                    "demonstracao": block_meta["sigla"],
                    "bloco_origem": block_key,
                    "conta_id": conta_id,
                    "nivel": nivel,
                    "descricao": descricao,
                    "conta_pai": conta_pai,
                    "ordem_conta": idx,
                    "dt_base": dt_base,
                    "dt_base_referencia": metadata["datas_base_map"].get(dt_base, dt_base),
                    "valor": pd.to_numeric(valor_item.get("@valor"), errors="coerce"),
                    "unidade_medida": metadata["unidade_medida"],
                }
            )

    df = pd.DataFrame(rows)
    if df.empty:
        return df
    return df.sort_values(
        by=["ordem_conta", "nivel", "descricao", "dt_base_referencia", "dt_base"],
        key=lambda s: s.map(_level_sort_key) if s.name == "nivel" else s,
        na_position="last",
    ).reset_index(drop=True)


def pivot_wide_cdsfn(df_long: pd.DataFrame) -> pd.DataFrame:
    if df_long.empty:
        return pd.DataFrame(
            columns=[
                "demonstracao",
                "nivel",
                "descricao",
                "conta_pai",
                "conta_id",
                "ordem_conta",
            ]
        )

    base_cols = ["demonstracao", "nivel", "descricao", "conta_pai", "conta_id", "ordem_conta"]
    dt_label_col = "dt_base_referencia"
    df_work = df_long.copy()
    df_work[dt_label_col] = df_work[dt_label_col].fillna(df_work["dt_base"]).fillna("sem_data")
    pivot = (
        df_work
        .pivot_table(
            index=base_cols,
            columns=dt_label_col,
            values="valor",
            aggfunc="first",
            observed=False,
        )
        .reset_index()
    )
    pivot.columns.name = None
    return pivot.sort_values(
        by=["ordem_conta", "nivel", "descricao"],
        key=lambda s: s.map(_level_sort_key) if s.name == "nivel" else s,
        na_position="last",
    ).reset_index(drop=True)


def _write_excel_display_sheet_cdsfn(
    workbook,
    *,
    sheet_name: str,
    df_view: pd.DataFrame,
    value_columns: list[str],
    column_labels: dict[str, str] | None = None,
    title: str | None = None,
    institution_label: str | None = None,
    cnpj: str | None = None,
    source_label: str | None = None,
    unit_label: str | None = None,
    timestamp_label: str | None = None,
    reference_label: str | None = None,
    document_periods_label: str | None = None,
):
    worksheet = workbook.add_worksheet(sheet_name[:31])

    ui_font = "IBM Plex Sans"
    context_enabled = any(
        bool(item)
        for item in [title, institution_label, cnpj, source_label, unit_label, timestamp_label, reference_label, document_periods_label]
    )
    header_row_idx = 0 if not context_enabled else 4
    data_start_row = header_row_idx + 1
    last_col_idx = max(len(value_columns), 0)

    fmt_title = workbook.add_format(
        {
            "bold": True,
            "font_name": ui_font,
            "font_size": 14,
            "font_color": "#111111",
            "align": "left",
            "valign": "vcenter",
        }
    )
    fmt_meta = workbook.add_format(
        {
            "font_name": ui_font,
            "font_size": 10,
            "font_color": "#50555f",
            "align": "left",
            "valign": "vcenter",
        }
    )
    fmt_header = workbook.add_format(
        {
            "bold": True,
            "font_name": ui_font,
            "font_size": 11,
            "font_color": "#FFFFFF",
            "bg_color": "#111111",
            "border": 1,
            "align": "center",
            "valign": "vcenter",
        }
    )
    value_format_cache: dict[tuple[bool, bool], Any] = {}
    text_format_cache: dict[tuple[bool, int], Any] = {}

    def _value_format(is_parent: bool, negative: bool):
        key = (is_parent, negative)
        cached = value_format_cache.get(key)
        if cached is not None:
            return cached
        props = {
            "font_name": ui_font,
            "font_size": 11,
            "border": 1,
            "align": "right",
            "valign": "top",
            "num_format": "#,##0.00",
        }
        if is_parent:
            props["bold"] = True
            props["bg_color"] = "#F6F6F6"
        if negative:
            props["font_color"] = "#7A1E2B"
        fmt = workbook.add_format(props)
        value_format_cache[key] = fmt
        return fmt

    def _text_format(is_parent: bool, depth: int):
        key = (is_parent, depth)
        cached = text_format_cache.get(key)
        if cached is not None:
            return cached
        props = {
            "font_name": ui_font,
            "font_size": 11,
            "border": 1,
            "align": "left",
            "valign": "top",
            "indent": min(depth, 15),
        }
        if is_parent:
            props["bold"] = True
            props["bg_color"] = "#F6F6F6"
        fmt = workbook.add_format(props)
        text_format_cache[key] = fmt
        return fmt

    if context_enabled:
        worksheet.merge_range(0, 0, 0, last_col_idx, title or sheet_name, fmt_title)
        institution_parts = []
        if institution_label:
            institution_parts.append(f"Instituição: {institution_label}")
        if cnpj:
            institution_parts.append(f"CNPJ: {cnpj}")
        worksheet.merge_range(1, 0, 1, last_col_idx, " | ".join(institution_parts) if institution_parts else "", fmt_meta)

        source_parts = []
        if source_label:
            source_parts.append(f"Fonte: {source_label}")
        if reference_label:
            source_parts.append(f"Ref. Info Contábil: {reference_label}")
        if document_periods_label:
            source_parts.append(f"Competências: {document_periods_label}")
        worksheet.merge_range(2, 0, 2, last_col_idx, " | ".join(source_parts) if source_parts else "", fmt_meta)

        footer_parts = []
        if unit_label:
            footer_parts.append(f"Unidade monetária: {unit_label}")
        if timestamp_label:
            footer_parts.append(f"Extraído em: {timestamp_label}")
        worksheet.merge_range(3, 0, 3, last_col_idx, " | ".join(footer_parts) if footer_parts else "", fmt_meta)
        worksheet.set_row(0, 24)
        worksheet.set_row(1, 18)
        worksheet.set_row(2, 18)
        worksheet.set_row(3, 18)

    headers = ["Conta"] + [column_labels.get(col, col) if column_labels else col for col in value_columns]
    for col_idx, header in enumerate(headers):
        worksheet.write(header_row_idx, col_idx, header, fmt_header)

    conta_width = len("Conta")
    value_widths = [len(str(header)) for header in headers[1:]]
    worksheet.set_row(header_row_idx, 22)

    for row_idx, (_, row) in enumerate(df_view.iterrows(), start=data_start_row):
        depth = int(row.get("depth", 0) or 0)
        has_children = bool(row.get("has_children")) or depth == 0
        nivel = str(row.get("nivel") or "").strip()
        descricao = str(row.get("descricao") or "")
        conta_text = f"{nivel} {descricao}".strip() if nivel else descricao
        conta_fmt = _text_format(has_children, depth)
        worksheet.write(row_idx, 0, conta_text, conta_fmt)
        conta_width = max(conta_width, len(conta_text) + max(depth * 2, 0))
        worksheet.set_row(row_idx, 22 if has_children else 20)

        for col_pos, value_col in enumerate(value_columns, start=1):
            valor = pd.to_numeric(row.get(value_col), errors="coerce")
            negative = bool(pd.notna(valor) and float(valor) < 0)
            cell_fmt = _value_format(has_children, negative)
            if pd.isna(valor):
                worksheet.write_blank(row_idx, col_pos, None, cell_fmt)
                value_widths[col_pos - 1] = max(value_widths[col_pos - 1], len("-"))
            else:
                worksheet.write_number(row_idx, col_pos, float(valor), cell_fmt)
                value_widths[col_pos - 1] = max(value_widths[col_pos - 1], len(_format_excel_display_value_cdsfn(valor)))

    worksheet.freeze_panes(data_start_row, 1)
    worksheet.hide_gridlines(2)
    worksheet.set_zoom(90)
    worksheet.set_column(0, 0, min(max(conta_width + 6, 24), 72))
    for idx, width in enumerate(value_widths, start=1):
        worksheet.set_column(idx, idx, min(max(width + 4, 16), 24))
    return worksheet


def build_excel_display_export_cdsfn(
    df_view: pd.DataFrame,
    value_columns: list[str],
    *,
    column_labels: dict[str, str] | None = None,
    sheet_name: str = "Tabela",
) -> bytes:
    output = BytesIO()
    with pd.ExcelWriter(output, engine="xlsxwriter") as writer:
        workbook = writer.book
        worksheet = _write_excel_display_sheet_cdsfn(
            workbook,
            sheet_name=sheet_name,
            df_view=df_view,
            value_columns=value_columns,
            column_labels=column_labels,
        )
        writer.sheets[sheet_name[:31]] = worksheet
    return output.getvalue()


def build_credit_package_excel_cdsfn(
    *,
    payloads: list[dict[str, Any]],
    institution_label: str,
    cnpj: str,
    periodos_ref_sel: list[str],
    column_labels: dict[str, str] | None = None,
    reference_label: str | None = None,
    document_periods: list[str] | None = None,
) -> bytes:
    output = BytesIO()
    periodos_label = ", ".join(str(item) for item in (document_periods or []) if str(item).strip())
    timestamp_label = _cdsfn_export_timestamp_gmt3()
    metadata = extract_metadata_cdsfn(payloads[0]) if payloads else {}
    unit_label = str(metadata.get("unidade_medida") or "N/D")
    source_label = "Documento 9011 JSON do Banco Central"

    with pd.ExcelWriter(output, engine="xlsxwriter") as writer:
        workbook = writer.book
        for block_key in [
            "BalancoPatrimonial",
            "DemonstracaoDoResultado",
            "DemonstracaoDoResultadoAbrangente",
            "DemonstracaoDosFluxosDeCaixa",
            "DemonstracaoDasMutacoesDoPatrimonioLiquido",
        ]:
            if not any(block_key in payload for payload in payloads):
                continue

            df_long = combine_normalized_blocks_cdsfn(payloads, block_key)
            df_view, _ = build_hierarchy_frame_cdsfn(df_long, block_key)
            df_view_render = df_view.copy()
            for periodo_ref in periodos_ref_sel:
                if periodo_ref not in df_view_render.columns:
                    df_view_render[periodo_ref] = pd.NA
            value_columns_render = list(periodos_ref_sel)
            block_meta = CDSFN_BLOCKS[block_key]
            worksheet = _write_excel_display_sheet_cdsfn(
                workbook,
                sheet_name=block_meta["sigla"],
                df_view=df_view_render,
                value_columns=value_columns_render,
                column_labels=column_labels,
                title=f"{block_meta['sigla']} - {block_meta['label']}",
                institution_label=institution_label,
                cnpj=cnpj,
                source_label=source_label,
                unit_label=unit_label,
                timestamp_label=timestamp_label,
                reference_label=reference_label,
                document_periods_label=periodos_label,
            )
            writer.sheets[block_meta["sigla"][:31]] = worksheet
    return output.getvalue()


_PPTX_FONT = "Calibri"
_PPTX_BLACK = "1F1F1F"
_PPTX_ORANGE = "EC7000"
_PPTX_DARK_GRAY = "6B6B6B"
_PPTX_LIGHT_GRAY = "E5E5E5"
_PPTX_OFF_WHITE = "F5F5F5"
_PPTX_RED = "C8102E"
_PPTX_WHITE = "FFFFFF"


def _pptx_rgb(hex_color: str) -> RGBColor:
    texto = str(hex_color or "").strip().replace("#", "")
    if len(texto) != 6:
        texto = _PPTX_BLACK
    return RGBColor(int(texto[0:2], 16), int(texto[2:4], 16), int(texto[4:6], 16))


def _pptx_set_shape_fill(shape: Any, color: str) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = _pptx_rgb(color)


def _pptx_set_no_line(shape: Any) -> None:
    try:
        shape.line.fill.background()
    except Exception:
        shape.line.color.rgb = _pptx_rgb(_PPTX_OFF_WHITE)


def _pptx_set_slide_background(slide: Any, color: str) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _pptx_rgb(color)


def _pptx_clear_text_frame(text_frame: Any, *, word_wrap: bool = True) -> Any:
    text_frame.clear()
    text_frame.margin_left = 0
    text_frame.margin_right = 0
    text_frame.margin_top = 0
    text_frame.margin_bottom = 0
    text_frame.word_wrap = word_wrap
    return text_frame.paragraphs[0]


def _pptx_add_textbox(
    slide: Any,
    text: Any,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    font_size: float = 11,
    color: str = _PPTX_BLACK,
    bold: bool = False,
    align: Any = PP_ALIGN.LEFT,
    font_name: str = _PPTX_FONT,
) -> Any:
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    paragraph = _pptx_clear_text_frame(shape.text_frame)
    paragraph.alignment = align
    paragraph.space_after = Pt(0)
    paragraph.space_before = Pt(0)
    run = paragraph.add_run()
    run.text = str(text or "")
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bool(bold)
    run.font.color.rgb = _pptx_rgb(color)
    return shape


def _pptx_set_cell_text(
    cell: Any,
    text: Any,
    *,
    font_size: float = 8,
    color: str = _PPTX_BLACK,
    bold: bool = False,
    align: Any = PP_ALIGN.LEFT,
    valign: Any = MSO_ANCHOR.MIDDLE,
    word_wrap: bool = False,
    margin: float = 0.025,
) -> None:
    cell.vertical_anchor = valign
    cell.margin_left = Inches(margin)
    cell.margin_right = Inches(margin)
    cell.margin_top = Inches(0.01)
    cell.margin_bottom = Inches(0.01)
    paragraph = _pptx_clear_text_frame(cell.text_frame, word_wrap=word_wrap)
    paragraph.alignment = align
    paragraph.space_after = Pt(0)
    paragraph.space_before = Pt(0)
    run = paragraph.add_run()
    run.text = str(text or "")
    run.font.name = _PPTX_FONT
    run.font.size = Pt(font_size)
    run.font.bold = bool(bold)
    run.font.color.rgb = _pptx_rgb(color)


def _pptx_set_cell_fill(cell: Any, color: str) -> None:
    cell.fill.solid()
    cell.fill.fore_color.rgb = _pptx_rgb(color)


def _pptx_set_cell_border(cell: Any, color: str = _PPTX_LIGHT_GRAY, width: str = "6350") -> None:
    tc_pr = cell._tc.get_or_add_tcPr()
    for edge in ("lnL", "lnR", "lnT", "lnB"):
        ln = tc_pr.find(qn(f"a:{edge}"))
        if ln is None:
            ln = OxmlElement(f"a:{edge}")
            tc_pr.append(ln)
        ln.set("w", width)
        for child in list(ln):
            ln.remove(child)
        solid_fill = OxmlElement("a:solidFill")
        srgb = OxmlElement("a:srgbClr")
        srgb.set("val", str(color).replace("#", ""))
        solid_fill.append(srgb)
        ln.append(solid_fill)
        dash = OxmlElement("a:prstDash")
        dash.set("val", "solid")
        ln.append(dash)


def _pptx_period_label(ref: Any, column_labels: dict[str, str] | None = None) -> str:
    texto = str(ref or "").strip()
    if column_labels and texto in column_labels:
        return str(column_labels[texto])
    return _reference_display_label(texto)


def _pptx_reference_label(referencia: Any) -> str:
    texto = str(referencia or "").strip()
    upper = texto.upper()
    if upper == "A":
        return "Anual (A)"
    if upper == "S":
        return "Semestral (S)"
    return texto or "N/D"


def _pptx_unit_label(unidade: Any) -> str:
    texto = str(unidade or "").strip()
    digits = _digits_only(texto)
    if digits == "1":
        return "R$"
    if digits == "1000":
        return "R$ mil"
    if digits == "1000000":
        return "R$ milhões"
    if not texto:
        return "N/D"
    if texto.upper().startswith("R$"):
        return texto
    return f"R$ ({texto})"


def _pptx_format_number(value: Any) -> str:
    valor = pd.to_numeric(value, errors="coerce")
    if pd.isna(valor):
        return "—"
    casas = 0 if float(valor).is_integer() else 2
    return formatar_numero_br(float(valor), casas=casas)


def _pptx_format_percent(value: Any) -> str:
    valor = pd.to_numeric(value, errors="coerce")
    if pd.isna(valor):
        return "n.m."
    return formatar_numero_br(float(valor), casas=1, sufixo="%")


def _pptx_value_is_negative(value: Any) -> bool:
    valor = pd.to_numeric(value, errors="coerce")
    return bool(pd.notna(valor) and float(valor) < 0)


def _pptx_abbreviate_account_text(text: Any, max_chars: int) -> str:
    original = re.sub(r"\s+", " ", str(text or "")).strip()
    if len(original) <= max_chars:
        return original

    abbreviated = original
    normalized = _normalize_text(original)
    if "PROVISAO" in normalized and "PERDAS ESPERADAS" in normalized and "OPERACOES DE CREDITO" in normalized:
        abbreviated = re.sub(
            r"\(?-?\)?\s*Provis[aã]o.*?perdas esperadas.*?opera[cç][oõ]es de cr[eé]dito",
            "(-) Provisão perdas esperadas op. crédito",
            abbreviated,
            flags=re.IGNORECASE,
        )
    replacements = [
        ("Demonstração das Mutações do Patrimônio Líquido", "DMPL"),
        ("Demonstração do Resultado Abrangente", "DRA"),
        ("Demonstração dos Fluxos de Caixa", "DFC"),
        ("Demonstração do Resultado", "DRE"),
        ("Balanço Patrimonial", "BP"),
        ("Patrimônio líquido", "Pat. Líquido"),
        ("Patrimônio Líquido", "Pat. Líquido"),
        ("Operações de crédito", "Op. crédito"),
        ("Operações de Crédito", "Op. crédito"),
        ("Aplicações interfinanceiras de liquidez", "Aplic. interfin. liquidez"),
        ("Instrumentos financeiros derivativos", "Instr. financ. derivativos"),
        ("Resultado com perda esperada", "Res. perda esperada"),
        ("Resultado abrangente", "Res. abrangente"),
        ("Atividades operacionais", "Ativ. operacionais"),
        ("Atividades de investimento", "Ativ. investimento"),
        ("Atividades de financiamento", "Ativ. financiamento"),
        ("Caixa líquido", "Caixa líq."),
        ("Aumento líquido", "Aumento líq."),
        ("Redução líquida", "Redução líq."),
        ("no início do período", "início período"),
        ("no fim do período", "fim período"),
        ("do período", "período"),
        ("para perdas esperadas", "perdas esperadas"),
        ("operações", "op."),
        ("Operações", "Op."),
        ("financeiros", "financ."),
        ("financeiras", "financ."),
        ("financeiro", "financ."),
        ("financeira", "financ."),
        ("permanentes", "perm."),
        ("participações", "partic."),
        ("controladas", "control."),
        ("coligadas", "colig."),
        ("recursos", "rec."),
        ("obrigações", "obrig."),
        ("crédito", "créd."),
        ("Crédito", "Créd."),
    ]
    for src, dst in replacements:
        abbreviated = abbreviated.replace(src, dst)
        if len(abbreviated) <= max_chars:
            return abbreviated

    if len(abbreviated) <= max_chars:
        return abbreviated
    if max_chars <= 3:
        return abbreviated[:max_chars]
    return abbreviated[: max_chars - 3].rstrip() + "..."


def _pptx_normalize_demonstrativos(demonstrativos: dict[str, pd.DataFrame]) -> dict[str, pd.DataFrame]:
    normalized: dict[str, pd.DataFrame] = {}
    for chave, df in (demonstrativos or {}).items():
        if df is None:
            continue
        chave_texto = str(chave or "").strip()
        if chave_texto in CDSFN_BLOCKS:
            sigla = CDSFN_BLOCKS[chave_texto]["sigla"]
        else:
            sigla = chave_texto.upper()
        normalized[sigla] = df
    return normalized


def _pptx_is_major_account(row: pd.Series, block_sigla: str) -> bool:
    depth = int(row.get("depth", 0) or 0)
    if depth != 0:
        return False
    desc_norm = _normalize_text(row.get("descricao") or "")
    major_terms = (
        "ATIVO",
        "PASSIVO",
        "PATRIMONIO LIQUIDO",
        "RESULTADO",
        "FLUXO",
        "CAIXA LIQUIDO",
    )
    if block_sigla in {"BP", "DRE", "DRA", "DFC", "DMPL"} and any(term in desc_norm for term in major_terms):
        return True
    return bool(row.get("has_children")) and depth == 0


def _pptx_row_height_estimate(conta_text: str, comparative: bool, depth: int, is_parent: bool) -> float:
    return 0.27 if is_parent else 0.25


def _pptx_chunk_table_rows(rows: list[dict[str, Any]], comparative: bool) -> list[list[dict[str, Any]]]:
    max_height = 5.25
    chunks: list[list[dict[str, Any]]] = []
    current: list[dict[str, Any]] = []
    used = 0.0
    for row in rows:
        row_height = float(row.get("_height", 0.28))
        if current and used + row_height > max_height:
            chunks.append(current)
            current = []
            used = 0.0
        current.append(row)
        used += row_height
        hard_limit = 18 if comparative else 21
        if len(current) >= hard_limit:
            chunks.append(current)
            current = []
            used = 0.0
    if current:
        chunks.append(current)
    return chunks or [[]]


def _pptx_add_content_header(
    slide: Any,
    title: str,
    subtitle: str,
    chip_label: str,
) -> None:
    title_size = 20 if len(str(title or "")) > 54 else 22
    _pptx_add_textbox(slide, title, 0.4, 0.25, 10.55, 0.48, font_size=title_size, bold=True)
    _pptx_add_textbox(slide, subtitle, 0.4, 0.78, 10.65, 0.42, font_size=10.5, color=_PPTX_DARK_GRAY)
    chip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(11.25), Inches(0.28), Inches(1.68), Inches(0.38))
    _pptx_set_shape_fill(chip, _PPTX_OFF_WHITE)
    _pptx_set_no_line(chip)
    chip.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = _pptx_clear_text_frame(chip.text_frame)
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = chip_label
    run.font.name = _PPTX_FONT
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = _pptx_rgb(_PPTX_BLACK)


def _pptx_add_footer(slide: Any, page: int, total: int) -> None:
    _pptx_add_textbox(
        slide,
        "Fonte: Banco Central do Brasil — Documento 9011 | Elaboração: Toma Conta",
        0.4,
        7.12,
        8.6,
        0.18,
        font_size=7.5,
        color=_PPTX_DARK_GRAY,
    )
    _pptx_add_textbox(
        slide,
        f"Página {page} de {total}",
        11.4,
        7.12,
        1.55,
        0.18,
        font_size=7.5,
        color=_PPTX_DARK_GRAY,
        align=PP_ALIGN.RIGHT,
    )


def _pptx_build_rows_for_demonstrativo(
    df: pd.DataFrame,
    *,
    periodo_atual: str,
    periodo_anterior: str | None,
    block_sigla: str,
) -> list[dict[str, Any]]:
    if periodo_atual not in df.columns:
        raise ValueError(f"Demonstrativo {block_sigla} sem coluna do período atual '{periodo_atual}'.")
    comparative = bool(periodo_anterior)
    if comparative and periodo_anterior not in df.columns:
        raise ValueError(f"Demonstrativo {block_sigla} sem coluna do período anterior '{periodo_anterior}'.")

    rows: list[dict[str, Any]] = []
    for _, row in df.iterrows():
        depth = int(row.get("depth", 0) or 0)
        nivel = str(row.get("nivel") or "").strip()
        descricao = str(row.get("descricao") or "").strip()
        conta_text = f"{nivel} {descricao}".strip() if nivel else descricao
        if not conta_text:
            continue
        is_parent = bool(row.get("has_children")) or depth == 0
        atual = pd.to_numeric(row.get(periodo_atual), errors="coerce")
        anterior = pd.to_numeric(row.get(periodo_anterior), errors="coerce") if periodo_anterior else pd.NA
        var_abs = atual - anterior if comparative and pd.notna(atual) and pd.notna(anterior) else pd.NA
        if comparative and pd.notna(var_abs) and pd.notna(anterior) and float(anterior) != 0:
            var_pct = float(var_abs) / abs(float(anterior)) * 100.0
        else:
            var_pct = pd.NA
        var_pct_nm = comparative and (pd.isna(anterior) or float(anterior) == 0)
        rows.append(
            {
                "conta": f"{'  ' * min(max(depth, 0), 6)}{conta_text}",
                "atual": atual,
                "anterior": anterior,
                "var_abs": var_abs,
                "var_pct": var_pct,
                "var_pct_nm": var_pct_nm,
                "is_parent": is_parent,
                "is_major": _pptx_is_major_account(row, block_sigla),
                "depth": depth,
                "_height": _pptx_row_height_estimate(conta_text, comparative, depth, is_parent),
            }
        )
    return rows


def _pptx_add_table_slide(
    prs: Presentation,
    *,
    title: str,
    subtitle: str,
    chip_label: str,
    headers: list[str],
    rows: list[dict[str, Any]],
    chunk_index: int,
    chunk_total: int,
    account_max_chars: int = 92,
    body_font_size: float = 7.8,
    header_font_size: float = 8.5,
    table_y: float = 1.26,
    table_h_max: float = 5.72,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    suffix = f" ({chunk_index}/{chunk_total})" if chunk_total > 1 else ""
    _pptx_add_content_header(slide, f"{title}{suffix}", subtitle, chip_label)

    comparative = len(headers) > 2
    x, y, w = 0.4, table_y, 12.53
    header_h = 0.34
    row_heights = [float(row.get("_height", 0.28)) for row in rows]
    table_h = min(table_h_max, header_h + sum(row_heights))
    table = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(x), Inches(y), Inches(w), Inches(table_h)).table
    widths = [6.65, 1.55, 1.55, 1.55, 1.23] if comparative else [9.15, 3.38]
    for idx, width in enumerate(widths[: len(headers)]):
        table.columns[idx].width = Inches(width)

    table.rows[0].height = Inches(header_h)
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        _pptx_set_cell_fill(cell, _PPTX_BLACK)
        _pptx_set_cell_border(cell, _PPTX_BLACK, "6350")
        _pptx_set_cell_text(cell, header, font_size=header_font_size, color=_PPTX_WHITE, bold=True, align=PP_ALIGN.CENTER)

    for row_idx, row in enumerate(rows, start=1):
        bg = _PPTX_OFF_WHITE if row_idx % 2 == 0 else _PPTX_WHITE
        if row.get("is_parent"):
            bg = _PPTX_OFF_WHITE
        table.rows[row_idx].height = Inches(row_heights[row_idx - 1])

        conta_cell = table.cell(row_idx, 0)
        _pptx_set_cell_fill(conta_cell, bg)
        _pptx_set_cell_border(conta_cell)
        conta_color = _PPTX_ORANGE if row.get("is_major") else _PPTX_BLACK
        _pptx_set_cell_text(
            conta_cell,
            _pptx_abbreviate_account_text(row.get("conta", ""), account_max_chars),
            font_size=body_font_size if not row.get("is_parent") else max(body_font_size, 7.8),
            color=conta_color,
            bold=bool(row.get("is_parent")),
            align=PP_ALIGN.LEFT,
        )

        value_specs: list[tuple[str, Any, bool]] = [
            ("atual", row.get("atual"), False),
        ]
        if comparative:
            value_specs.extend(
                [
                    ("anterior", row.get("anterior"), False),
                    ("var_abs", row.get("var_abs"), False),
                    ("var_pct", row.get("var_pct"), bool(row.get("var_pct_nm"))),
                ]
            )
        for col_idx, (_, value, percent_nm) in enumerate(value_specs, start=1):
            cell = table.cell(row_idx, col_idx)
            _pptx_set_cell_fill(cell, bg)
            _pptx_set_cell_border(cell)
            if col_idx == len(value_specs) and comparative:
                texto = "n.m." if percent_nm or pd.isna(value) else _pptx_format_percent(value)
                negative = False if percent_nm or pd.isna(value) else _pptx_value_is_negative(value)
            else:
                texto = _pptx_format_number(value)
                negative = _pptx_value_is_negative(value)
            _pptx_set_cell_text(
                cell,
                texto,
                font_size=body_font_size,
                color=_PPTX_RED if negative else _PPTX_BLACK,
                bold=bool(row.get("is_parent")),
                align=PP_ALIGN.RIGHT,
            )


def _pptx_fit_single_slide_row_height(rows: list[dict[str, Any]], available_h: float, header_h: float, default_h: float) -> float:
    if not rows:
        return default_h
    exact_h = max(0.01, (available_h - header_h) / max(len(rows), 1))
    return min(default_h, exact_h)


def _pptx_headers(current_label: str, previous_label: str | None) -> list[str]:
    headers = ["Conta", current_label]
    if previous_label:
        headers.extend([previous_label, "Var. Abs.", "Var. %"])
    return headers


def _pptx_add_compact_table_to_slide(
    slide: Any,
    *,
    x: float,
    y: float,
    w: float,
    h: float,
    headers: list[str],
    rows: list[dict[str, Any]],
    account_max_chars: int,
    body_font_size: float = 7.5,
    header_font_size: float = 7.5,
    default_row_h: float = 0.17,
) -> None:
    comparative = len(headers) > 2
    header_h = 0.23
    row_h = _pptx_fit_single_slide_row_height(rows, h, header_h, default_row_h)
    table = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(x), Inches(y), Inches(w), Inches(h)).table
    widths = [0.48, 0.135, 0.135, 0.145, 0.105] if comparative else [0.72, 0.28]
    for idx, proportion in enumerate(widths[: len(headers)]):
        table.columns[idx].width = Inches(w * proportion)

    table.rows[0].height = Inches(header_h)
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        _pptx_set_cell_fill(cell, _PPTX_BLACK)
        _pptx_set_cell_border(cell, _PPTX_BLACK, "6350")
        _pptx_set_cell_text(
            cell,
            header,
            font_size=header_font_size,
            color=_PPTX_WHITE,
            bold=True,
            align=PP_ALIGN.CENTER,
            margin=0.015,
        )

    for row_idx, row in enumerate(rows, start=1):
        table.rows[row_idx].height = Inches(row_h)
        bg = _PPTX_OFF_WHITE if row_idx % 2 == 0 or row.get("is_parent") else _PPTX_WHITE

        conta_cell = table.cell(row_idx, 0)
        _pptx_set_cell_fill(conta_cell, bg)
        _pptx_set_cell_border(conta_cell)
        _pptx_set_cell_text(
            conta_cell,
            _pptx_abbreviate_account_text(row.get("conta", ""), account_max_chars),
            font_size=body_font_size,
            color=_PPTX_ORANGE if row.get("is_major") else _PPTX_BLACK,
            bold=bool(row.get("is_parent")),
            align=PP_ALIGN.LEFT,
            margin=0.015,
        )

        value_specs: list[tuple[Any, bool]] = [(row.get("atual"), False)]
        if comparative:
            value_specs.extend(
                [
                    (row.get("anterior"), False),
                    (row.get("var_abs"), False),
                    (row.get("var_pct"), bool(row.get("var_pct_nm"))),
                ]
            )
        for col_idx, (value, percent_nm) in enumerate(value_specs, start=1):
            cell = table.cell(row_idx, col_idx)
            _pptx_set_cell_fill(cell, bg)
            _pptx_set_cell_border(cell)
            if comparative and col_idx == len(value_specs):
                texto = "n.m." if percent_nm or pd.isna(value) else _pptx_format_percent(value)
                negative = False if percent_nm or pd.isna(value) else _pptx_value_is_negative(value)
            else:
                texto = _pptx_format_number(value)
                negative = _pptx_value_is_negative(value)
            _pptx_set_cell_text(
                cell,
                texto,
                font_size=body_font_size,
                color=_PPTX_RED if negative else _PPTX_BLACK,
                bold=bool(row.get("is_parent")),
                align=PP_ALIGN.RIGHT,
                margin=0.012,
            )


def _pptx_add_missing_panel(slide: Any, x: float, y: float, w: float, h: float) -> None:
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    _pptx_set_shape_fill(box, _PPTX_OFF_WHITE)
    _pptx_set_no_line(box)
    box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = _pptx_clear_text_frame(box.text_frame)
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = "Demonstração não disponível para os filtros selecionados."
    run.font.name = _PPTX_FONT
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = _pptx_rgb(_PPTX_DARK_GRAY)


def _pptx_rows_or_none(
    demonstrativos: dict[str, pd.DataFrame],
    sigla: str,
    *,
    periodo_atual: str,
    periodo_anterior: str | None,
) -> list[dict[str, Any]] | None:
    df = demonstrativos.get(sigla)
    if not isinstance(df, pd.DataFrame) or df.empty:
        return None
    return _pptx_build_rows_for_demonstrativo(
        df,
        periodo_atual=periodo_atual,
        periodo_anterior=periodo_anterior,
        block_sigla=sigla,
    )


def _pptx_add_bp_dre_combined_slide(
    prs: Presentation,
    *,
    demonstrativos: dict[str, pd.DataFrame],
    periodo_atual: str,
    periodo_anterior: str | None,
    current_label: str,
    previous_label: str | None,
    subtitle: str,
    chip_label: str,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _pptx_add_content_header(slide, "BP + DRE — visão comparativa", subtitle, chip_label)

    headers = _pptx_headers(current_label, previous_label)
    left_x, right_x = 0.4, 6.78
    panel_w = 5.98
    title_y = 1.22
    table_y = 1.55
    table_h = 5.34
    _pptx_add_textbox(
        slide,
        "BP — Balanço Patrimonial",
        left_x,
        title_y,
        panel_w,
        0.22,
        font_size=11.5,
        bold=True,
        color=_PPTX_BLACK,
    )
    _pptx_add_textbox(
        slide,
        "DRE — Demonstração do Resultado",
        right_x,
        title_y,
        panel_w,
        0.22,
        font_size=11.5,
        bold=True,
        color=_PPTX_BLACK,
    )
    divider = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.56), Inches(1.22), Inches(0.01), Inches(5.72))
    _pptx_set_shape_fill(divider, _PPTX_LIGHT_GRAY)
    _pptx_set_no_line(divider)

    bp_rows = _pptx_rows_or_none(
        demonstrativos,
        "BP",
        periodo_atual=periodo_atual,
        periodo_anterior=periodo_anterior,
    )
    dre_rows = _pptx_rows_or_none(
        demonstrativos,
        "DRE",
        periodo_atual=periodo_atual,
        periodo_anterior=periodo_anterior,
    )
    account_chars = 34 if previous_label else 48
    if bp_rows:
        _pptx_add_compact_table_to_slide(
            slide,
            x=left_x,
            y=table_y,
            w=panel_w,
            h=table_h,
            headers=headers,
            rows=bp_rows,
            account_max_chars=account_chars,
        )
    else:
        _pptx_add_missing_panel(slide, left_x, table_y, panel_w, table_h)
    if dre_rows:
        _pptx_add_compact_table_to_slide(
            slide,
            x=right_x,
            y=table_y,
            w=panel_w,
            h=table_h,
            headers=headers,
            rows=dre_rows,
            account_max_chars=account_chars,
        )
    else:
        _pptx_add_missing_panel(slide, right_x, table_y, panel_w, table_h)


def _pptx_add_single_slide_demonstrativo(
    prs: Presentation,
    *,
    title: str,
    subtitle: str,
    chip_label: str,
    headers: list[str],
    rows: list[dict[str, Any]],
    account_max_chars: int = 96,
) -> None:
    header_h = 0.34
    table_h_max = 5.72
    row_h = _pptx_fit_single_slide_row_height(rows, table_h_max, header_h, 0.22)
    rows_one_slide = [dict(row, _height=row_h) for row in rows]
    _pptx_add_table_slide(
        prs,
        title=title,
        subtitle=subtitle,
        chip_label=chip_label,
        headers=headers,
        rows=rows_one_slide,
        chunk_index=1,
        chunk_total=1,
        account_max_chars=account_max_chars,
        body_font_size=7.5,
        header_font_size=8.0,
        table_h_max=table_h_max,
    )


def _pptx_add_missing_demonstrativo_slide(
    prs: Presentation,
    *,
    title: str,
    subtitle: str,
    chip_label: str,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _pptx_add_content_header(slide, title, subtitle, chip_label)
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(1.75), Inches(12.53), Inches(1.1))
    _pptx_set_shape_fill(box, _PPTX_OFF_WHITE)
    _pptx_set_no_line(box)
    paragraph = _pptx_clear_text_frame(box.text_frame)
    paragraph.alignment = PP_ALIGN.LEFT
    run = paragraph.add_run()
    run.text = "Demonstração não disponível para os filtros selecionados."
    run.font.name = _PPTX_FONT
    run.font.size = Pt(17)
    run.font.bold = True
    run.font.color.rgb = _pptx_rgb(_PPTX_DARK_GRAY)


def _pptx_find_metric_value(
    demonstrativos: dict[str, pd.DataFrame],
    *,
    siglas: tuple[str, ...],
    terms: tuple[str, ...],
    periodo_atual: str,
) -> Any:
    termos_norm = tuple(_normalize_text(term) for term in terms)
    for sigla in siglas:
        df = demonstrativos.get(sigla)
        if not isinstance(df, pd.DataFrame) or df.empty or periodo_atual not in df.columns:
            continue
        for _, row in df.iterrows():
            desc_norm = _normalize_text(row.get("descricao") or "")
            if all(term in desc_norm for term in termos_norm):
                valor = pd.to_numeric(row.get(periodo_atual), errors="coerce")
                if pd.notna(valor):
                    return valor
    return pd.NA


def _pptx_add_summary_slide(
    prs: Presentation,
    *,
    demonstrativos: dict[str, pd.DataFrame],
    periodo_atual: str,
    referencia_label: str,
    unidade_label: str,
    chip_label: str,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _pptx_add_content_header(
        slide,
        "Sumário executivo",
        f"Referência: {referencia_label} | Unidade: {unidade_label}",
        chip_label,
    )

    metrics = [
        ("Ativo Total", ("BP",), ("ATIVO",)),
        ("Operações de Crédito", ("BP", "DRE"), ("OPERACOES DE CREDITO",)),
        ("Provisão para Perdas Esperadas", ("BP", "DRE"), ("PERDAS ESPERADAS",)),
        ("Patrimônio Líquido", ("BP",), ("PATRIMONIO LIQUIDO",)),
    ]
    cards: list[tuple[str, Any]] = []
    for label, siglas, terms in metrics:
        valor = _pptx_find_metric_value(demonstrativos, siglas=siglas, terms=terms, periodo_atual=periodo_atual)
        cards.append((label, valor))

    for idx, (label, valor) in enumerate(cards):
        x = 0.48 + idx * 3.18
        y = 1.55
        card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(2.88), Inches(0.96))
        _pptx_set_shape_fill(card, _PPTX_OFF_WHITE)
        _pptx_set_no_line(card)
        _pptx_add_textbox(slide, label.upper(), x + 0.16, y + 0.14, 2.55, 0.18, font_size=8.5, color=_PPTX_DARK_GRAY)
        value_text = _pptx_format_number(valor) if pd.notna(valor) else "n.d."
        _pptx_add_textbox(
            slide,
            value_text,
            x + 0.16,
            y + 0.39,
            2.55,
            0.32,
            font_size=18,
            color=_PPTX_RED if pd.notna(valor) and _pptx_value_is_negative(valor) else _PPTX_BLACK,
            bold=True,
        )


def build_balanco_dre_dmpl_individual_pptx(
    *,
    institution_name: str,
    cnpj: str,
    periodo_atual: str,
    periodo_anterior: str | None = None,
    referencia: str | None = None,
    remessa: str | None = None,
    unidade: str | None = None,
    demonstrativos: dict[str, pd.DataFrame],
    column_labels: dict[str, str] | None = None,
    document_periods: list[str] | None = None,
) -> bytes:
    periodo_atual = str(periodo_atual or "").strip()
    if not periodo_atual:
        raise ValueError("periodo_atual é obrigatório para gerar o PPTX.")
    periodo_anterior = str(periodo_anterior or "").strip() or None
    demonstrativos_norm = _pptx_normalize_demonstrativos(demonstrativos)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    current_label = _pptx_period_label(periodo_atual, column_labels)
    previous_label = _pptx_period_label(periodo_anterior, column_labels) if periodo_anterior else None
    referencia_label = _pptx_reference_label(referencia)
    unidade_label = _pptx_unit_label(unidade)
    cnpj_digits = _digits_only(cnpj)
    cnpj_label = cnpj_digits or str(cnpj or "").strip() or "N/D"
    institution_label = str(institution_name or "").strip() or "Instituição não informada"
    remessa_label = str(remessa or "").strip() or "N/D"
    document_periods_label = ", ".join(str(item) for item in (document_periods or []) if str(item).strip())

    cover = prs.slides.add_slide(prs.slide_layouts[6])
    _pptx_set_slide_background(cover, _PPTX_BLACK)
    orange_bar = cover.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.82), Inches(2.12), Inches(0.055), Inches(1.28))
    _pptx_set_shape_fill(orange_bar, _PPTX_ORANGE)
    _pptx_set_no_line(orange_bar)
    _pptx_add_textbox(cover, "Balanço, DRE e DMPL", 1.05, 2.02, 8.8, 0.62, font_size=42, color=_PPTX_WHITE, bold=True)
    _pptx_add_textbox(cover, institution_label, 1.07, 2.82, 9.0, 0.34, font_size=18, color=_PPTX_LIGHT_GRAY)
    _pptx_add_textbox(cover, f"CNPJ: {cnpj_label}", 1.07, 3.18, 5.2, 0.27, font_size=14, color=_PPTX_LIGHT_GRAY)
    _pptx_add_textbox(cover, f"Período: {current_label} | Referência: {referencia_label}", 1.07, 3.5, 6.8, 0.27, font_size=14, color=_PPTX_LIGHT_GRAY)
    _pptx_add_textbox(cover, "Banco Central do Brasil — Documento 9011", 1.07, 4.05, 5.3, 0.25, font_size=11, color=_PPTX_DARK_GRAY)
    _pptx_add_textbox(cover, "Toma Conta | Análise Institucional", 9.45, 6.88, 3.35, 0.22, font_size=10.5, color=_PPTX_LIGHT_GRAY, align=PP_ALIGN.RIGHT)

    _pptx_add_summary_slide(
        prs,
        demonstrativos=demonstrativos_norm,
        periodo_atual=periodo_atual,
        referencia_label=referencia_label,
        unidade_label=unidade_label,
        chip_label=current_label,
    )

    block_order = [
        ("BP", "BP — Balanço Patrimonial"),
        ("DRE", "DRE — Demonstração do Resultado"),
        ("DMPL", "DMPL — Demonstração das Mutações do Patrimônio Líquido"),
        ("DRA", "DRA — Demonstração do Resultado Abrangente"),
        ("DFC", "DFC — Demonstração dos Fluxos de Caixa"),
    ]
    subtitle_parts = [f"Unidade: {unidade_label}", f"Referência: {referencia_label}", f"Remessa: {remessa_label}"]
    if document_periods_label:
        subtitle_parts.append(f"Competências: {document_periods_label}")
    subtitle = " | ".join(subtitle_parts)

    _pptx_add_bp_dre_combined_slide(
        prs,
        demonstrativos=demonstrativos_norm,
        periodo_atual=periodo_atual,
        periodo_anterior=periodo_anterior,
        current_label=current_label,
        previous_label=previous_label,
        subtitle=subtitle,
        chip_label=current_label,
    )

    for sigla, title in block_order:
        if sigla in {"BP", "DRE"}:
            continue
        df = demonstrativos_norm.get(sigla)
        if not isinstance(df, pd.DataFrame) or df.empty:
            _pptx_add_missing_demonstrativo_slide(prs, title=title, subtitle=subtitle, chip_label=current_label)
            continue
        rows = _pptx_build_rows_for_demonstrativo(
            df,
            periodo_atual=periodo_atual,
            periodo_anterior=periodo_anterior,
            block_sigla=sigla,
        )
        if not rows:
            _pptx_add_missing_demonstrativo_slide(prs, title=title, subtitle=subtitle, chip_label=current_label)
            continue
        headers = _pptx_headers(current_label, previous_label if periodo_anterior else None)
        if sigla == "DFC":
            _pptx_add_single_slide_demonstrativo(
                prs,
                title=title,
                subtitle=subtitle,
                chip_label=current_label,
                headers=headers,
                rows=rows,
                account_max_chars=86 if periodo_anterior else 118,
            )
            continue
        chunks = _pptx_chunk_table_rows(rows, comparative=bool(periodo_anterior))
        for idx, chunk in enumerate(chunks, start=1):
            _pptx_add_table_slide(
                prs,
                title=title,
                subtitle=subtitle,
                chip_label=current_label,
                headers=headers,
                rows=chunk,
                chunk_index=idx,
                chunk_total=len(chunks),
                account_max_chars=82 if periodo_anterior else 110,
            )

    total_slides = len(prs.slides)
    for idx, slide in enumerate(prs.slides, start=1):
        if idx == 1:
            continue
        _pptx_add_footer(slide, idx, total_slides)

    output = BytesIO()
    prs.save(output)
    return output.getvalue()


def build_excel_export_cdsfn(
    metadata: dict[str, Any],
    df_long: pd.DataFrame,
    df_wide: pd.DataFrame,
) -> bytes:
    output = BytesIO()
    metadata_rows = [
        {"campo": "cnpj", "valor": metadata.get("cnpj")},
        {"campo": "codigo_documento", "valor": metadata.get("codigo_documento")},
        {"campo": "tipo_remessa", "valor": metadata.get("tipo_remessa")},
        {"campo": "unidade_medida", "valor": metadata.get("unidade_medida")},
        {"campo": "data_base", "valor": metadata.get("data_base")},
    ]
    for item in metadata.get("datas_base_referencia", []):
        metadata_rows.append({"campo": f"dt_ref_{item.get('id')}", "valor": item.get("data")})

    with pd.ExcelWriter(output, engine="xlsxwriter") as writer:
        pd.DataFrame(metadata_rows).to_excel(writer, sheet_name="Metadata", index=False)
        df_long.to_excel(writer, sheet_name="Long", index=False)
        df_wide.to_excel(writer, sheet_name="Wide", index=False)
    return output.getvalue()
