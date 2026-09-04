"""Export dos painéis do SCR para PPTX com gráficos nativos do Office.

Nativo importa: o gráfico sai como ``<c:lineChart>`` de verdade, com os dados
embutidos numa planilha do próprio arquivo. Quem receber o deck consegue mudar
cor, escala, série e até os números sem voltar aqui — o que uma imagem colada
não permite.

Layout: quatro painéis por slide, em quadrantes de tamanho igual, no formato dos
decks. Cada quadrante tem título, subtítulo, fonte e o gráfico.

Duas exigências de leitura que o módulo garante:

* **Rótulo apenas no último período.** Os demais pontos ficam sem rótulo. Isso é
  feito ponto a ponto, com ``showVal`` no último — o valor continua vindo do
  dado, então segue correto se alguém editar a série no Office.
* **Percentual com duas casas** no eixo de valores e nos rótulos, via
  ``numFmt`` ``0.00%``. Os valores são gravados como fração (0,0463), que é o
  que o formato percentual do Office espera.
"""

from __future__ import annotations

from io import BytesIO
from typing import Any, Dict, List, Optional, Sequence, Tuple

import pandas as pd
from lxml import etree
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import (
    XL_CHART_TYPE,
    XL_DATA_LABEL_POSITION,
    XL_LEGEND_POSITION,
    XL_TICK_MARK,
)
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

from .sgs_credit_analytics import (
    ITAU_BLACK,
    cor_do_rotulo,
    eixo_datas_adaptativo,
    escala_de_eixo,
    formatar_competencia,
)

# Slide 16:9.
SLIDE_LARGURA = Inches(13.333)
SLIDE_ALTURA = Inches(7.5)

MARGEM = Inches(0.38)
GUTTER_H = Inches(0.30)
GUTTER_V = Inches(0.34)
TOPO_CONTEUDO = Inches(0.30)

ALTURA_TITULO = Inches(0.30)
ALTURA_SUBTITULO = Inches(0.24)
ALTURA_FONTE = Inches(0.20)

FONTE_TITULO_PT = 13
FONTE_SUBTITULO_PT = 10
FONTE_FONTE_PT = 7.5
FONTE_EIXO_PT = 8
# O eixo dos meses tem doze marcas e é o primeiro a ficar apinhado; um
# ponto a menos que o eixo de valores resolve sem perder leitura.
FONTE_EIXO_CATEGORIA_PT = 7
# A legenda vai no mesmo corpo do eixo: ela ocupa a faixa da direita, e cada
# ponto a mais ali é largura que sai do gráfico — a ponto de os meses do eixo
# ficarem colados um no outro.
FONTE_LEGENDA_PT = 7
FONTE_ROTULO_PT = 8
# Na coluna empilhada o rótulo tem que caber na largura da barra: com 12
# meses num gráfico de meia lâmina, a barra tem cerca de 0,34 in e um valor
# de seis dígitos a 8 pt passa disso, e o PowerPoint corta as pontas.
FONTE_ROTULO_COLUNA_PT = 7
# Abaixo desta fatia do total, o rótulo não cabe na altura da barra.
PARTICIPACAO_MINIMA_ROTULO = 0.07
# Acima disto a barra fica mais estreita que o próprio valor e o rótulo vai
# para a vertical, que é o que evita o corte no meio do número.
MAXIMO_COLUNAS_ROTULO_DEITADO = 14
# Muitas séries no mesmo gráfico: o rótulo encolhe para os nomes caberem
# empilhados na faixa da direita sem se sobrepor.
FONTE_ROTULO_DENSO_PT = 6.5

COR_TITULO = RGBColor(0x11, 0x11, 0x11)
COR_SUBTITULO = RGBColor(0x3C, 0x3C, 0x3C)
COR_FONTE = RGBColor(0x8F, 0x8F, 0x8F)
COR_EIXO = RGBColor(0x6F, 0x6F, 0x6F)
COR_GRADE = RGBColor(0xE6, 0xE6, 0xE6)

FORMATO_PERCENTUAL = "0.00%"
LARGURA_LINHA_PT = 1.75
LARGURA_LINHA_TOTAL_PT = 2.25

PAINEIS_POR_SLIDE = 4


def rotulo_mes(data_base: str) -> str:
    """``2026-06`` -> ``Jun/26``, no padrão dos eixos do deck."""
    texto = str(data_base)
    try:
        ano, mes = texto.split("-")[:2]
        return formatar_competencia(pd.Timestamp(year=int(ano), month=int(mes), day=1))
    except (ValueError, IndexError):
        return texto


def _hex_para_rgb(cor: str) -> RGBColor:
    limpo = str(cor).lstrip("#")
    return RGBColor(int(limpo[0:2], 16), int(limpo[2:4], 16), int(limpo[4:6], 16))


def _caixa_texto(
    slide,
    *,
    left: Emu,
    top: Emu,
    width: Emu,
    height: Emu,
    texto: str,
    tamanho: float,
    cor: RGBColor,
    negrito: bool = False,
    entrelinha: Optional[float] = None,
    espaco_entre_paragrafos: float = 0.0,
):
    """Caixa de texto sem preenchimento e sem contorno.

    Um texto com linhas em branco vira vários parágrafos, para o bloco de
    leitura dos dados sair no deck como sai na tela.
    """
    caixa = slide.shapes.add_textbox(left, top, width, height)
    quadro = caixa.text_frame
    quadro.word_wrap = True
    quadro.margin_left = quadro.margin_right = 0
    quadro.margin_top = quadro.margin_bottom = 0
    blocos = [parte.strip() for parte in str(texto).split("\n\n") if parte.strip()] or [""]
    for indice, bloco in enumerate(blocos):
        paragrafo = quadro.paragraphs[0] if indice == 0 else quadro.add_paragraph()
        paragrafo.alignment = PP_ALIGN.LEFT
        paragrafo.line_spacing = entrelinha or 1.0
        paragrafo.space_before = Pt(espaco_entre_paragrafos)
        paragrafo.space_after = Pt(0)
        corrida = paragrafo.add_run()
        corrida.text = bloco
        corrida.font.size = Pt(tamanho)
        corrida.font.bold = negrito
        corrida.font.color.rgb = cor
    return caixa


def _definir_tick_lbl_pos(eixo, valor: str) -> None:
    """Onde os rótulos do eixo de categorias são desenhados."""
    elemento = eixo._element
    no = elemento.find(qn("c:tickLblPos"))
    if no is None:
        no = etree.SubElement(elemento, qn("c:tickLblPos"))
    no.set("val", valor)
    _ordenar_filhos(elemento, ORDEM_CAT_AX)


def _definir_intervalo_rotulos(category_axis, total_categorias: int) -> int:
    """Mantém todas as categorias; os meses intermediários têm rótulo vazio."""
    intervalo = 1
    elemento = category_axis._element
    for tag in ("c:tickLblSkip", "c:tickMarkSkip"):
        no = elemento.find(qn(tag))
        if no is None:
            no = etree.SubElement(elemento, qn(tag))
        no.set("val", str(intervalo))
    _ordenar_filhos(elemento, ORDEM_CAT_AX)
    return intervalo


# Ordem exigida para os filhos de <c:dLbl> (CT_DLbl).
ORDEM_DLBL = (
    "c:idx", "c:delete", "c:layout", "c:tx", "c:numFmt", "c:spPr", "c:txPr",
    "c:dLblPos", "c:showLegendKey", "c:showVal", "c:showCatName",
    "c:showSerName", "c:showPercent", "c:showBubbleSize", "c:separator",
    "c:extLst",
)


def _rotular_apenas_ultimo_ponto(
    serie,
    indice_ultimo: int,
    formato_numero: str = FORMATO_PERCENTUAL,
    com_nome: bool = True,
) -> None:
    """Liga o rótulo só no último ponto da série.

    Usa ``showVal`` em vez de texto literal: o rótulo continua ligado ao dado,
    então quem editar a série no Office vê o número acompanhar.
    """
    if indice_ultimo < 0:
        return
    dLbl = serie.points[indice_ultimo].data_label._get_or_add_dLbl()

    formato = dLbl.find(qn("c:numFmt"))
    if formato is None:
        formato = etree.Element(qn("c:numFmt"))
        dLbl.insert(0, formato)
    formato.set("formatCode", formato_numero)
    formato.set("sourceLinked", "0")

    # Nome da série junto do valor: o rótulo do último ponto identifica a linha
    # sozinho, como na tela, e continua ligado ao dado.
    mostrar = {
        "c:showLegendKey": "0",
        "c:showVal": "1",
        "c:showCatName": "0",
        "c:showSerName": "1" if com_nome else "0",
        "c:showPercent": "0",
        "c:showBubbleSize": "0",
    }
    for tag, valor in mostrar.items():
        no = dLbl.find(qn(tag))
        if no is None:
            no = etree.SubElement(dLbl, qn(tag))
        no.set("val", valor)
    _ordenar_filhos(dLbl, ORDEM_DLBL)


def _rotular_todos_os_pontos(
    serie, total: int, formato_numero: str = FORMATO_PERCENTUAL,
    com_nome: bool = True,
) -> None:
    for indice in range(total):
        _rotular_apenas_ultimo_ponto(serie, indice, formato_numero, com_nome)


def _layout_do_rotulo(rotulo, x: float, y: float) -> None:
    """Fixa a posição de um rótulo dentro do quadro do gráfico.

    ``x`` e ``y`` são frações do gráfico, medidas da borda. É o que permite
    empurrar todos os rótulos finais para a faixa à direita da área de
    plotagem, escalonados, em vez de deixá-los sobre as próprias linhas.
    """
    elemento = rotulo._dLbl
    existente = elemento.find(qn("c:layout"))
    if existente is not None:
        elemento.remove(existente)
    layout = etree.Element(qn("c:layout"))
    elemento.insert(0, layout)
    manual = etree.SubElement(layout, qn("c:manualLayout"))
    _definir(manual, "c:xMode", "edge")
    _definir(manual, "c:yMode", "edge")
    for tag, valor in (("c:x", round(x, 4)), ("c:y", round(y, 4))):
        no = etree.SubElement(manual, qn(tag))
        no.set("val", str(valor))
    _ordenar_filhos(elemento, ORDEM_DLBL)


# Fração da largura reservada, à direita, para o rótulo do último ponto.
GUTTER_ROTULO_LINHA = 0.28
GUTTER_ROTULO_LINHA_MINIMO = 0.16
# Piso da faixa de rótulo e largura dos números do eixo da direita, ambos em
# polegadas: é o que a régua e o nome da série ocupam de fato no papel.
LARGURA_MINIMA_FAIXA_IN = 0.95
LARGURA_REGUA_DIREITA_IN = 0.62


def gutter_para(
    finais: Sequence[Tuple[str, float]], largura_emu: int, rotulo_fn=None,
    recuo_in: float = 0.0,
) -> float:
    """Faixa à direita: a régua do eixo mais o maior nome de série.

    Gráfico com nomes curtos — "PF inad; 7,8%" — não precisa de um terço da
    largura reservado, e essa sobra vinha do eixo dos meses, que ficava
    apinhado. A conta é em polegadas: numa lâmina inteira, uma fração fixa
    reservava quatro polegadas para um rótulo de uma.

    O nome entra inteiro quando cabe numa linha dentro do teto da faixa. Quando
    não cabe, entra pela metade, que é o que ele ocupa quebrado em duas — o
    Office quebra o rótulo só na borda do gráfico, então a faixa precisa
    comportar a largura que ele vai mesmo ter.
    """
    if not finais:
        return GUTTER_ROTULO_LINHA_MINIMO
    rotulo_fn = rotulo_fn or (lambda nome: nome)
    corpo = corpo_do_rotulo(len(finais))
    maior = max(len(f"{rotulo_fn(nome)}; 00,0%") for nome, _ in finais)
    largura_in = max(largura_emu / 914400.0, 0.5)
    uma_linha_in = maior * (LARGURA_MEDIA_CARACTERE * corpo / 72)
    teto_in = GUTTER_ROTULO_LINHA * largura_in - recuo_in
    necessario_in = uma_linha_in if uma_linha_in <= teto_in else uma_linha_in / 2
    total_in = max(necessario_in, LARGURA_MINIMA_FAIXA_IN) + recuo_in
    return min(GUTTER_ROTULO_LINHA, total_in / largura_in)


def corpo_do_rotulo(total_series: int) -> float:
    """Corpo do rótulo final: encolhe quando há muita série para empilhar."""
    return FONTE_ROTULO_PT if total_series <= 5 else FONTE_ROTULO_DENSO_PT


def _escalonar_no_gutter(
    finais: Sequence[Tuple[str, float]],
    altura_grafico_emu: int,
    largura_grafico_emu: int,
    rotulo_fn=None,
    gutter: float = GUTTER_ROTULO_LINHA,
    secundarias: Sequence[str] = (),
    escalas: Optional[Dict[bool, Tuple[float, float]]] = None,
) -> Dict[str, float]:
    """Posição vertical de cada rótulo na faixa da direita, sem sobreposição.

    Mesma regra da tela: parte da altura do próprio valor, empurra para baixo
    quem colidiria com o vizinho e recentraliza o conjunto. A altura de cada
    rótulo acompanha o número de linhas que o nome da série ocupa na faixa —
    "Financiamento à exportação" quebra em duas e precisa do dobro do espaço.
    """
    if not finais:
        return {}
    rotulo_fn = rotulo_fn or (lambda nome: nome)
    altura_in = max(altura_grafico_emu / 914400.0, 0.5)
    largura_faixa_in = largura_grafico_emu / 914400.0 * gutter
    corpo = corpo_do_rotulo(len(finais))
    # A caixa do rótulo tem margem própria e quebra antes do que a largura
    # bruta sugere; o fator segura a estimativa do lado conservador, porque
    # subestimar a altura de um rótulo faz ele entrar no de baixo.
    por_linha = max(
        int(largura_faixa_in * 0.58 / (LARGURA_MEDIA_CARACTERE * corpo / 72)), 8
    )
    unidade = (corpo * FATOR_LINHA_CALIBRI / 72) / altura_in

    def alturas(nome: str) -> float:
        texto = f"{rotulo_fn(nome)}; 00,0%"
        return max(1, -(-len(texto) // por_linha)) * unidade

    topo, base = 0.04, 0.04 + ALTURA_PLOT_LINHA
    # Cada série é lida na régua do próprio eixo, a mesma que o gráfico usa.
    # Normalizar pelo intervalo dos valores finais jogava o rótulo do eixo da
    # direita para o rodapé, longe da linha que ele nomeia.
    secundarias = set(secundarias)
    escalas = dict(escalas or {})
    for eh_secundaria in (False, True):
        if eh_secundaria in escalas:
            continue
        grupo = [
            valor for nome, valor in finais
            if (nome in secundarias) is eh_secundaria
        ]
        if grupo:
            escalas[eh_secundaria] = escala_de_eixo(grupo)

    def alvo(nome: str, valor: float) -> float:
        piso, teto = escalas[nome in secundarias]
        amplitude = (teto - piso) or 1.0
        fracao = min(max((teto - valor) / amplitude, 0.0), 1.0)
        return topo + fracao * (base - topo - unidade)

    alvos = {nome: alvo(nome, valor) for nome, valor in finais}

    ordem = sorted(alvos, key=lambda nome: alvos[nome])
    necessario = sum(alturas(nome) for nome in ordem)
    def empilhar(fator: float = 1.0) -> Dict[str, float]:
        """Encosta um rótulo no outro, de cima para baixo, na ordem das séries."""
        posicao, resultado = topo, {}
        for nome in ordem:
            resultado[nome] = posicao
            posicao += alturas(nome) * fator
        return resultado

    if necessario > base - topo:
        # Não cabe nem encostando um no outro: comprime na proporção da altura
        # de cada rótulo. Repartir a faixa em fatias iguais dava a um nome de
        # duas linhas o mesmo espaço de um de uma, e os dois se sobrepunham.
        return empilhar((base - topo) / necessario)

    posicoes = dict(alvos)
    for anterior, atual in zip(ordem, ordem[1:]):
        posicoes[atual] = max(posicoes[atual], posicoes[anterior] + alturas(anterior))
    fundo = posicoes[ordem[-1]] + alturas(ordem[-1])
    if fundo > base:
        folga_acima = posicoes[ordem[0]] - topo
        if fundo - base > folga_acima:
            # Não há de onde subir: encosta todos a partir do topo. Cabe, já
            # que a soma das alturas foi conferida acima.
            return empilhar()
        for nome in posicoes:
            posicoes[nome] -= fundo - base
    return posicoes


def _posicoes_escalonadas_rotulos(
    tabela: pd.DataFrame, ordem: Sequence[str]
) -> Dict[str, XL_DATA_LABEL_POSITION]:
    """Alterna acima/abaixo quando os pontos finais formam um aglomerado."""
    finais: List[Tuple[str, float]] = []
    for nome in ordem:
        validos = pd.to_numeric(tabela[nome], errors="coerce").dropna()
        if not validos.empty:
            finais.append((nome, float(validos.iloc[-1])))
    if len(finais) < 2:
        return {nome: XL_DATA_LABEL_POSITION.RIGHT for nome, _ in finais}

    valores = [valor for _, valor in finais]
    amplitude = max(valores) - min(valores)
    referencia = max(abs(valor) for valor in valores) or 1.0
    distancia_minima = max(amplitude * 0.10, referencia * 0.012)
    ordenados = sorted(finais, key=lambda item: item[1], reverse=True)
    posicoes = {nome: XL_DATA_LABEL_POSITION.RIGHT for nome, _ in ordenados}

    inicio = 0
    while inicio < len(ordenados):
        fim = inicio + 1
        while (
            fim < len(ordenados)
            and abs(ordenados[fim - 1][1] - ordenados[fim][1]) <= distancia_minima
        ):
            fim += 1
        grupo = ordenados[inicio:fim]
        if len(grupo) > 1:
            alternativas = (
                (XL_DATA_LABEL_POSITION.ABOVE, XL_DATA_LABEL_POSITION.RIGHT)
                if len(grupo) == 2
                else (
                    XL_DATA_LABEL_POSITION.ABOVE,
                    XL_DATA_LABEL_POSITION.RIGHT,
                    XL_DATA_LABEL_POSITION.BELOW,
                )
            )
            for posicao, (nome, _) in enumerate(grupo):
                posicoes[nome] = alternativas[posicao % len(alternativas)]
        inicio = fim
    return posicoes


def _estilizar_eixos(
    chart, total_categorias: int, formato_numero: str = FORMATO_PERCENTUAL,
    base_zero: bool = False,
    escala: Optional[Tuple[float, float]] = None,
) -> None:
    eixo_valor = chart.value_axis
    if escala is not None:
        # Escala fixa: o rótulo do último ponto é posicionado por coordenada e
        # precisa da mesma régua que o gráfico desenha.
        eixo_valor.minimum_scale, eixo_valor.maximum_scale = escala
    elif base_zero:
        # Mesma âncora da tela: barra medida a partir do zero, para a variação
        # aparecer na proporção certa.
        eixo_valor.minimum_scale = 0
    eixo_valor.tick_labels.number_format = formato_numero
    eixo_valor.tick_labels.number_format_is_linked = False
    eixo_valor.tick_labels.font.size = Pt(FONTE_EIXO_PT)
    eixo_valor.tick_labels.font.color.rgb = COR_EIXO
    eixo_valor.has_major_gridlines = True
    eixo_valor.major_gridlines.format.line.color.rgb = COR_GRADE
    eixo_valor.major_gridlines.format.line.width = Pt(0.5)
    eixo_valor.format.line.color.rgb = COR_GRADE
    eixo_valor.major_tick_mark = XL_TICK_MARK.NONE
    eixo_valor.minor_tick_mark = XL_TICK_MARK.NONE

    eixo_categoria = chart.category_axis
    eixo_categoria.tick_labels.font.size = Pt(FONTE_EIXO_CATEGORIA_PT)
    # Meses sempre no rodapé do gráfico. Com o padrão, uma série que cruza o
    # zero empurra os rótulos para o meio, por cima das próprias linhas.
    _definir_tick_lbl_pos(eixo_categoria, "low")
    eixo_categoria.tick_labels.font.color.rgb = COR_EIXO
    eixo_categoria.has_major_gridlines = False
    eixo_categoria.format.line.color.rgb = COR_GRADE
    eixo_categoria.major_tick_mark = XL_TICK_MARK.NONE
    eixo_categoria.minor_tick_mark = XL_TICK_MARK.NONE
    _definir_intervalo_rotulos(eixo_categoria, total_categorias)


# Fração da largura reservada, à direita, para o rótulo do último ponto. Sem
# ela o rótulo é desenhado dentro da área de plotagem, sobre as próprias
# linhas, e cortado na borda do gráfico.
# Na coluna empilhada o rótulo do último mês fica dentro da fatia, mas o
# texto sobra para a direita da barra e era cortado na borda do gráfico.
GUTTER_ROTULO_COLUNA = 0.24
GUTTER_LEGENDA_MAXIMO = 0.30
# Piso da faixa de legenda, em polegadas: o quadrado da cor mais um nome curto.
LARGURA_MINIMA_LEGENDA_IN = 1.05
# Altura da área de plotagem. O resto do quadro é o eixo de categorias e a
# legenda. A coluna precisa de mais folga embaixo: a legenda tem mais itens e
# entrava por cima dos meses.
ALTURA_PLOT_LINHA = 0.86
ALTURA_PLOT_COLUNA = 0.84


# Ordem exigida para os filhos de <c:legend> (CT_Legend).
ORDEM_LEGENDA = (
    "c:legendPos", "c:legendEntry", "c:layout", "c:overlay", "c:spPr",
    "c:txPr", "c:extLst",
)


def _layout_manual_da_legenda(
    legend, gutter: float, total_series: int, altura_emu: int
) -> None:
    """Prende a legenda à faixa da direita, na altura que as entradas pedem.

    Sem caixa explícita o Office desenha a legenda numa área própria, menor
    que o quadro, e corta as entradas que não couberem: o card de mix da
    carteira mostrava sete das dez séries, e as três escondidas eram
    justamente as de nome comprido. A caixa acompanha o número de séries e
    fica centrada — cobrindo a altura inteira, três entradas ficavam espalhadas
    de ponta a ponta.
    """
    altura_in = max(altura_emu / 914400.0, 0.5)
    linha = (FONTE_LEGENDA_PT * FATOR_LINHA_CALIBRI / 72) * 1.25
    altura = min(0.96, (total_series * linha + 0.08) / altura_in)
    elemento = legend._element
    existente = elemento.find(qn("c:layout"))
    if existente is not None:
        elemento.remove(existente)
    layout = etree.SubElement(elemento, qn("c:layout"))
    manual = etree.SubElement(layout, qn("c:manualLayout"))
    _definir(manual, "c:xMode", "edge")
    _definir(manual, "c:yMode", "edge")
    for tag, valor in (
        ("c:x", round(1.0 - gutter, 4)), ("c:y", round((1.0 - altura) / 2, 4)),
        ("c:w", round(gutter, 4)), ("c:h", round(altura, 4)),
    ):
        no = etree.SubElement(manual, qn(tag))
        no.set("val", str(valor))
    _ordenar_filhos(elemento, ORDEM_LEGENDA)


def _layout_manual_do_plot(plot_area, gutter: float, altura: float) -> None:
    """Encolhe a área de plotagem para abrir espaço de rótulo à direita."""
    existente = plot_area.find(qn("c:layout"))
    if existente is not None:
        plot_area.remove(existente)
    layout = etree.Element(qn("c:layout"))
    plot_area.insert(0, layout)
    manual = etree.SubElement(layout, qn("c:manualLayout"))
    _definir(manual, "c:layoutTarget", "inner")
    _definir(manual, "c:xMode", "edge")
    _definir(manual, "c:yMode", "edge")
    for tag, valor in (
        ("c:x", 0.06), ("c:y", 0.04),
        ("c:w", round(0.94 - gutter, 4)), ("c:h", altura),
    ):
        no = etree.SubElement(manual, qn(tag))
        no.set("val", str(valor))


def _proximo_ax_id(chart, deslocamento: int) -> int:
    """Id de eixo inédito dentro do gráfico."""
    existentes = {
        int(no.get("val"))
        for no in chart._chartSpace.iter(qn("c:axId"))
        if no.get("val")
    }
    candidato = (max(existentes) if existentes else 100_000_000) + deslocamento
    while candidato in existentes:
        candidato += 1
    return candidato


# Ordem exigida para os filhos de <c:chartSpace> (CT_ChartSpace). O spPr entra
# depois de <c:chart> e antes de <c:txPr>; anexado no fim, ficaria depois de
# <c:externalData> e o PowerPoint recusaria o arquivo.
ORDEM_CHART_SPACE = (
    "c:date1904", "c:lang", "c:roundedCorners", "c:style", "c:clrMapOvr",
    "c:pivotSource", "c:protection", "c:chart", "c:spPr", "c:txPr",
    "c:externalData", "c:printSettings", "c:userShapes", "c:extLst",
)


def _sem_preenchimento_nem_contorno(elemento, sequencia: Sequence[str]) -> None:
    """Aplica ``noFill`` no preenchimento e na linha, na posição do schema."""
    existente = elemento.find(qn("c:spPr"))
    if existente is not None:
        elemento.remove(existente)
    sp_pr = etree.SubElement(elemento, qn("c:spPr"))
    etree.SubElement(sp_pr, qn("a:noFill"))
    linha = etree.SubElement(sp_pr, qn("a:ln"))
    etree.SubElement(linha, qn("a:noFill"))
    _ordenar_filhos(elemento, sequencia)


# Ordem exigida pelo schema para os filhos de <c:ser> num gráfico de linhas
# (CT_LineSer). O PowerPoint recusa o arquivo inteiro se a sequência não bater.
ORDEM_LINE_SER = (
    "c:idx", "c:order", "c:tx", "c:spPr", "c:marker", "c:dPt", "c:dLbls",
    "c:trendline", "c:errBars", "c:cat", "c:val", "c:smooth", "c:extLst",
)

# Elementos que só existem em série de barra (CT_BarSer) e não podem viajar
# junto quando a série vira linha.
FILHOS_SO_DE_BARRA = ("c:invertIfNegative", "c:pictureOptions", "c:shape")

# Grupos de gráfico reconhecidos dentro de <c:plotArea>. O schema manda que
# TODOS venham antes de qualquer eixo.
GRUPOS_DE_GRAFICO = (
    "c:areaChart", "c:area3DChart", "c:lineChart", "c:line3DChart",
    "c:stockChart", "c:radarChart", "c:scatterChart", "c:pieChart",
    "c:pie3DChart", "c:doughnutChart", "c:barChart", "c:bar3DChart",
    "c:ofPieChart", "c:surfaceChart", "c:surface3DChart", "c:bubbleChart",
)


# Ordem exigida para os filhos de <c:catAx> (CT_CatAx). O python-pptx emite
# noMultiLvlLbl antes de tickLblSkip/tickMarkSkip, que só entram depois.
ORDEM_CAT_AX = (
    "c:axId", "c:scaling", "c:delete", "c:axPos", "c:majorGridlines",
    "c:minorGridlines", "c:title", "c:numFmt", "c:majorTickMark",
    "c:minorTickMark", "c:tickLblPos", "c:spPr", "c:txPr", "c:crossAx",
    "c:crosses", "c:crossesAt", "c:auto", "c:lblAlgn", "c:lblOffset",
    "c:tickLblSkip", "c:tickMarkSkip", "c:noMultiLvlLbl", "c:extLst",
)


def _ordenar_filhos(elemento, sequencia: Sequence[str]) -> None:
    """Reordena os filhos para a sequência que o schema exige.

    Fora de ordem, o PowerPoint recusa o arquivo inteiro e oferece reparo em
    vez de abrir o deck.
    """
    posicao = {qn(nome): indice for indice, nome in enumerate(sequencia)}
    for filho in sorted(elemento, key=lambda e: posicao.get(e.tag, len(posicao))):
        elemento.append(filho)


def _ordenar_line_ser(ser) -> None:
    """Reordena os filhos de uma série para a sequência de CT_LineSer."""
    for nome in FILHOS_SO_DE_BARRA:
        for filho in ser.findall(qn(nome)):
            ser.remove(filho)
    _ordenar_filhos(ser, ORDEM_LINE_SER)


def _indice_apos_grupos(plot_area) -> int:
    """Posição logo depois do último grupo de gráfico do plotArea."""
    tags = {qn(nome) for nome in GRUPOS_DE_GRAFICO}
    ultimo = 0
    for indice, filho in enumerate(plot_area):
        if filho.tag in tags:
            ultimo = indice + 1
    return ultimo


def _mover_para_eixo_secundario(
    chart, indices: Sequence[int], formato_numero: str, cor: RGBColor | None = None,
    escala: Optional[Tuple[float, float]] = None,
) -> bool:
    """Tira as séries de ``indices`` do gráfico de colunas e as põe em linha,
    num segundo eixo de valores à direita.

    Sem isto, um card de volume (R$ bi) mais prazo (meses) era achatado num
    eixo único: a linha de prazo, entre 23 e 34, virava uma reta no rodapé de
    um eixo que ia até 350.

    A montagem respeita a sequência do schema em dois pontos que o PowerPoint
    não perdoa: o ``<c:lineChart>`` entra antes de qualquer eixo, e os filhos
    de cada ``<c:ser>`` movida são reordenados para CT_LineSer.
    """
    if not indices:
        return False
    plot_area = chart._chartSpace.chart.plotArea
    # O grupo de origem é a coluna nos cards de volume e prazo, e a própria
    # linha quando uma das séries do gráfico de linhas vai para a direita.
    grupo = plot_area.find(qn("c:barChart"))
    if grupo is None:
        grupo = plot_area.find(qn("c:lineChart"))
    if grupo is None:
        return False
    bar_chart = grupo

    series = bar_chart.findall(qn("c:ser"))
    mover = [series[i] for i in indices if 0 <= i < len(series)]
    if not mover or len(mover) == len(series):
        return False

    id_categoria = _proximo_ax_id(chart, 1)
    id_valor = _proximo_ax_id(chart, 2)

    line_chart = etree.Element(qn("c:lineChart"))
    plot_area.insert(_indice_apos_grupos(plot_area), line_chart)
    agrupamento = etree.SubElement(line_chart, qn("c:grouping"))
    agrupamento.set("val", "standard")
    varia = etree.SubElement(line_chart, qn("c:varyColors"))
    varia.set("val", "0")
    for elemento in mover:
        bar_chart.remove(elemento)
        line_chart.append(elemento)
        # A série que vem de um gráfico de linhas já traz o marcador. Criar um
        # segundo deixa dois <c:marker> na mesma série e o PowerPoint recusa o
        # arquivo.
        marcador = elemento.find(qn("c:marker"))
        if marcador is None:
            marcador = etree.SubElement(elemento, qn("c:marker"))
        simbolo = marcador.find(qn("c:symbol"))
        if simbolo is None:
            simbolo = etree.SubElement(marcador, qn("c:symbol"))
        simbolo.set("val", "none")
        _ordenar_line_ser(elemento)
    marcador_grupo = etree.SubElement(line_chart, qn("c:marker"))
    marcador_grupo.set("val", "1")
    for identificador in (id_categoria, id_valor):
        no = etree.SubElement(line_chart, qn("c:axId"))
        no.set("val", str(identificador))

    # Eixo de categorias espelhado e oculto: o Office exige o par, mas quem lê
    # o slide só deve ver um eixo de meses.
    cat_ax = etree.SubElement(plot_area, qn("c:catAx"))
    _definir(cat_ax, "c:axId", str(id_categoria))
    escala_cat = etree.SubElement(cat_ax, qn("c:scaling"))
    _definir(escala_cat, "c:orientation", "minMax")
    _definir(cat_ax, "c:delete", "1")
    _definir(cat_ax, "c:axPos", "b")
    _definir(cat_ax, "c:crossAx", str(id_valor))

    val_ax = etree.SubElement(plot_area, qn("c:valAx"))
    _definir(val_ax, "c:axId", str(id_valor))
    escala_val = etree.SubElement(val_ax, qn("c:scaling"))
    _definir(escala_val, "c:orientation", "minMax")
    _definir(val_ax, "c:delete", "0")
    _definir(val_ax, "c:axPos", "r")
    formato = etree.SubElement(val_ax, qn("c:numFmt"))
    formato.set("formatCode", formato_numero)
    formato.set("sourceLinked", "0")
    _definir(val_ax, "c:majorTickMark", "none")
    _definir(val_ax, "c:minorTickMark", "none")
    _definir(val_ax, "c:tickLblPos", "nextTo")
    # Sem spPr o Office desenha o eixo secundário com a linha preta grossa do
    # estilo padrão, ao lado do eixo da esquerda, que é um fio cinza.
    sp_pr = etree.SubElement(val_ax, qn("c:spPr"))
    etree.SubElement(sp_pr, qn("a:noFill"))
    linha_eixo = etree.SubElement(sp_pr, qn("a:ln"))
    linha_eixo.set("w", str(Pt(0.5).emu // 1))
    preenchimento_linha = etree.SubElement(linha_eixo, qn("a:solidFill"))
    tinta = etree.SubElement(preenchimento_linha, qn("a:srgbClr"))
    tinta.set("val", str(COR_GRADE))
    if cor is not None:
        # Eixo da direita na cor da linha que ele mede: é o que diz ao leitor
        # que "meses" se lê à direita.
        val_ax.append(_texto_do_eixo(cor))
    _definir(val_ax, "c:crossAx", str(id_categoria))
    _definir(val_ax, "c:crosses", "max")
    # Ancorado em zero, como na tela: sem isso um prazo que varia 1 mês em 47
    # preenche a altura do gráfico igual a um que varia 11 em 27. O teto também
    # é fixo, porque o rótulo da série vai por coordenada nesta mesma régua.
    # CT_Scaling é sequência: logBase, orientation, max, min — nessa ordem.
    piso, teto = escala if escala is not None else (0.0, None)
    if teto is not None:
        escala_val.append(_no_com_valor("c:max", _texto_de_numero(teto)))
    escala_val.append(_no_com_valor("c:min", _texto_de_numero(piso)))
    return True


def _texto_de_numero(valor: float) -> str:
    """Número sem notação científica nem zeros à toa, como o Office grava."""
    return f"{valor:.10f}".rstrip("0").rstrip(".") or "0"


def _texto_do_eixo(cor: RGBColor):
    """<c:txPr> com a cor pedida, para tingir os rótulos de um eixo."""
    tx_pr = etree.Element(qn("c:txPr"))
    etree.SubElement(tx_pr, qn("a:bodyPr"))
    etree.SubElement(tx_pr, qn("a:lstStyle"))
    paragrafo = etree.SubElement(tx_pr, qn("a:p"))
    propriedades = etree.SubElement(paragrafo, qn("a:pPr"))
    fonte = etree.SubElement(propriedades, qn("a:defRPr"))
    fonte.set("sz", str(int(FONTE_EIXO_PT * 100)))
    fonte.set("b", "1")
    preenchimento = etree.SubElement(fonte, qn("a:solidFill"))
    valor = etree.SubElement(preenchimento, qn("a:srgbClr"))
    valor.set("val", str(cor))
    etree.SubElement(paragrafo, qn("a:endParaRPr"))
    return tx_pr


def _girar_rotulo(rotulo) -> None:
    """Põe o rótulo na vertical, lendo de baixo para cima.

    Com 27 categorias a barra fica mais estreita que o texto e o Office corta
    o rótulo no meio ("11," em vez de "11,42%"). Na vertical ele cabe.
    """
    tx_pr = rotulo._dLbl.find(qn("c:txPr"))
    if tx_pr is None:
        return
    corpo = tx_pr.find(qn("a:bodyPr"))
    if corpo is None:
        corpo = etree.Element(qn("a:bodyPr"))
        tx_pr.insert(0, corpo)
    corpo.set("rot", "-5400000")
    corpo.set("vert", "horz")


def _preencher_rotulo(rotulo, cor_hex: str) -> None:
    """Fundo sólido no rótulo, na cor da série, sem contorno."""
    elemento = rotulo._dLbl
    existente = elemento.find(qn("c:spPr"))
    if existente is not None:
        elemento.remove(existente)
    sp_pr = etree.SubElement(elemento, qn("c:spPr"))
    preenchimento = etree.SubElement(sp_pr, qn("a:solidFill"))
    valor = etree.SubElement(preenchimento, qn("a:srgbClr"))
    valor.set("val", str(cor_hex).lstrip("#").upper())
    linha = etree.SubElement(sp_pr, qn("a:ln"))
    etree.SubElement(linha, qn("a:noFill"))
    _ordenar_filhos(elemento, ORDEM_DLBL)


def _no_com_valor(tag: str, valor: str):
    no = etree.Element(qn(tag))
    no.set("val", valor)
    return no


def _definir(pai, tag: str, valor: str):
    no = etree.SubElement(pai, qn(tag))
    no.set("val", valor)
    return no


# Cabeçalho compacto do card no deck: título e unidade na mesma caixa, com
# espaço reservado para duas linhas. Alguns títulos passam de 70 caracteres e,
# com uma linha só, transbordavam sobre o gráfico. A fonte é igual em todos os
# cards e vive uma vez por seção, na faixa de leitura.
ALTURA_TITULO_COMPACTO = Inches(0.46)


def _adicionar_painel(
    slide,
    painel: Any,
    *,
    left: Emu,
    top: Emu,
    width: Emu,
    height: Emu,
    rotulo_serie_fn,
    cabecalho_compacto: bool = False,
) -> Dict[str, Any]:
    """Desenha um card: cabeçalho e gráfico nativo."""
    if cabecalho_compacto:
        unidade = str(getattr(painel, "subtitulo", "") or "").strip()
        rotulo = f"{painel.titulo}  ·  {unidade}" if unidade else painel.titulo
        _caixa_texto(
            slide, left=left, top=top, width=width, height=ALTURA_TITULO_COMPACTO,
            texto=rotulo, tamanho=FONTE_TITULO_PT, cor=COR_TITULO, negrito=True,
        )
        altura_cabecalho = ALTURA_TITULO_COMPACTO
    else:
        _caixa_texto(
            slide, left=left, top=top, width=width, height=ALTURA_TITULO,
            texto=painel.titulo, tamanho=FONTE_TITULO_PT, cor=COR_TITULO, negrito=True,
        )
        _caixa_texto(
            slide, left=left, top=top + ALTURA_TITULO, width=width, height=ALTURA_SUBTITULO,
            texto=painel.subtitulo, tamanho=FONTE_SUBTITULO_PT, cor=COR_SUBTITULO,
        )
        _caixa_texto(
            slide, left=left, top=top + ALTURA_TITULO + ALTURA_SUBTITULO,
            width=width, height=ALTURA_FONTE,
            texto=painel.fonte, tamanho=FONTE_FONTE_PT, cor=COR_FONTE,
        )
        altura_cabecalho = ALTURA_TITULO + ALTURA_SUBTITULO + ALTURA_FONTE

    topo_grafico = top + altura_cabecalho
    altura_grafico = height - altura_cabecalho

    # `data_base` e `serie` chegam como categóricas com TODAS as categorias da
    # série histórica. Sem virar texto antes, o pivot recria as data-bases
    # ausentes como linhas vazias e o eixo do gráfico ganha meses fantasma.
    plano = painel.series.copy()
    plano["data_base"] = plano["data_base"].astype(str)
    plano["serie"] = plano["serie"].astype(str)
    tabela = plano.pivot_table(
        index="data_base", columns="serie", values="valor", aggfunc="first",
        observed=True,
    ).sort_index()
    ordem_categorias = getattr(painel, "ordem_categorias", None)
    if ordem_categorias:
        presentes = [categoria for categoria in ordem_categorias if categoria in tabela.index]
        tabela = tabela.reindex(presentes)
    indices = [str(idx) for idx in tabela.index]
    meses_validos = all(
        len(indice) >= 7 and indice[4] == "-" and indice[5:7].isdigit()
        for indice in indices
    )
    if meses_validos:
        datas = [pd.Timestamp(f"{indice[:7]}-01") for indice in indices]
        selecionadas, _ = eixo_datas_adaptativo(datas)
        meses_selecionados = {data.strftime("%Y-%m") for data in selecionadas}
        categorias = [
            rotulo_mes(indice) if indice[:7] in meses_selecionados else "\u00a0"
            for indice in indices
        ]
    else:
        categorias = indices

    dados = CategoryChartData()
    dados.categories = categorias
    ordem = [s for s in painel.ordem_series if s in tabela.columns]
    for nome in ordem:
        coluna = tabela[nome]
        valores = [None if pd.isna(v) else float(v) for v in coluna]
        dados.add_series(rotulo_serie_fn(nome), valores)

    tipo_grafico = getattr(painel, "tipo_grafico", "line")
    if tipo_grafico == "column_stacked":
        chart_type = XL_CHART_TYPE.COLUMN_STACKED
    elif tipo_grafico == "column_line":
        # Entra como colunas; as séries do eixo secundário viram linha logo
        # depois, com o segundo eixo de valores.
        chart_type = XL_CHART_TYPE.COLUMN_CLUSTERED
    else:
        chart_type = XL_CHART_TYPE.LINE
    formato_numero = getattr(painel, "formato_numero", FORMATO_PERCENTUAL)
    series_secundarias = set(getattr(painel, "series_secundarias", ()) or ())
    grafico = slide.shapes.add_chart(
        chart_type, left, topo_grafico, width, altura_grafico, dados
    ).chart
    grafico.has_title = False
    # Sem o spPr explícito, o Office desenha o contorno do estilo padrão em
    # volta do gráfico. O deck não tem moldura em lugar nenhum.
    _sem_preenchimento_nem_contorno(grafico._chartSpace, ORDEM_CHART_SPACE)

    # Gráfico de linhas não leva legenda: o nome da série já vai no rótulo do
    # último ponto, e a legenda embaixo disputava espaço com os meses do eixo.
    # A coluna empilhada mantém a legenda, porque lá o rótulo leva só o valor.
    # Uma série só não precisa de legenda: o título do card já diz o que é, e
    # o Office ainda lista as categorias quando a legenda fica ligada.
    grafico.has_legend = tipo_grafico == "column_stacked" and len(ordem) > 1
    if grafico.has_legend:
        # Legenda à direita, na faixa que já fica reservada. Embaixo, um mix de
        # dez produtos ocupava quatro fileiras e esmagava as barras.
        grafico.legend.position = XL_LEGEND_POSITION.RIGHT
        grafico.legend.include_in_layout = False
        grafico.legend.font.size = Pt(FONTE_LEGENDA_PT)
        grafico.legend.font.color.rgb = COR_EIXO

    plot = grafico.plots[0]
    plot.has_data_labels = False
    if tipo_grafico == "column_stacked":
        # O PowerPoint recorta o rótulo pela largura da barra. Com o vão padrão
        # de 150% a barra fica com 0,15 in num gráfico de meia lâmina e o valor
        # sai cortado nas duas pontas; encostando as barras ela dobra de
        # largura e o rótulo cabe.
        plot.gap_width = 25
        plot.overlap = 100

    indice_ultimo = len(categorias) - 1
    posicoes_rotulos = _posicoes_escalonadas_rotulos(tabela, ordem)
    totais_por_categoria = (
        tabela[ordem].abs().sum(axis=1).replace(0, float("nan"))
        if tipo_grafico == "column_stacked" else None
    )
    # Valor final de cada série, para escalonar os rótulos na faixa da direita.
    finais: List[Tuple[str, float]] = []
    for nome in ordem:
        validos = pd.to_numeric(tabela[nome], errors="coerce").dropna()
        if not validos.empty:
            finais.append((nome, float(validos.iloc[-1])))
    finais_de_linha = [
        (nome, valor) for nome, valor in finais
        if tipo_grafico == "line" or nome in series_secundarias
    ]
    # Régua de cada eixo, calculada aqui e aplicada tanto ao eixo quanto à
    # coordenada do rótulo — os dois precisam ler a mesma escala.
    escalas: Dict[bool, Tuple[float, float]] = {}
    if tipo_grafico in {"line", "column_line"}:
        for eh_secundaria in (False, True):
            colunas = [
                nome for nome in ordem
                if (nome in series_secundarias) is eh_secundaria
            ]
            if not colunas:
                continue
            valores = pd.to_numeric(
                tabela[colunas].stack(), errors="coerce"
            ).dropna()
            if valores.empty:
                continue
            # Barra e eixo secundário ancoram no zero por natureza: a barra
            # codifica o valor no comprimento e a régua da direita foi
            # ancorada de propósito, para o prazo não encher a altura toda.
            # A linha do eixo primário fica com a regra automática.
            escalas[eh_secundaria] = escala_de_eixo(
                valores.tolist(),
                ancorar_zero=(
                    True if eh_secundaria or tipo_grafico == "column_line"
                    else None
                ),
            )
    # Com eixo secundário, a régua da direita ocupa a borda da área de
    # plotagem: o rótulo precisa começar depois dela, senão cai sobre os
    # próprios números do eixo. A régua ocupa uma largura em polegadas, não
    # uma fração do gráfico.
    largura_in = max(int(width) / 914400.0, 0.5)
    recuo_in = (
        min(0.10 * largura_in, LARGURA_REGUA_DIREITA_IN)
        if series_secundarias else 0.012 * largura_in
    )
    recuo_eixo = recuo_in / largura_in
    gutter_linha = gutter_para(
        finais_de_linha, int(width), rotulo_serie_fn, recuo_in=recuo_in
    )
    alturas_rotulo = (
        {} if tipo_grafico == "column_stacked"
        else _escalonar_no_gutter(
            finais_de_linha, int(altura_grafico), int(width), rotulo_serie_fn,
            gutter=gutter_linha, secundarias=series_secundarias,
            escalas=escalas,
        )
    )
    # O Office ancora o rótulo pelo centro da caixa, e não pela borda esquerda:
    # a coordenada é o meio da faixa, senão o nome mais longo cresce para a
    # esquerda e cai por cima dos números do eixo.
    x_rotulo = (
        0.06 + (0.94 - gutter_linha) + recuo_eixo
        + (gutter_linha - recuo_eixo) / 2
    )
    for posicao, nome in enumerate(ordem):
        serie = plot.series[posicao]
        serie.smooth = False
        cor = _hex_para_rgb(painel.cores.get(nome, "#8F8F8F"))
        eh_total = nome in painel.tracejadas
        eh_secundaria = nome in series_secundarias
        eh_coluna = tipo_grafico == "column_stacked" or (
            tipo_grafico == "column_line" and not eh_secundaria
        )
        if eh_coluna:
            serie.format.fill.solid()
            serie.format.fill.fore_color.rgb = cor
            serie.format.line.color.rgb = cor
        else:
            linha = serie.format.line
            linha.color.rgb = cor
            linha.width = Pt(LARGURA_LINHA_TOTAL_PT if eh_total else LARGURA_LINHA_PT)
            if eh_total:
                linha.dash_style = MSO_LINE_DASH_STYLE.DASH

        coluna = tabela[nome]
        ultimo_valido = indice_ultimo
        while ultimo_valido >= 0 and pd.isna(coluna.iloc[ultimo_valido]):
            ultimo_valido -= 1
        rotular_todos = bool(getattr(painel, "rotular_todos_pontos", False))
        # Fatia pequena demais não comporta o rótulo e o texto acaba por cima
        # da fatia vizinha. Mesma regra da tela, onde a trava de tamanho
        # uniforme esconde o que não cabe.
        participacao = None
        if eh_coluna and totais_por_categoria is not None:
            participacao = (
                pd.to_numeric(coluna, errors="coerce") / totais_por_categoria
            )
        # Nome da série no rótulo só onde ele cabe. Numa fatia de coluna
        # empilhada o texto é mais largo que a barra e vira uma pilha ilegível
        # à direita; lá o nome fica na legenda e o rótulo leva só o valor.
        com_nome = not eh_coluna
        if rotular_todos:
            _rotular_todos_os_pontos(serie, len(coluna), formato_numero, com_nome)
            indices_rotulados = [
                indice for indice, valor in enumerate(coluna) if pd.notna(valor)
            ]
        elif (
            participacao is not None
            and ultimo_valido >= 0
            and float(participacao.iloc[ultimo_valido] or 0) < PARTICIPACAO_MINIMA_ROTULO
        ):
            indices_rotulados = []
        else:
            _rotular_apenas_ultimo_ponto(
                serie, ultimo_valido, formato_numero, com_nome
            )
            indices_rotulados = [ultimo_valido] if ultimo_valido >= 0 else []
        # A posição do rótulo é validada pelo tipo do gráfico: coluna aceita
        # apenas ctr/inBase/inEnd/outEnd. Escrever "r" ou "t" numa série de
        # coluna produz um arquivo que o PowerPoint recusa inteiro, oferecendo
        # reparo — foi o que quebrou o deck de Concessões.
        for indice_rotulo in indices_rotulados:
            rotulo = serie.points[indice_rotulo].data_label
            if eh_coluna:
                rotulo.position = (
                    # No meio da fatia, e não na borda de cima: com a caixa
                    # centrada na borda, a fatia de baixo jogava metade do
                    # rótulo abaixo do eixo, por cima do nome do mês.
                    XL_DATA_LABEL_POSITION.CENTER
                    if tipo_grafico == "column_stacked"
                    else XL_DATA_LABEL_POSITION.OUTSIDE_END
                )
            elif not rotular_todos:
                rotulo.position = posicoes_rotulos.get(
                    nome, XL_DATA_LABEL_POSITION.RIGHT
                )
            if (
                not eh_coluna
                and nome in alturas_rotulo
                and indice_rotulo == ultimo_valido
            ):
                _layout_do_rotulo(rotulo, x_rotulo, alturas_rotulo[nome])
            rotulo.font.size = Pt(
                FONTE_ROTULO_COLUNA_PT if eh_coluna
                else corpo_do_rotulo(len(finais))
            )
            if eh_coluna and len(categorias) > MAXIMO_COLUNAS_ROTULO_DEITADO:
                _girar_rotulo(rotulo)
            rotulo.font.bold = True
            cor_serie = painel.cores.get(nome, "#8F8F8F")
            if eh_coluna:
                # Rótulo de barra ganha fundo sólido na cor da própria série,
                # com o texto na tinta que contrasta com ela. Fica ancorado na
                # série mesmo quando o valor cai fora da barra.
                _preencher_rotulo(rotulo, cor_serie)
                rotulo.font.color.rgb = _hex_para_rgb(cor_do_rotulo(cor_serie))
            else:
                # Rótulo de linha na cor da própria linha: é o que liga o valor
                # à série sem precisar de legenda.
                rotulo.font.color.rgb = _hex_para_rgb(cor_serie)

    _estilizar_eixos(
        grafico, len(categorias), formato_numero,
        base_zero=tipo_grafico in {"column_stacked", "column_line"},
        escala=escalas.get(False),
    )
    empilhado = tipo_grafico == "column_stacked"
    if empilhado and not grafico.has_legend:
        # Sem legenda a faixa da direita não serve para nada: o rótulo da barra
        # é desenhado dentro dela. A área de plotagem toma a largura inteira.
        gutter = 0.06
    elif empilhado:
        # A faixa da direita abriga a legenda: precisa caber o maior nome de
        # série, senão a legenda é desenhada por cima das barras.
        maior = max((len(rotulo_serie_fn(nome)) for nome in ordem), default=10)
        largura_in = max(int(width) / 914400.0, 0.5)
        necessario = (maior + 3) * (LARGURA_MEDIA_CARACTERE * FONTE_LEGENDA_PT / 72)
        # Como na faixa de rótulo, a conta é em polegadas: o piso em fração da
        # largura roubava um quarto do gráfico de quem tem nome curto, e os
        # meses do eixo ficavam apinhados a ponto de "Jul/26" encostar no
        # vizinho.
        gutter = min(
            GUTTER_LEGENDA_MAXIMO,
            max(LARGURA_MINIMA_LEGENDA_IN, necessario) / largura_in,
        )
    else:
        gutter = gutter_linha
    _layout_manual_do_plot(
        grafico._chartSpace.chart.plotArea,
        gutter,
        ALTURA_PLOT_COLUNA if empilhado else ALTURA_PLOT_LINHA,
    )
    if grafico.has_legend:
        _layout_manual_da_legenda(
            grafico.legend, gutter, len(ordem), int(altura_grafico)
        )

    secundario = False
    if tipo_grafico in {"column_line", "line"} and series_secundarias:
        indices = [i for i, nome in enumerate(ordem) if nome in series_secundarias]
        cor_secundaria = next(
            (
                _hex_para_rgb(painel.cores[nome])
                for nome in ordem
                if nome in series_secundarias and nome in painel.cores
            ),
            None,
        )
        secundario = _mover_para_eixo_secundario(
            grafico,
            indices,
            getattr(painel, "formato_secundario", formato_numero),
            cor_secundaria,
            escala=escalas.get(True),
        )

    return {
        "titulo": painel.titulo,
        "series": len(ordem),
        "categorias": len(categorias),
        "rotulos": len(ordem),
        "eixo_secundario": secundario,
    }


def exportar_paineis_pptx(
    paineis: Sequence[Any],
    *,
    rotulo_serie_fn=None,
    titulo_deck: Optional[str] = None,
) -> Tuple[bytes, Dict[str, Any]]:
    """Monta o PPTX com quatro painéis por slide e devolve os bytes.

    ``rotulo_serie_fn`` traduz o nome cru da série para o rótulo da legenda
    (por exemplo, "Mais de 1 a 2 salários mínimos" -> "1 a 2").
    """
    if not paineis:
        raise ValueError("nenhum painel para exportar")

    rotulo_serie_fn = rotulo_serie_fn or (lambda nome: nome)

    prs = Presentation()
    prs.slide_width = SLIDE_LARGURA
    prs.slide_height = SLIDE_ALTURA
    layout_branco = prs.slide_layouts[6]

    largura_quadrante = Emu(int((SLIDE_LARGURA - 2 * MARGEM - GUTTER_H) / 2))
    altura_quadrante = Emu(
        int((SLIDE_ALTURA - 2 * MARGEM - TOPO_CONTEUDO - GUTTER_V) / 2)
    )

    resumo: List[Dict[str, Any]] = []
    slides = 0
    for inicio in range(0, len(paineis), PAINEIS_POR_SLIDE):
        bloco = paineis[inicio:inicio + PAINEIS_POR_SLIDE]
        slide = prs.slides.add_slide(layout_branco)
        slides += 1

        if titulo_deck:
            _caixa_texto(
                slide, left=MARGEM, top=Inches(0.16),
                width=Emu(int(SLIDE_LARGURA - 2 * MARGEM)), height=Inches(0.3),
                texto=titulo_deck, tamanho=11, cor=COR_SUBTITULO,
            )

        for posicao, painel in enumerate(bloco):
            coluna, linha = posicao % 2, posicao // 2
            left = Emu(int(MARGEM + coluna * (largura_quadrante + GUTTER_H)))
            top = Emu(int(MARGEM + TOPO_CONTEUDO + linha * (altura_quadrante + GUTTER_V)))
            resumo.append(_adicionar_painel(
                slide, painel,
                left=left, top=top,
                width=largura_quadrante, height=altura_quadrante,
                rotulo_serie_fn=rotulo_serie_fn,
            ))

    buffer = BytesIO()
    prs.save(buffer)
    return buffer.getvalue(), {
        "slides": slides,
        "paineis": len(paineis),
        "paineis_por_slide": PAINEIS_POR_SLIDE,
        "formato_percentual": FORMATO_PERCENTUAL,
        "detalhe": resumo,
    }


# =============================================================================
# DECK POR SEÇÃO
# =============================================================================
# Um deck contínuo: a leitura dos dados de cada aba na faixa de topo e os
# gráficos daquela aba na grade abaixo. Sem slide divisor, sem régua, sem
# moldura — só caixas de texto e gráficos nativos do Office.

FONTE_CAPA_PT = 26
FONTE_CAPA_SUB_PT = 12
FONTE_SECAO_PT = 20
FONTE_COMENTARIO_PT = 10.5
FONTE_NOME_ABA_PT = 10
ENTRELINHA_COMENTARIO = 1.5

LARGURA_UTIL = Emu(int(SLIDE_LARGURA - 2 * MARGEM))
ALTURA_UTIL = Emu(int(SLIDE_ALTURA - 2 * MARGEM - TOPO_CONTEUDO))
COLUNAS_GRADE = 2
LINHAS_GRADE = 2
LARGURA_CELULA = Emu(int((LARGURA_UTIL - GUTTER_H) / COLUNAS_GRADE))

ALTURA_NOME_ABA = Inches(0.30)
ALTURA_TITULO_SECAO = Inches(0.34)
RESPIRO_GRADE = Inches(0.16)
MAXIMO_GRAFICOS_POR_SLIDE = COLUNAS_GRADE * LINHAS_GRADE
# Uma linha de respiro entre o título e o primeiro parágrafo, no lugar do vão
# de 0,80 in que separava os dois em slides diferentes.
RESPIRO_TITULO = Inches(0.26)
# O PowerPoint conta "1,5 linhas" sobre a altura natural da fonte, não sobre o
# corpo: em Calibri isso dá 1,22x o corpo. Usar o corpo direto subdimensionava
# a caixa e a linha de fontes entrava por cima do último parágrafo.
FATOR_LINHA_CALIBRI = 1.22


def altura_de_linha(corpo_pt: float, entrelinha: float = ENTRELINHA_COMENTARIO) -> int:
    return int(Inches(corpo_pt * FATOR_LINHA_CALIBRI * entrelinha / 72))

# Largura média de caractere ~0,48 em. Serve para estimar quantas linhas o
# comentário ocupa e decidir entre faixa e lâmina própria.
LARGURA_MEDIA_CARACTERE = 0.48
# Acima disso o comentário passa a 8 pt e continua na mesma lâmina dos
# gráficos, em vez de ocupar uma lâmina só para ele.
MAXIMO_LINHAS_CONFORTAVEIS = 5
FONTE_COMENTARIO_COMPACTA_PT = 8

LARGURA_LEITURA = Inches(9.6)


def _caracteres_por_linha(largura_emu: int, corpo_pt: float) -> int:
    polegadas = largura_emu / 914400.0
    return max(int(polegadas / (LARGURA_MEDIA_CARACTERE * corpo_pt / 72.0)), 1)


def linhas_do_comentario(texto: str, largura_emu: int, corpo_pt: float) -> int:
    """Linhas que o comentário ocupa numa caixa dessa largura."""
    por_linha = _caracteres_por_linha(largura_emu, corpo_pt)
    return sum(
        max(1, -(-len(bloco) // por_linha))
        for bloco in str(texto).split("\n\n")
        if bloco.strip()
    )


def _slide_em_branco(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _slide_de_capa(prs, *, titulo: str, subtitulo: str, rodape: str) -> None:
    slide = _slide_em_branco(prs)
    _caixa_texto(
        slide, left=MARGEM, top=Inches(2.5),
        width=LARGURA_UTIL, height=Inches(1.0),
        texto=titulo, tamanho=FONTE_CAPA_PT, cor=COR_TITULO, negrito=True,
    )
    _caixa_texto(
        slide, left=MARGEM, top=Inches(3.5),
        width=LARGURA_LEITURA, height=Inches(0.5),
        texto=subtitulo, tamanho=FONTE_CAPA_SUB_PT, cor=COR_SUBTITULO,
    )
    _caixa_texto(
        slide, left=MARGEM, top=Inches(6.6),
        width=LARGURA_LEITURA, height=Inches(0.4),
        texto=rodape, tamanho=FONTE_FONTE_PT, cor=COR_FONTE,
    )


def _bloco_de_leitura(
    slide, *, titulo: str, texto: str, fontes: str, linhas: int, largura: int,
    corpo: float = FONTE_COMENTARIO_PT,
) -> int:
    """Título e comentário no topo do slide. Devolve a altura ocupada."""
    _caixa_texto(
        slide, left=MARGEM, top=MARGEM,
        width=largura, height=ALTURA_TITULO_SECAO,
        texto=titulo, tamanho=FONTE_SECAO_PT, cor=COR_TITULO, negrito=True,
    )
    topo_texto = int(MARGEM + ALTURA_TITULO_SECAO + RESPIRO_TITULO)
    altura_texto = linhas * altura_de_linha(corpo)
    _caixa_texto(
        slide, left=MARGEM, top=Emu(topo_texto),
        width=largura, height=Emu(altura_texto),
        texto=texto, tamanho=corpo, cor=COR_SUBTITULO,
        entrelinha=ENTRELINHA_COMENTARIO, espaco_entre_paragrafos=0,
    )
    altura = int(ALTURA_TITULO_SECAO) + int(RESPIRO_TITULO) + altura_texto
    if fontes:
        _caixa_texto(
            slide, left=MARGEM, top=Emu(int(MARGEM) + altura + int(Inches(0.04))),
            width=largura, height=Inches(0.18),
            texto=fontes, tamanho=FONTE_FONTE_PT, cor=COR_FONTE,
        )
        altura += int(Inches(0.04)) + int(Inches(0.18))
    return altura


def _slide_de_leitura(prs, *, titulo: str, texto: str, fontes: str) -> None:
    """Comentário longo demais para a faixa ganha lâmina própria."""
    slide = _slide_em_branco(prs)
    linhas = linhas_do_comentario(texto, int(LARGURA_LEITURA), FONTE_COMENTARIO_PT)
    _bloco_de_leitura(
        slide, titulo=titulo, texto=texto, fontes=fontes,
        linhas=linhas, largura=int(LARGURA_LEITURA),
    )


def _desenhar_grade(
    slide, paineis: Sequence[Any], *, altura_cabecalho: int, rotulo_serie_fn
) -> List[Dict[str, Any]]:
    """Grade de duas colunas que preenche a altura livre do slide.

    A altura da célula sai do que sobra abaixo do cabeçalho, dividido pelo
    número de linhas que a seção precisa. Slide com dois gráficos usa a altura
    inteira; com três ou quatro, duas linhas. Célula de altura fixa deixava
    metade do slide vazia sempre que a faixa de leitura era curta.
    """
    if not paineis:
        return []
    topo_grade = int(MARGEM) + altura_cabecalho + int(RESPIRO_GRADE)
    disponivel = int(SLIDE_ALTURA) - int(MARGEM) - topo_grade
    linhas = -(-len(paineis) // COLUNAS_GRADE)
    altura_celula = int((disponivel - (linhas - 1) * int(GUTTER_V)) / linhas)

    resumo: List[Dict[str, Any]] = []
    for posicao, painel in enumerate(paineis):
        coluna, linha = posicao % COLUNAS_GRADE, posicao // COLUNAS_GRADE
        # Linha incompleta fica centralizada, em vez de encostar à esquerda.
        nesta_linha = min(len(paineis) - linha * COLUNAS_GRADE, COLUNAS_GRADE)
        # Gráfico sozinho na linha ocupa a largura inteira, em vez de metade
        # dela com o outro lado vazio.
        largura = LARGURA_UTIL if nesta_linha == 1 else LARGURA_CELULA
        ocupado = nesta_linha * int(largura) + (nesta_linha - 1) * int(GUTTER_H)
        recuo = int((int(LARGURA_UTIL) - ocupado) / 2)
        resumo.append(_adicionar_painel(
            slide, painel,
            left=Emu(int(MARGEM) + recuo + coluna * (int(largura) + int(GUTTER_H))),
            top=Emu(topo_grade + linha * (altura_celula + int(GUTTER_V))),
            width=largura, height=Emu(altura_celula),
            rotulo_serie_fn=rotulo_serie_fn,
            cabecalho_compacto=True,
        ))
    return resumo


def exportar_deck_por_secao(
    secoes: Sequence[Tuple[str, Optional[Any], Sequence[Any]]],
    *,
    titulo_deck: str,
    subtitulo_capa: str = "",
    rodape_capa: str = "",
    rotulo_serie_fn=None,
) -> Tuple[bytes, Dict[str, Any]]:
    """Deck contínuo de várias abas, na ordem em que aparecem na tela.

    ``secoes`` é uma sequência de ``(titulo, leitura, paineis)``, onde
    ``leitura`` é ``(texto, fontes)`` ou ``None``.
    """
    if not secoes:
        raise ValueError("nenhuma seção para exportar")

    rotulo_serie_fn = rotulo_serie_fn or (lambda nome: nome)
    prs = Presentation()
    prs.slide_width = SLIDE_LARGURA
    prs.slide_height = SLIDE_ALTURA

    _slide_de_capa(
        prs, titulo=titulo_deck, subtitulo=subtitulo_capa, rodape=rodape_capa
    )

    resumo: List[Dict[str, Any]] = []
    slides = 1
    for titulo, leitura, paineis in secoes:
        texto, fontes = leitura if leitura else ("", "")
        texto = str(texto or "")
        corpo = FONTE_COMENTARIO_PT
        linhas_faixa = (
            linhas_do_comentario(texto, int(LARGURA_UTIL), corpo)
            if texto.strip() else 0
        )
        if linhas_faixa > MAXIMO_LINHAS_CONFORTAVEIS:
            # Comentário longo encolhe para 8 pt e fica na mesma lâmina dos
            # gráficos, em vez de ocupar uma lâmina só para ele.
            corpo = FONTE_COMENTARIO_COMPACTA_PT
            linhas_faixa = linhas_do_comentario(texto, int(LARGURA_UTIL), corpo)

        restantes = list(paineis)
        primeiro = True
        while restantes or (primeiro and texto.strip()):
            slide = _slide_em_branco(prs)
            slides += 1
            if primeiro and texto.strip():
                altura_cabecalho = _bloco_de_leitura(
                    slide, titulo=titulo, texto=texto, fontes=fontes,
                    linhas=linhas_faixa, largura=int(LARGURA_UTIL), corpo=corpo,
                )
            else:
                _caixa_texto(
                    slide, left=MARGEM, top=MARGEM,
                    width=LARGURA_UTIL, height=ALTURA_NOME_ABA,
                    texto=titulo, tamanho=FONTE_NOME_ABA_PT, cor=COR_FONTE,
                )
                altura_cabecalho = int(ALTURA_NOME_ABA)
            bloco = restantes[:MAXIMO_GRAFICOS_POR_SLIDE]
            restantes = restantes[MAXIMO_GRAFICOS_POR_SLIDE:]
            resumo.extend(_desenhar_grade(
                slide, bloco,
                altura_cabecalho=altura_cabecalho,
                rotulo_serie_fn=rotulo_serie_fn,
            ))
            primeiro = False

    buffer = BytesIO()
    prs.save(buffer)
    return buffer.getvalue(), {
        "slides": slides,
        "paineis": len(resumo),
        "secoes": len(secoes),
        "paineis_por_slide": MAXIMO_GRAFICOS_POR_SLIDE,
        "detalhe": resumo,
    }
